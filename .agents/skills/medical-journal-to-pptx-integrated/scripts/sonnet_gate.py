#!/usr/bin/env python3
"""
qa_gate.py — one enforced quality gate for the medical-journal-to-pptx pipeline.

Why this exists
---------------
Older versions of this skill relied on a long human-style checklist in SKILL.md
that the model was supposed to read and self-attest ("[ ] every slide has
notes", "[ ] no Chinese in visible text", ...). A strong model (Opus) can hold
all of those rules in its head and honestly verify them. A weaker model tends to
tick the boxes without really checking, or forget a rule under context pressure.

This script moves that judgment OUT of the model's head and INTO deterministic
code. Instead of trusting the model to remember ~40 rules, we let it run one
command that mechanically checks them and prints, for every failure, the exact
fix to make. The model's job becomes the thing models are reliably good at:
read a concrete error, fix that one thing, run again until green. That is why a
Sonnet-class model using this gate can reach the quality that previously needed
Opus — the invariants are enforced, not remembered.

Modes
-----
  spec  <deck_spec.json>              content + asset checks (run before build)
  pptx  <deck.pptx> [--spec S]        built-file checks (run after build)
  all   <deck_spec.json> --pptx P     both (recommended final gate)

Exit code is non-zero if ANY check fails, so the pipeline cannot advance past a
broken deck. WARN items do not fail the build but should be looked at.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

MODE_LIMITS = {"lite": (8, 16), "full": (40, 55), "smoke": (1, 55)}
TABLE_MARGIN_MIN = 8
POSTPROCESS_SUFFIX = ".postprocess.json"

PANEL_FRAGMENT_RE = re.compile(r"(_panel[_\-]?[a-z0-9]+|panel_[a-z0-9]+)\.png$", re.I)
CJK_RE = re.compile("[　-〿぀-ヿ㐀-䶿一-鿿＀-￯]")
EMOJI_RE = re.compile(
    "["
    "\U0001F300-\U0001FAFF"
    "\U00002600-\U000027BF"
    "\U00002190-\U000021FF"
    "\U00002B00-\U00002BFF"
    "\U0001F1E6-\U0001F1FF"
    "✅⚠❤➡⭐→ℹ⃣️"
    "]"
)
KEYCAP_RE = re.compile("[0-9]️?⃣")


class Report:
    def __init__(self):
        self.items = []

    def ok(self, c, m): self.items.append(("PASS", c, m, ""))
    def warn(self, c, m, fix=""): self.items.append(("WARN", c, m, fix))
    def fail(self, c, m, fix=""): self.items.append(("FAIL", c, m, fix))

    @property
    def n_fail(self): return sum(1 for lv, *_ in self.items if lv == "FAIL")
    @property
    def n_warn(self): return sum(1 for lv, *_ in self.items if lv == "WARN")

    def render(self):
        icon = {"PASS": "[PASS]", "WARN": "[WARN]", "FAIL": "[FAIL]"}
        lines = []
        for lv, c, m, fix in self.items:
            lines.append(f"{icon[lv]} {c}: {m}")
            if fix and lv != "PASS":
                lines.append(f"        -> fix: {fix}")
        n_pass = sum(1 for lv, *_ in self.items if lv == "PASS")
        lines.append("")
        lines.append(f"SUMMARY: {self.n_fail} failing, {self.n_warn} warnings, {n_pass} passing.")
        if self.n_fail:
            lines.append("RESULT: NOT READY -- fix every [FAIL] above and re-run qa_gate.")
        elif self.n_warn:
            lines.append("RESULT: PASSES gate (review [WARN] items; they do not block).")
        else:
            lines.append("RESULT: PASSES gate cleanly.")
        return "\n".join(lines)


def visible_text_fields(slide):
    out = []
    for key in ("title", "caption", "authors", "citation", "subtitle", "kicker"):
        v = slide.get(key)
        if isinstance(v, str) and v.strip():
            out.append((key, v))
    for key in ("bullets", "items", "body"):
        v = slide.get(key)
        if isinstance(v, list):
            for i, item in enumerate(v):
                if isinstance(item, str) and item.strip():
                    out.append((f"{key}[{i}]", item))
        elif isinstance(v, str) and v.strip():
            out.append((key, v))
    return out


def figure_number(slide):
    blob = f"{slide.get('caption','')} {slide.get('title','')}"
    if re.search(r"\btable\b", blob, re.I):
        return None
    m = re.search(r"\bfig(?:ure)?\.?\s*([0-9]+)", blob, re.I)
    return f"Figure {m.group(1)}" if m else None


def referenced_panel_letters(notes):
    letters = set()
    for m in re.finditer(r"【\s*([A-Za-z])\s*[圖图]", notes):
        letters.add(m.group(1).upper())
    for m in re.finditer(r"\bpanel\s*([A-Za-z])\b", notes, re.I):
        letters.add(m.group(1).upper())
    return letters


def read_sidecar(image_path):
    side = image_path.with_suffix(image_path.suffix + POSTPROCESS_SUFFIX)
    if side.exists():
        try:
            return json.loads(side.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return None


def looks_like_table(image_path, sidecar):
    if "table" in image_path.name.lower():
        return True
    if sidecar and str(sidecar.get("asset_type", "")).lower() == "table":
        return True
    return False


def table_split_group(image_path):
    m = re.match(r"(Table[_\-]?\d+)[A-Z]\.png$", image_path.name, re.I)
    return m.group(1).lower() if m else None


def check_spec(spec_path, rep, *, content_mode="full", style="standard"):
    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    spec_dir = spec_path.parent
    slides = spec.get("slides", [])
    meta = spec.get("meta", {})

    n = len(slides)
    minimum, maximum = MODE_LIMITS[content_mode]
    if minimum <= n <= maximum:
        rep.ok("slide_count", f"{n} slides (target {minimum}-{maximum}; {content_mode})")
    else:
        rep.fail("slide_count", f"{n} slides is outside {minimum}-{maximum}",
                 f"Adjust the deck to the {content_mode} mode slide budget.")

    types = [s.get("type") for s in slides]
    if types.count("title") == 1:
        rep.ok("title_slide", "exactly one title slide")
    else:
        rep.fail("title_slide", f"found {types.count('title')} title slides",
                 "The deck needs exactly one 'title' slide as slide 1.")
    required_types = ("outline", "references", "thanks") if content_mode == "full" else ("thanks",)
    for req in required_types:
        if req in types:
            rep.ok(f"has_{req}", f"'{req}' slide present")
        else:
            (rep.warn if req == "outline" else rep.fail)(
                f"has_{req}", f"no '{req}' slide found",
                f"Add a '{req}' slide (see references/slide_structure.md).")

    title = next((s for s in slides if s.get("type") == "title"), None)
    if title:
        missing = [k for k in ("authors", "citation") if not str(title.get(k, "")).strip()]
        if missing:
            rep.fail("title_metadata", f"title slide missing {missing}",
                     "Fill authors and citation from the PDF front matter.")
        else:
            rep.ok("title_metadata", "title slide has authors + citation")

    cjk_hits = []
    for idx, s in enumerate(slides, 1):
        for field, text in visible_text_fields(s):
            if CJK_RE.search(text):
                sample = "".join(CJK_RE.findall(text)[:6])
                cjk_hits.append(f"slide {idx} ({s.get('type')}) {field}: {sample}")
    if isinstance(meta.get("footer_label"), str) and CJK_RE.search(meta["footer_label"]):
        cjk_hits.append("meta.footer_label")
    if cjk_hits:
        rep.fail("visible_text_english", f"{len(cjk_hits)} visible fields contain CJK",
                 "Slide-visible text must be English only (titles, bullets, captions, "
                 "outline items, authors, citation, footer). Move Chinese into notes. "
                 "First offenders: " + " | ".join(cjk_hits[:5]))
    else:
        rep.ok("visible_text_english", "no CJK in any slide-visible text")

    no_notes = [i for i, s in enumerate(slides, 1) if not str(s.get("notes", "")).strip()]
    if no_notes:
        rep.fail("notes_present", f"{len(no_notes)} slides have empty notes: {no_notes[:10]}",
                 "Every slide needs Traditional-Chinese speaker notes.")
    else:
        rep.ok("notes_present", "all slides have notes")

    total_emoji = sum(len(EMOJI_RE.findall(str(s.get("notes", "")))) for s in slides)
    if total_emoji == 0 and slides:
        rep.fail("notes_emoji", "notes contain zero scan/lead emojis",
                 "Notes must use emoji scaffolding so they are scannable while "
                 "speaking. See references/notes_style.md.")
    else:
        rep.ok("notes_emoji", f"{total_emoji} scaffolding emoji across notes")

    literal_bold = [i for i, s in enumerate(slides, 1) if "**" in str(s.get("notes", ""))]
    if literal_bold:
        rep.warn("notes_bold_markup",
                 f"{len(literal_bold)} slides use **bold** markup in notes source",
                 "Expected in source (builder converts to real bold). Confirm the BUILT "
                 "pptx shows no literal '**' via `qa_gate.py pptx`.")

    outline = next((s for s in slides if s.get("type") == "outline"), None)
    if outline:
        items = outline.get("items") or []
        keycapped = sum(1 for it in items if KEYCAP_RE.search(str(it)))
        ranged = sum(1 for it in items if re.search(r"[Ss]lides?\s*\d+\s*[-–]\s*\d+", str(it)))
        if items and keycapped >= max(1, len(items) - 1):
            rep.ok("outline_markers", f"{keycapped}/{len(items)} outline items numbered")
        else:
            rep.warn("outline_markers", "outline items are not numbered with keycap emoji",
                     "Prefix each outline item like '1️⃣ Background - Slides 3-7'.")
        if items and ranged < len(items):
            rep.warn("outline_ranges", f"only {ranged}/{len(items)} items have a Slides X-Y hint",
                     "Add 'Slides A-B' to each outline item once the sequence is fixed.")

    fig_numbers = {}
    for idx, s in enumerate(slides, 1):
        if s.get("type") != "figure":
            continue
        img = s.get("image")
        if not img:
            rep.fail("figure_image", f"slide {idx} is a figure with no 'image'",
                     "Point the figure slide at its final asset.")
            continue
        candidate = Path(img).expanduser()
        img_path = candidate if candidate.is_absolute() else spec_dir / candidate
        name = Path(img).name

        if not img_path.exists():
            rep.fail("image_exists", f"slide {idx} image not found: {img}",
                     "Build/relocate the asset so the spec path resolves.")
            continue

        if PANEL_FRAGMENT_RE.search(name):
            rep.fail("raw_panel_ref", f"slide {idx} references raw panel crop {name}",
                     "Recompose panels into ONE figure image before referencing it. "
                     "One figure = one slide = one image.")

        if img_path.suffix.lower() == ".png" and read_sidecar(img_path) is None:
            rep.fail("asset_sidecar", f"slide {idx} asset {name} has no {POSTPROCESS_SUFFIX} sidecar",
                     "Run the asset through postprocess_assets.py so it writes a sidecar; "
                     "build_deck.py requires it.")

        labels = s.get("panel_labels") or []
        if len(labels) > 1 and not s.get("panel_geometry_exception"):
            sidecar = read_sidecar(img_path) or {}
            has_geo = (len(s.get("panel_label_x_fracs") or []) >= len(labels)
                       or len(s.get("panel_boxes") or []) >= len(labels)
                       or len(sidecar.get("panel_boxes") or []) >= len(labels)
                       or bool(sidecar.get("native_labels")))
            if not has_geo:
                rep.fail("panel_geometry",
                         f"slide {idx} has {len(labels)} panel_labels but no geometry",
                         "Add panel_label_x_fracs or panel_boxes, OR use the native "
                         "post-build label flow (recompose_panels_banded.py + "
                         "add_panel_labels.py) and omit panel_labels from the spec.")

        lab_set = {str(l).upper() for l in labels}
        if lab_set:
            ghost = referenced_panel_letters(str(s.get("notes", ""))) - lab_set
            if ghost:
                rep.fail("panel_notes_ref",
                         f"slide {idx} notes mention panel(s) {sorted(ghost)} not in {sorted(lab_set)}",
                         "Only reference panels that exist, using the panel-label style.")

        notes = str(s.get("notes", ""))
        if notes and not notes.lstrip().startswith("【"):
            rep.warn("figure_notes_marker", f"slide {idx} figure notes don't open with a bracket header",
                     "Open figure/table notes with the image-caption header for scannability.")

        side = read_sidecar(img_path)
        if looks_like_table(img_path, side) and side:
            if side.get("command") in {"trim", "labels"} and isinstance(side.get("margin"), int) \
                    and side["margin"] < TABLE_MARGIN_MIN:
                rep.fail("table_margin", f"slide {idx} table {name} margin={side['margin']} (< {TABLE_MARGIN_MIN}px)",
                         "Re-trim with --margin 12 (safety band 8-24px) so text/gridlines "
                         "never touch the edge.")

        fn = figure_number(s)
        if fn:
            fig_numbers.setdefault(fn, []).append(idx)

    dupes = {k: v for k, v in fig_numbers.items() if len(v) > 1}
    if dupes:
        rep.fail("one_figure_one_slide",
                 "same figure on multiple slides: "
                 + "; ".join(f"{k} on slides {v}" for k, v in dupes.items()),
                 "Merge each paper Figure into ONE slide with ONE recomposed image.")
    elif fig_numbers:
        rep.ok("one_figure_one_slide", f"{len(fig_numbers)} figures each on one slide")

    groups = {}
    for s in slides:
        img = s.get("image")
        if not img:
            continue
        g = table_split_group(Path(img))
        if g:
            groups.setdefault(g, []).append(s)
    try:
        from PIL import Image
        for g, gslides in groups.items():
            if len(gslides) < 2:
                continue
            widths, screen = {}, set()
            for s in gslides:
                image = Path(s["image"]).expanduser()
                p = image if image.is_absolute() else spec_dir / image
                if p.exists():
                    with Image.open(p) as im:
                        widths[Path(s["image"]).name] = im.width
                screen.add(s.get("image_width_in"))
            if len(set(widths.values())) > 1:
                rep.fail("table_split_px_width", f"{g}: unequal pixel widths {widths}",
                         "Normalize split parts to the same canvas width (same-width).")
            elif None in screen or len(screen) > 1:
                rep.warn("table_split_screen_width", f"{g}: image_width_in differs {screen}",
                         "Give every split part the SAME image_width_in so they render at "
                         "equal on-screen width (SKILL section 4A).")
            else:
                rep.ok("table_split_width", f"{g}: split parts share width")
    except ImportError:
        pass

    flat = []
    content = [s for s in slides if s.get("type") == "content"]
    for idx, s in enumerate(slides, 1):
        if s.get("type") != "content":
            continue
        body_items = []
        for key in ("bullets", "body", "items"):
            v = s.get(key)
            if isinstance(v, list):
                body_items = [str(x) for x in v]
                break
            if isinstance(v, str):
                body_items = [ln for ln in v.splitlines() if ln.strip()]
                break
        if len(body_items) < 3:
            continue
        joined = "\n".join(body_items)
        has_label = any(re.search(r"[A-Za-z][\w /-]{2,}:", ln) for ln in body_items)
        has_logic = ("→" in joined) or ("✅" in joined) or ("⚠" in joined)
        if not has_label and not has_logic:
            flat.append(idx)
    if content:
        if len(flat) > len(content) / 2:
            rep.fail("content_structure",
                     f"{len(flat)}/{len(content)} content slides are flat bullet lists (slides {flat[:10]})",
                     "Rewrite as teaching blocks: 2-4 short 'Label:' sections with a "
                     "supporting line, using -> for consequence and one take-home line. "
                     "SKILL section 6A.")
        elif flat:
            rep.warn("content_structure",
                     f"{len(flat)} content slides look like flat bullet lists (slides {flat})",
                     "Consider adding section labels / a consequence or take-home line.")
        else:
            rep.ok("content_structure", "content slides use structured teaching blocks")

    logo = meta.get("logo_path")
    if logo:
        if (spec_dir / logo).exists() or Path(logo).exists():
            rep.ok("logo_path", "meta.logo_path resolves")
        else:
            rep.warn("logo_path", f"meta.logo_path {logo!r} does not resolve",
                     "Remove it (builder falls back to bundled logo) or fix the path.")
    else:
        rep.ok("logo_path", "no meta.logo_path (builder uses bundled default logo)")


def check_pptx(pptx_path, rep, spec_path=None, *, content_mode="full", style="standard"):
    try:
        from pptx import Presentation
    except ImportError:
        rep.fail("pptx_import", "python-pptx not installed",
                 "pip install python-pptx --break-system-packages")
        return

    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)
    n = len(slides)
    minimum, maximum = MODE_LIMITS[content_mode]
    if minimum <= n <= maximum:
        rep.ok("pptx_slide_count", f"{n} slides")
    else:
        rep.fail("pptx_slide_count", f"{n} slides outside {minimum}-{maximum}",
                 f"Rebuild the deck within the {content_mode} mode slide budget.")

    empty, literal, emoji = [], [], 0
    for i, sl in enumerate(slides, 1):
        txt = sl.notes_slide.notes_text_frame.text if sl.has_notes_slide else ""
        if not txt.strip():
            empty.append(i)
            continue
        if "**" in txt:
            literal.append(i)
        emoji += len(EMOJI_RE.findall(txt))
    if empty:
        rep.fail("pptx_notes_present", f"slides with no notes: {empty[:10]}",
                 "Every slide must carry speaker notes.")
    else:
        rep.ok("pptx_notes_present", "all slides have notes")
    if literal:
        rep.fail("pptx_literal_bold", f"slides show literal '**' in notes: {literal[:10]}",
                 "Builder should convert **term** to real bold. Rebuild with current "
                 "build_deck.py; do not hand-paste ** into notes.")
    else:
        rep.ok("pptx_literal_bold", "no literal '**' markup in built notes")
    if emoji == 0:
        rep.fail("pptx_notes_emoji", "built notes contain zero emoji",
                 "Notes lost their emoji scaffolding during authoring.")
    else:
        rep.ok("pptx_notes_emoji", f"{emoji} note emoji in built deck")

    body_without_pic = []
    for i, sl in enumerate(slides, 1):
        if i == 1 or i == n:
            continue
        pics = [sh for sh in sl.shapes if sh.shape_type == 13]
        if not pics:
            body_without_pic.append(i)
    if body_without_pic:
        rep.warn("pptx_logo_on_slides",
                 f"{len(body_without_pic)} body slides have no picture (logo?): {body_without_pic[:10]}",
                 "Every non-title/non-thanks slide should show the upper-right logo. "
                 "Check meta.logo_path / bundled fallback.")
    else:
        rep.ok("pptx_logo_on_slides", "every body slide carries at least the logo picture")

    cjk_slides = []
    for i, sl in enumerate(slides, 1):
        for sh in sl.shapes:
            if sh.has_text_frame and CJK_RE.search(sh.text_frame.text or ""):
                cjk_slides.append(i)
                break
    if cjk_slides:
        rep.fail("pptx_visible_english", f"CJK found in visible text on slides {cjk_slides[:10]}",
                 "Slide-visible text must be English. Move Chinese to notes and rebuild.")
    else:
        rep.ok("pptx_visible_english", "no CJK in built slide-visible text")

    if spec_path is not None:
        expected_slides = json.loads(spec_path.read_text(encoding="utf-8")).get("slides", [])
        missing = []
        checked = 0
        for index, (slide_spec, built_slide) in enumerate(zip(expected_slides, slides), start=1):
            if slide_spec.get("type") != "figure":
                continue
            labels = slide_spec.get("panel_labels") or []
            image_name = slide_spec.get("image")
            if not labels and isinstance(image_name, str):
                image_path = Path(image_name).expanduser()
                if not image_path.is_absolute():
                    image_path = spec_path.parent / image_path
                sidecar = read_sidecar(image_path) or {}
                if sidecar.get("native_labels"):
                    labels = sidecar.get("labels") or []
            if not isinstance(labels, list) or not labels:
                continue
            checked += len(labels)
            visible = {
                shape.text_frame.text.strip().upper()
                for shape in built_slide.shapes
                if getattr(shape, "has_text_frame", False)
            }
            absent = [str(label).strip().upper() for label in labels
                      if str(label).strip().upper() not in visible]
            if absent:
                missing.append(f"slide {index}: {', '.join(absent)}")
        if missing:
            rep.fail("pptx_panel_labels", "; ".join(missing),
                     "Render the spec panel labels or run add_panel_labels.py before final QA.")
        elif checked:
            rep.ok("pptx_panel_labels", f"all {checked} expected native panel labels are visible")


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    sub = ap.add_subparsers(dest="mode", required=True)
    sp = sub.add_parser("spec"); sp.add_argument("spec")
    pp = sub.add_parser("pptx"); pp.add_argument("pptx"); pp.add_argument("--spec", default=None)
    al = sub.add_parser("all"); al.add_argument("spec"); al.add_argument("--pptx", required=True)
    for parser in (sp, pp, al):
        parser.add_argument("--content-mode", choices=tuple(MODE_LIMITS), default="full")
        parser.add_argument("--style", choices=("standard", "nice"), default="standard")
    args = ap.parse_args()
    rep = Report()
    if args.mode == "spec":
        check_spec(Path(args.spec).resolve(), rep, content_mode=args.content_mode, style=args.style)
    elif args.mode == "pptx":
        check_pptx(Path(args.pptx).resolve(), rep, Path(args.spec).resolve() if args.spec else None,
                   content_mode=args.content_mode, style=args.style)
    elif args.mode == "all":
        check_spec(Path(args.spec).resolve(), rep, content_mode=args.content_mode, style=args.style)
        check_pptx(Path(args.pptx).resolve(), rep, Path(args.spec).resolve(),
                   content_mode=args.content_mode, style=args.style)
    print(rep.render())
    sys.exit(1 if rep.n_fail else 0)


if __name__ == "__main__":
    main()
