#!/usr/bin/env python3
"""Synthetic regression tests for slide-aware multipanel layout selection."""

from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from PIL import Image, ImageDraw

import build_deck
import image_polarity
from postprocess_assets import content_bbox as content_bbox_for_test
from recompose_panels_banded import clean_panel_edges, residual_edge_review


SCRIPT = Path(__file__).with_name("recompose_panels_banded.py")
POSTPROCESS = Path(__file__).with_name("postprocess_assets.py")
ADD_PANEL_LABELS = Path(__file__).with_name("add_panel_labels.py")


class PanelLayoutTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.directory = Path(self.temporary.name)
        self.addCleanup(self.temporary.cleanup)

    def compose(self, width: int, height: int, *extra: str) -> dict:
        inputs = []
        for index in range(4):
            path = self.directory / f"panel_{index}_{width}x{height}.png"
            Image.new("RGB", (width, height), (25 + 35 * index, 80, 120)).save(path)
            inputs.append(path)

        output = self.directory / f"result_{width}x{height}_{len(list(self.directory.glob('result_*.png')))}.png"
        command = [
            sys.executable,
            str(SCRIPT),
            str(output),
            "--inputs",
            *(str(path) for path in inputs),
            "--labels",
            "A,B,C,D",
            "--geometry",
            str(self.directory / "geometry.json"),
            "--no-trim",
            *extra,
        ]
        subprocess.run(command, check=True, capture_output=True, text=True)
        return json.loads(Path(str(output) + ".postprocess.json").read_text())

    def compose_inputs(self, inputs: list[Path], *extra: str) -> tuple[Path, dict, dict]:
        output = self.directory / f"custom_{len(list(self.directory.glob('custom_*.png')))}.png"
        geometry = self.directory / "custom_geometry.json"
        labels = ",".join(chr(ord("A") + index) for index in range(len(inputs)))
        command = [
            sys.executable,
            str(SCRIPT),
            str(output),
            "--inputs",
            *(str(path) for path in inputs),
            "--labels",
            labels,
            "--geometry",
            str(geometry),
            *extra,
        ]
        subprocess.run(command, check=True, capture_output=True, text=True)
        sidecar = json.loads(Path(str(output) + ".postprocess.json").read_text())
        return output, sidecar, json.loads(geometry.read_text())

    def source_panels(self, source: Image.Image, boxes: list[tuple[int, int, int, int]]) -> list[Path]:
        source_path = self.directory / f"source_{len(list(self.directory.glob('source_*.png')))}.png"
        source.save(source_path)
        paths = []
        for index, box in enumerate(boxes):
            label = chr(ord("A") + index)
            path = self.directory / f"{source_path.stem}_panel_{label}.png"
            source.crop(box).save(path)
            Path(str(path) + ".postprocess.json").write_text(json.dumps({
                "source": str(source_path),
                "crop_box_px": list(box),
                "source_label_placement": "embedded",
                "embedded_label": label,
            }))
            paths.append(path)
        return paths

    def test_portrait_panels_use_single_horizontal_row(self) -> None:
        result = self.compose(600, 800)

        self.assertEqual(result["layout_mode"], "auto")
        self.assertEqual((result["rows"], result["cols"]), (1, 4))
        self.assertEqual(len(result["source_inputs"]), 4)
        candidates = {candidate["cols"]: candidate for candidate in result["layout_candidates"]}
        self.assertTrue(all(candidate["safety_margin_px"] == 16 for candidate in candidates.values()))
        self.assertTrue(all(
            candidate["composite_width_px"] == candidate["unpadded_width_px"] + 32
            and candidate["composite_height_px"] == candidate["unpadded_height_px"] + 32
            for candidate in candidates.values()
        ))
        self.assertGreater(
            candidates[4]["min_panel_area_sq_in"],
            candidates[2]["min_panel_area_sq_in"],
        )

    def test_wide_panels_use_multiple_rows(self) -> None:
        result = self.compose(1600, 400)

        self.assertEqual((result["rows"], result["cols"]), (2, 2))

    def test_clinical_grid_rows_fill_right_edge_without_external_canvas(self) -> None:
        inputs = []
        for index, size in enumerate(((301, 241), (319, 243), (337, 245),
                                      (307, 239), (323, 247), (341, 251))):
            path = self.directory / f"rounding-{index}.png"
            Image.new("RGB", size, (40 + index * 20, 75, 115)).save(path)
            inputs.append(path)

        output, sidecar, _ = self.compose_inputs(
            inputs, "--asset-type", "clinical-image", "--cols", "3", "--no-trim"
        )

        width, _ = sidecar["padded_size_px"]
        boxes = sidecar["panel_boxes_px"]
        self.assertEqual(sidecar["safety_margin_px"], 0)
        self.assertEqual(boxes[2]["x"] + boxes[2]["w"], width)
        self.assertEqual(boxes[5]["x"] + boxes[5]["w"], width)
        with Image.open(output) as composed:
            rgb = composed.convert("RGB")
            self.assertNotEqual(rgb.getpixel((width - 1, boxes[2]["h"] // 2)), (6, 20, 40))
            self.assertNotEqual(
                rgb.getpixel((width - 1, boxes[5]["y"] + boxes[5]["h"] // 2)),
                (6, 20, 40),
            )

    def test_figure19_nonshared_horizontal_seams_are_automatically_required(self) -> None:
        source = self.directory / "figure19-source.png"
        image = Image.new("RGB", (300, 140), (10, 10, 10))
        draw = ImageDraw.Draw(image)
        draw.rectangle((0, 0, 99, 69), fill=(30, 45, 60))
        draw.rectangle((100, 0, 199, 70), fill=(110, 25, 35))
        draw.rectangle((200, 0, 299, 72), fill=(35, 75, 130))
        draw.rectangle((0, 70, 99, 139), fill=(185, 190, 195))
        draw.rectangle((100, 71, 199, 139), fill=(210, 105, 45))
        draw.rectangle((200, 73, 299, 139), fill=(215, 205, 55))
        image.save(source)

        reviews = {}
        review_specs = {
            "top_ab": ("x", (0, 71), (95, 105), 100),
            "top_bc": ("x", (0, 73), (195, 205), 200),
            "bottom_de": ("x", (70, 140), (95, 105), 100),
            "bottom_ef": ("x", (71, 140), (195, 205), 200),
            "left_ad": ("y", (0, 100), (65, 76), 70),
            "middle_be": ("y", (100, 200), (66, 77), 71),
            "right_cf": ("y", (200, 300), (68, 79), 73),
        }
        for name, (axis, band, search, selected) in review_specs.items():
            report = self.directory / f"{name}.json"
            overlay = self.directory / f"{name}.png"
            result = subprocess.run([
                sys.executable, str(POSTPROCESS), "seam-review",
                str(source), str(report), str(overlay),
                "--axis", axis,
                "--band", str(band[0]), str(band[1]),
                "--search", str(search[0]), str(search[1]),
                "--selected", str(selected),
                "--tolerance", "0",
            ], capture_output=True, text=True)
            self.assertEqual(result.returncode, 0, result.stderr + result.stdout)
            reviews[name] = report

        boxes = [
            (0, 0, 100, 70), (100, 0, 200, 71), (200, 0, 300, 73),
            (0, 70, 100, 140), (100, 71, 200, 140), (200, 73, 300, 140),
        ]
        bindings = [
            (("right", "top_ab"), ("bottom", "left_ad")),
            (("left", "top_ab"), ("right", "top_bc"), ("bottom", "middle_be")),
            (("left", "top_bc"), ("bottom", "right_cf")),
            (("right", "bottom_de"), ("top", "left_ad")),
            (("left", "bottom_de"), ("right", "bottom_ef"), ("top", "middle_be")),
            (("left", "bottom_ef"), ("top", "right_cf")),
        ]
        panels = []
        for index, (box, panel_bindings) in enumerate(zip(boxes, bindings)):
            panel = self.directory / f"figure19-panel-{index}.png"
            command = [
                sys.executable, str(POSTPROCESS), "panel-crop",
                str(source), str(panel),
                "--box", *(str(value) for value in box),
                "--label", chr(ord("A") + index),
                "--label-placement", "absent",
            ]
            for edge, report_name in panel_bindings:
                command.extend(("--seam-review", str(reviews[report_name]),
                                "--seam-edge", edge,
                                "--require-seam-edge", edge))
            result = subprocess.run(command, capture_output=True, text=True)
            self.assertEqual(result.returncode, 0, result.stderr + result.stdout)
            panels.append(panel)

        output, sidecar, _ = self.compose_inputs(
            panels, "--asset-type", "clinical-image", "--cols", "3", "--no-trim"
        )

        topology = sidecar["source_seam_topology"]
        self.assertEqual(topology["schema"], "medical-journal-source-seam-topology/v1")
        self.assertEqual(len(topology["groups"]), 1)
        self.assertEqual(len(topology["groups"][0]["adjacencies"]), 7)
        self.assertEqual(
            topology["panel_required_edges"],
            [
                {"panel_index": 0, "edges": ["bottom", "right"]},
                {"panel_index": 1, "edges": ["bottom", "left", "right"]},
                {"panel_index": 2, "edges": ["bottom", "left"]},
                {"panel_index": 3, "edges": ["right", "top"]},
                {"panel_index": 4, "edges": ["left", "right", "top"]},
                {"panel_index": 5, "edges": ["left", "top"]},
            ],
        )
        handled, failures = image_polarity._deterministic_helper_evidence(output, sidecar)
        self.assertTrue(handled)
        self.assertFalse(failures)

        tampered = dict(sidecar)
        tampered.pop("source_seam_topology")
        handled, failures = image_polarity._deterministic_helper_evidence(output, tampered)
        self.assertTrue(handled)
        self.assertTrue(any("source-seam topology" in failure for failure in failures))

    def test_clinical_same_source_panel_crops_cannot_bypass_seam_evidence(self) -> None:
        source = self.directory / "legacy-common-seam-source.png"
        image = Image.new("RGB", (120, 80), (25, 50, 75))
        ImageDraw.Draw(image).rectangle((0, 40, 119, 79), fill=(170, 185, 200))
        image.save(source)
        panels = []
        for index, box in enumerate(((0, 0, 120, 40), (0, 40, 120, 80))):
            panel = self.directory / f"legacy-common-seam-{index}.png"
            result = subprocess.run([
                sys.executable, str(POSTPROCESS), "panel-crop",
                str(source), str(panel),
                "--box", *(str(value) for value in box),
                "--label", chr(ord("A") + index),
                "--label-placement", "absent",
            ], capture_output=True, text=True)
            self.assertEqual(result.returncode, 0, result.stderr + result.stdout)
            panels.append(panel)

        result = subprocess.run([
            sys.executable, str(SCRIPT), str(self.directory / "legacy-bypass.png"),
            "--inputs", *(str(path) for path in panels),
            "--labels", "A,B",
            "--geometry", str(self.directory / "legacy-bypass-geometry.json"),
            "--asset-type", "clinical-image", "--cols", "1", "--no-trim",
        ], capture_output=True, text=True)

        self.assertNotEqual(result.returncode, 0)
        self.assertIn("automatically required source-seam evidence", result.stderr)

    def test_explicit_columns_remain_manual_override(self) -> None:
        result = self.compose(600, 800, "--cols", "2")

        self.assertEqual(result["layout_mode"], "manual")
        self.assertEqual((result["rows"], result["cols"]), (2, 2))
        self.assertEqual(len(result["layout_candidates"]), 1)

    def test_nice_style_dimensions_are_respected(self) -> None:
        result = self.compose(
            600,
            800,
            "--slide-box-w-in",
            "12.13",
            "--slide-box-h-in",
            "4.95",
        )

        self.assertEqual((result["rows"], result["cols"]), (1, 4))
        self.assertEqual(result["slide_box_w_in"], 12.13)
        self.assertEqual(result["slide_box_h_in"], 4.95)

    def test_white_gray_and_antialiased_rims_are_removed_with_a_hard_limit(self) -> None:
        image = Image.new("RGB", (120, 100), (8, 8, 8))
        draw = ImageDraw.Draw(image)
        draw.rectangle((0, 0, 119, 1), fill=(250, 250, 250))
        draw.line((0, 2, 119, 2), fill=(143, 143, 143))
        draw.rectangle((0, 98, 119, 99), fill=(116, 116, 116))
        draw.line((0, 0, 0, 99), fill=(255, 255, 255))
        draw.line((119, 0, 119, 99), fill=(144, 144, 144))

        cleaned, edges = clean_panel_edges(image, max_edge_px=4)

        self.assertEqual(edges, {"top": 3, "bottom": 2, "left": 1, "right": 1})
        self.assertEqual(cleaned.size, (118, 95))
        self.assertEqual(cleaned.getpixel((0, 0)), (8, 8, 8))

    def test_single_pdf_boundary_column_is_removed_before_long_dark_canvas(self) -> None:
        for boundary in ((255, 255, 255), (177, 175, 176), (164, 162, 163)):
            with self.subTest(boundary=boundary):
                image = Image.new("RGB", (120, 100), (34, 30, 31))
                ImageDraw.Draw(image).line((119, 0, 119, 99), fill=boundary)

                cleaned, edges = clean_panel_edges(image, max_edge_px=4)

                self.assertEqual(
                    edges,
                    {"top": 0, "bottom": 0, "left": 0, "right": 1},
                )
                self.assertEqual(cleaned.size, (119, 100))
                self.assertEqual(cleaned.getpixel((118, 50)), (34, 30, 31))

    def test_light_image_region_thicker_than_budget_is_not_cropped(self) -> None:
        image = Image.new("RGB", (120, 100), (7, 7, 7))
        ImageDraw.Draw(image).rectangle((0, 0, 119, 11), fill=(245, 245, 245))

        cleaned, edges = clean_panel_edges(image, max_edge_px=4)

        self.assertEqual(edges["top"], 0)
        self.assertEqual(cleaned.size, image.size)
        self.assertEqual(cleaned.getpixel((60, 0)), (245, 245, 245))

    def test_five_pixel_white_edge_is_reported_for_review_without_auto_crop(self) -> None:
        image = Image.new("RGB", (120, 100), (7, 7, 7))
        ImageDraw.Draw(image).rectangle((0, 95, 119, 99), fill=(252, 252, 252))

        cleaned, edges = clean_panel_edges(image, max_edge_px=4)
        review = residual_edge_review(cleaned, max_edge_px=4)

        self.assertEqual(edges["bottom"], 0)
        self.assertEqual(review["status"], "needs-review")
        self.assertEqual(review["candidates"]["bottom"]["depth_px"], 5)

    def test_disabled_cleanup_reports_two_pixel_white_edge_for_review(self) -> None:
        image = Image.new("RGB", (120, 100), (7, 7, 7))
        ImageDraw.Draw(image).rectangle((0, 0, 1, 99), fill=(252, 252, 252))

        review = residual_edge_review(
            image,
            max_edge_px=4,
            include_bounded=True,
        )

        self.assertEqual(review["status"], "needs-review")
        self.assertEqual(review["candidates"]["left"]["depth_px"], 2)
        self.assertEqual(
            review["candidates"]["left"]["reason"],
            "full-edge-near-white-band-survives-disabled-cleanup",
        )

    def test_broad_white_background_is_not_misreported_as_a_narrow_edge(self) -> None:
        image = Image.new("RGB", (120, 100), (252, 252, 252))

        review = residual_edge_review(image, max_edge_px=4)

        self.assertEqual(review, {"status": "clear", "candidates": {}})

    def test_white_seam_before_uniform_gray_mri_content_stops_at_content(self) -> None:
        image = Image.new("RGB", (120, 100), (90, 90, 90))
        draw = ImageDraw.Draw(image)
        draw.rectangle((0, 0, 119, 1), fill=(255, 255, 255))
        draw.line((0, 2, 119, 2), fill=(88, 88, 88))

        cleaned, edges = clean_panel_edges(image, max_edge_px=4)

        self.assertEqual(edges["top"], 2)
        self.assertEqual(cleaned.size, (120, 98))
        self.assertEqual(cleaned.getpixel((60, 0)), (88, 88, 88))

    def test_bright_anatomy_touching_edge_and_dark_background_are_preserved(self) -> None:
        image = Image.new("RGB", (120, 100), (5, 5, 5))
        ImageDraw.Draw(image).rectangle((42, 0, 68, 24), fill=(250, 250, 250))

        cleaned, edges = clean_panel_edges(image, max_edge_px=4)

        self.assertEqual(edges, {"top": 0, "bottom": 0, "left": 0, "right": 0})
        self.assertEqual(cleaned.tobytes(), image.tobytes())

    def test_dark_frame_bordering_brighter_content_is_preserved(self) -> None:
        image = Image.new("RGB", (120, 100), (145, 145, 145))
        ImageDraw.Draw(image).rectangle((0, 0, 2, 99), fill=(18, 18, 18))

        cleaned, edges = clean_panel_edges(image, max_edge_px=4)

        self.assertEqual(edges["left"], 0)
        self.assertEqual(cleaned.tobytes(), image.tobytes())

    def test_colored_clinical_scale_at_image_edge_is_never_classified_as_rim(self) -> None:
        image = Image.new("RGB", (120, 100), (4, 4, 4))
        draw = ImageDraw.Draw(image)
        for y in range(image.height):
            draw.line((0, y, 3, y), fill=(min(255, y * 2), 255 - min(255, y * 2), 180))

        cleaned, edges = clean_panel_edges(image, max_edge_px=4)

        self.assertEqual(edges["left"], 0)
        self.assertEqual(cleaned.tobytes(), image.tobytes())

    def test_single_figure_trimmer_removes_bounded_dark_canvas_hairline(self) -> None:
        source = self.directory / "single_figure.png"
        output = self.directory / "single_figure_trimmed.png"
        image = Image.new("RGB", (160, 120), (5, 5, 5))
        draw = ImageDraw.Draw(image)
        draw.rectangle((38, 20, 120, 97), fill=(145, 145, 145))
        draw.line((0, 116, 79, 116), fill=(44, 44, 44))
        draw.line((80, 116, 159, 116), fill=(56, 56, 56))
        draw.line((0, 117, 159, 117), fill=(105, 105, 105))
        draw.line((0, 118, 159, 118), fill=(170, 170, 170))
        draw.line((0, 119, 159, 119), fill=(230, 230, 230))
        image.save(source)

        subprocess.run([
            sys.executable,
            str(POSTPROCESS),
            "trim",
            str(source),
            str(output),
            "--asset-type",
            "figure",
            "--bg-aware",
            "off",
        ], check=True, capture_output=True, text=True)

        metadata = json.loads(Path(str(output) + ".postprocess.json").read_text())
        self.assertEqual(metadata["edge_trim_px"]["bottom"], 4)
        self.assertEqual(metadata["max_edge_px"], 4)
        self.assertEqual(metadata["safety_margin_px"], 16)
        self.assertEqual(metadata["unpadded_size_px"], [160, 116])
        self.assertEqual(metadata["padded_size_px"], [192, 148])
        with Image.open(output) as trimmed:
            self.assertEqual(trimmed.size, (192, 148))
            self.assertLess(trimmed.getpixel((0, trimmed.height - 1))[0], 20)

    def test_single_figure_default_adds_exact_safety_canvas_at_source_edges(self) -> None:
        source = self.directory / "edge_touching_figure.png"
        output = self.directory / "edge_touching_figure_trimmed.png"
        image = Image.new("RGB", (160, 120), (7, 7, 7))
        ImageDraw.Draw(image).rectangle((0, 24, 42, 96), fill=(210, 210, 210))
        image.save(source)

        subprocess.run([
            sys.executable,
            str(POSTPROCESS),
            "trim",
            str(source),
            str(output),
            "--asset-type",
            "figure",
            "--bg-aware",
            "off",
        ], check=True, capture_output=True, text=True)

        metadata = json.loads(Path(str(output) + ".postprocess.json").read_text())
        with Image.open(output) as padded:
            padded = padded.convert("RGB")
            self.assertEqual(padded.size, (192, 152))
            self.assertEqual(
                padded.crop((16, 16, 176, 136)).tobytes(),
                image.tobytes(),
            )
            self.assertTrue(all(padded.getpixel((x, 0)) == (7, 7, 7) for x in range(192)))
            self.assertTrue(all(padded.getpixel((0, y)) == (7, 7, 7) for y in range(152)))
        self.assertEqual(metadata["margin"], 16)
        self.assertEqual(metadata["safety_margin_px"], 16)
        self.assertEqual(metadata["padding_background"], [7, 7, 7])

    def test_intermediate_figure_defaults_to_zero_margin(self) -> None:
        source = self.directory / "intermediate_figure.png"
        output = self.directory / "intermediate_figure_trimmed.png"
        image = Image.new("RGB", (100, 80), (9, 9, 9))
        image.save(source)

        subprocess.run([
            sys.executable,
            str(POSTPROCESS),
            "trim",
            str(source),
            str(output),
            "--asset-type",
            "figure",
            "--intermediate",
            "--bg-aware",
            "off",
        ], check=True, capture_output=True, text=True)

        metadata = json.loads(Path(str(output) + ".postprocess.json").read_text())
        self.assertEqual(metadata["margin"], 0)
        self.assertEqual(metadata["safety_margin_px"], 0)
        with Image.open(output) as processed:
            self.assertEqual(processed.size, image.size)

    def test_labels_command_uses_same_default_safety_margin(self) -> None:
        source = self.directory / "labelled_figure.png"
        output = self.directory / "labelled_figure_trimmed.png"
        image = Image.new("RGB", (100, 90), "white")
        draw = ImageDraw.Draw(image)
        draw.rectangle((0, 0, 99, 79), outline="black", width=2)
        draw.text((44, 81), "A", fill="black")
        draw.point((50, 89), fill="black")
        image.save(source)

        subprocess.run([
            sys.executable,
            str(POSTPROCESS),
            "labels",
            str(source),
            str(output),
            "--labels",
            "A",
            "--cut-bottom-px",
            "10",
            "--asset-type",
            "figure",
            "--bg-aware",
            "off",
        ], check=True, capture_output=True, text=True)

        metadata = json.loads(Path(str(output) + ".postprocess.json").read_text())
        self.assertEqual(metadata["margin"], 16)
        self.assertEqual(metadata["safety_margin_px"], 16)
        with Image.open(output) as processed:
            self.assertEqual(processed.size, (132, 112))

    def test_table_trimmer_keeps_its_required_margin_without_panel_rim_cleanup(self) -> None:
        source = self.directory / "table.png"
        output = self.directory / "table_trimmed.png"
        image = Image.new("RGB", (160, 120), "white")
        ImageDraw.Draw(image).rectangle((35, 30, 130, 92), outline="black", width=2)
        image.save(source)

        subprocess.run([
            sys.executable,
            str(POSTPROCESS),
            "trim",
            str(source),
            str(output),
            "--asset-type",
            "table",
        ], check=True, capture_output=True, text=True)

        metadata = json.loads(Path(str(output) + ".postprocess.json").read_text())
        self.assertEqual(metadata["margin"], 16)
        self.assertEqual(metadata["safety_margin_px"], 16)
        self.assertEqual(metadata["padding_background"], [255, 255, 255])
        self.assertNotIn("edge_trim_px", metadata)
        with Image.open(output) as trimmed:
            bbox = content_bbox_for_test(trimmed)
            self.assertEqual(bbox[0], 16)
            self.assertEqual(bbox[1], 16)
            self.assertEqual(trimmed.width - bbox[2], 16)
            self.assertEqual(trimmed.height - bbox[3], 16)

    def test_legacy_recomposer_defaults_to_exact_outer_safety_margin(self) -> None:
        first = self.directory / "legacy_a.png"
        second = self.directory / "legacy_b.png"
        Image.new("RGB", (40, 30), (25, 60, 90)).save(first)
        Image.new("RGB", (40, 30), (90, 60, 25)).save(second)
        output = self.directory / "legacy_composite.png"

        subprocess.run([
            sys.executable,
            str(POSTPROCESS),
            "recompose-panels",
            str(output),
            "--inputs",
            str(first),
            str(second),
            "--cols",
            "2",
            "--fit",
            "stretch",
            "--bg",
            "#061428",
        ], check=True, capture_output=True, text=True)

        metadata = json.loads(Path(str(output) + ".postprocess.json").read_text())
        with Image.open(output) as composite:
            composite = composite.convert("RGB")
            self.assertEqual(composite.size, (130, 62))
            self.assertTrue(all(composite.getpixel((x, 0)) == (6, 20, 40) for x in range(130)))
            self.assertTrue(all(composite.getpixel((0, y)) == (6, 20, 40) for y in range(62)))
        self.assertEqual(metadata["margin"], 16)
        self.assertEqual(metadata["safety_margin_px"], 16)

    def test_split_table_rebuilds_exact_safety_margin_on_both_parts(self) -> None:
        source = self.directory / "Table_1.png"
        image = Image.new("RGB", (160, 240), "white")
        draw = ImageDraw.Draw(image)
        draw.rectangle((20, 10, 140, 230), outline="black", width=2)
        for y in (42, 80, 118, 156, 194):
            draw.line((20, y, 140, y), fill="black", width=2)
        image.save(source)
        top = self.directory / "Table_1A.png"
        bottom = self.directory / "Table_1B.png"

        subprocess.run([
            sys.executable,
            str(POSTPROCESS),
            "split-table",
            str(source),
            str(top),
            str(bottom),
            "--split-y",
            "130",
            "--repeat-header-y",
            "60",
        ], check=True, capture_output=True, text=True)

        for output in (top, bottom):
            metadata = json.loads(Path(str(output) + ".postprocess.json").read_text())
            with Image.open(output) as part:
                bbox = content_bbox_for_test(part)
                self.assertEqual(bbox[0], 16)
                self.assertEqual(bbox[1], 16)
                self.assertEqual(part.width - bbox[2], 16)
                self.assertEqual(part.height - bbox[3], 16)
            self.assertEqual(metadata["asset_type"], "table")
            self.assertEqual(metadata["safety_margin_px"], 16)
            self.assertEqual(metadata["table_safety_margin_px"], 16)
        with Image.open(top) as top_image, Image.open(bottom) as bottom_image:
            self.assertEqual(top_image.width, bottom_image.width)

    def test_split_embedded_frame_moves_only_its_own_vertical_panel_pair(self) -> None:
        source = Image.new("RGB", (360, 260), (5, 5, 5))
        draw = ImageDraw.Draw(source)
        draw.rectangle((18, 82, 72, 134), outline=(235, 235, 235), width=3)
        draw.text((37, 98), "A", fill="white")
        draw.rectangle((195, 68, 249, 118), outline=(235, 235, 235), width=3)
        draw.text((214, 84), "B", fill="white")
        draw.rectangle((92, 155, 150, 225), outline=(160, 160, 160), width=2)
        inputs = self.source_panels(source, [
            (0, 0, 180, 126),
            (180, 0, 360, 126),
            (0, 126, 180, 260),
            (180, 126, 360, 260),
        ])

        _, sidecar, _ = self.compose_inputs(inputs, "--no-trim")

        first, second, third, fourth = sidecar["panel_cleanup"]
        self.assertEqual(first["boundary_adjustments"][0]["side"], "bottom")
        self.assertEqual(third["boundary_adjustments"][0]["side"], "top")
        self.assertGreater(first["boundary_adjustments"][0]["adjusted_boundary_px"], 134)
        self.assertEqual(
            first["boundary_adjustments"][0]["effective_crop_box_px"][3],
            third["boundary_adjustments"][0]["effective_crop_box_px"][1],
        )
        self.assertEqual(second["boundary_adjustments"], [])
        self.assertEqual(fourth["boundary_adjustments"], [])
        self.assertEqual(first["label_overwritten_pixels"], 0)

    def test_full_width_lower_panel_keeps_one_boundary_for_both_upper_panels(self) -> None:
        source = Image.new("RGB", (360, 300), (5, 5, 5))
        draw = ImageDraw.Draw(source)
        draw.rectangle((0, 96, 49, 155), outline=(225, 225, 225), width=3)
        draw.text((14, 113), "A", fill="white")
        draw.rectangle((140, 94, 194, 150), outline=(225, 225, 225), width=3)
        draw.text((158, 110), "B", fill="white")
        draw.text((115, 178), "CBF", fill="white")
        inputs = self.source_panels(source, [
            (0, 0, 130, 145),
            (130, 0, 360, 145),
            (0, 145, 360, 300),
        ])

        _, sidecar, _ = self.compose_inputs(inputs, "--no-trim")

        records = [entry["boundary_adjustments"][0] for entry in sidecar["panel_cleanup"]]
        adjusted = {entry["adjusted_boundary_px"] for entry in records}
        self.assertEqual(len(adjusted), 1)
        self.assertGreater(next(iter(adjusted)), 155)
        self.assertEqual([entry["side"] for entry in records], ["bottom", "bottom", "top"])

    def test_colored_scale_crossing_seam_blocks_automatic_boundary_adjustment(self) -> None:
        source = Image.new("RGB", (180, 250), (5, 5, 5))
        draw = ImageDraw.Draw(source)
        draw.rectangle((10, 77, 63, 132), outline=(235, 235, 235), width=3)
        draw.text((27, 94), "A", fill="white")
        draw.rectangle((140, 120, 145, 155), fill=(240, 30, 70))
        inputs = self.source_panels(source, [
            (0, 0, 180, 120),
            (0, 120, 180, 250),
        ])

        _, sidecar, _ = self.compose_inputs(inputs, "--no-trim")

        self.assertEqual(sidecar["panel_cleanup"][0]["boundary_adjustments"], [])
        self.assertEqual(sidecar["panel_cleanup"][1]["boundary_adjustments"], [])

    def test_figure10_reviewed_horizontal_seam_cannot_be_shifted_for_false_label_frame(self) -> None:
        source = Image.new("RGB", (400, 620), (12, 12, 12))
        draw = ImageDraw.Draw(source)
        draw.rectangle((29, 312, 366, 597), fill=(92, 92, 92))
        # A bright rectangular acquisition feature crosses the reviewed seam
        # near the upper panel's right corner.  The legacy label-frame repair
        # mistook this for a clipped panel letter and expanded A into D.
        draw.rectangle((337, 294, 365, 332), outline=(235, 235, 235), width=3)
        inputs = self.source_panels(source, [
            (29, 25, 367, 312),
            (29, 312, 367, 598),
        ])
        for path, edge in zip(inputs, ("bottom", "top")):
            sidecar_path = Path(str(path) + ".postprocess.json")
            metadata = json.loads(sidecar_path.read_text())
            metadata["command"] = "panel-crop"
            metadata["required_seam_edges"] = [edge]
            metadata["seam_reviews"] = {edge: {"edge": edge}}
            sidecar_path.write_text(json.dumps(metadata))

        _, sidecar, _ = self.compose_inputs(
            inputs,
            "--asset-type", "clinical-image",
            "--cols", "1",
            "--no-trim",
        )

        upper, lower = sidecar["panel_cleanup"]
        self.assertEqual(upper["boundary_adjustments"], [])
        self.assertEqual(lower["boundary_adjustments"], [])
        adjacency = sidecar["source_seam_topology"]["groups"][0]["adjacencies"][0]
        self.assertEqual(adjacency["first_coordinate_px"], 312)
        self.assertEqual(adjacency["second_coordinate_px"], 312)

    def test_embedded_source_label_preserves_all_image_pixels_and_uses_no_duplicate(self) -> None:
        source = self.directory / "embedded_A.png"
        image = Image.new("RGB", (160, 120), (4, 4, 4))
        draw = ImageDraw.Draw(image)
        draw.rectangle((0, 84, 55, 119), fill=(215, 215, 215))
        draw.text((8, 91), "A", fill=(0, 0, 0))
        image.save(source)
        Path(str(source) + ".postprocess.json").write_text(json.dumps({
            "source_label_placement": "embedded",
            "embedded_label": "A",
        }))

        output, sidecar, geometry = self.compose_inputs([source], "--no-trim")

        with Image.open(output) as padded:
            padded = padded.convert("RGB")
            self.assertEqual(padded.size, (192, 152))
            self.assertEqual(padded.crop((16, 16, 176, 136)).tobytes(), image.tobytes())
        self.assertEqual(sidecar["safety_margin_px"], 16)
        self.assertEqual(sidecar["unpadded_size_px"], [160, 120])
        self.assertEqual(sidecar["padded_size_px"], [192, 152])
        self.assertEqual(sidecar["source_label_policy"], "preserve")
        self.assertFalse(sidecar["native_labels"])
        self.assertEqual(sidecar["embedded_labels"], ["A"])
        self.assertEqual(sidecar["panel_cleanup"][0]["label_overwritten_pixels"], 0)
        self.assertEqual(geometry[output.stem], [])
        self.assertEqual(sidecar["layout_candidates"][0]["band_px"], 0)

    def test_explicit_preserve_policy_works_without_extractor_metadata(self) -> None:
        source = self.directory / "legacy_embedded.png"
        Image.new("RGB", (160, 120), (6, 6, 6)).save(source)

        _, sidecar, _ = self.compose_inputs(
            [source], "--source-label-policy", "preserve", "--no-trim"
        )

        self.assertEqual(sidecar["source_label_policy"], "preserve")
        self.assertFalse(sidecar["native_labels"])

    def test_verified_exterior_label_margin_can_be_removed_without_overwriting(self) -> None:
        source = self.directory / "margin_A.png"
        image = Image.new("RGB", (160, 120), (10, 10, 10))
        draw = ImageDraw.Draw(image)
        draw.rectangle((20, 15, 135, 80), fill=(170, 170, 170))
        draw.text((9, 100), "A", fill=(255, 255, 255))
        image.save(source)
        Path(str(source) + ".postprocess.json").write_text(json.dumps({
            "source_panel_label": {
                "placement": "external-margin",
                "box_px": [6, 96, 28, 120],
            },
        }))

        output, sidecar, geometry = self.compose_inputs([source], "--no-trim")

        self.assertEqual(sidecar["source_label_policy"], "crop-safe-margin")
        self.assertTrue(sidecar["native_labels"])
        self.assertEqual(sidecar["panel_cleanup"][0]["label_action"], "cropped-exterior-margin")
        self.assertEqual(sidecar["panel_cleanup"][0]["label_margin_crop"]["side"], "bottom")
        self.assertEqual(geometry[next(iter(geometry))][0]["label"], "A")
        self.assertEqual(sidecar["native_label_color"], "#8FA8C8")
        self.assertEqual(sidecar["safety_margin_px"], 16)
        with Image.open(output) as composite:
            width, height = composite.size
        candidate = sidecar["layout_candidates"][0]
        self.assertEqual(candidate["composite_width_px"], width)
        self.assertEqual(candidate["composite_height_px"], height)
        anchor = geometry[output.stem][0]
        expected_right = 16 + candidate["unpadded_width_px"]
        panel_bottom = 16 + candidate["unpadded_height_px"] - candidate["band_px"]
        self.assertAlmostEqual(anchor["fx_right"] * width, expected_right, delta=1.0)
        self.assertAlmostEqual(
            anchor["fy_center"] * height - panel_bottom,
            (0.06 + 0.0525) / candidate["fit_in_per_px"],
            delta=1.0,
        )

    def test_native_label_stamping_uses_padded_geometry(self) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        source = self.directory / "stamp_margin_A.png"
        image = Image.new("RGB", (160, 120), (248, 248, 248))
        ImageDraw.Draw(image).rectangle((0, 0, 159, 89), fill=(10, 10, 10))
        ImageDraw.Draw(image).text((9, 99), "A", fill=(0, 0, 0))
        image.save(source)
        Path(str(source) + ".postprocess.json").write_text(json.dumps({
            "source_panel_label": {
                "placement": "external-margin",
                "box_px": [6, 96, 28, 116],
                "image_box_px": [0, 0, 160, 90],
            },
        }))
        output, _, _ = self.compose_inputs([source], "--no-trim")

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(str(output), Inches(0.75), Inches(1.0), width=Inches(10.0))
        raw_pptx = self.directory / "raw_labels.pptx"
        labelled_pptx = self.directory / "labelled.pptx"
        spec = self.directory / "label_spec.json"
        specification = {
            "slides": [{"type": "figure", "image": str(output)}],
        }
        spec.write_text(json.dumps(specification))
        presentation.save(raw_pptx)
        normalized = Presentation(raw_pptx)
        manifest = build_deck.make_build_manifest(
            specification, spec, "standard", normalized
        )
        build_deck._set_core_manifest_properties(normalized, manifest)
        normalized.save(raw_pptx)
        build_deck._embed_manifest_part(raw_pptx, manifest)

        subprocess.run([
            sys.executable,
            str(ADD_PANEL_LABELS),
            str(raw_pptx),
            str(labelled_pptx),
            "--spec",
            str(spec),
            "--geometry",
            str(self.directory / "custom_geometry.json"),
            "--label-pt",
            "18",
        ], check=True, capture_output=True, text=True)

        labelled = Presentation(labelled_pptx)
        label_shapes = [
            shape for shape in labelled.slides[0].shapes
            if getattr(shape, "has_text_frame", False) and shape.text == "A"
        ]
        self.assertEqual(len(label_shapes), 1)
        run = label_shapes[0].text_frame.paragraphs[0].runs[0]
        self.assertAlmostEqual(run.font.size.pt, 18.0)
        self.assertEqual(str(run.font.color.rgb), "8FA8C8")

    def test_source_geometry_overrides_stale_embedded_flag_for_exterior_label(self) -> None:
        source = self.directory / "stale_embedded_A.png"
        image = Image.new("RGB", (160, 120), (248, 248, 248))
        ImageDraw.Draw(image).rectangle((0, 0, 159, 89), fill=(10, 10, 10))
        ImageDraw.Draw(image).text((9, 99), "A", fill=(0, 0, 0))
        image.save(source)
        Path(str(source) + ".postprocess.json").write_text(json.dumps({
            "source_label_placement": "embedded",
            "embedded_label": "A",
            "source_label_bbox_pt": [6, 96, 28, 116],
            "source_image_bbox_pt": [0, 0, 160, 90],
            "source_crop_bbox_pt": [0, 0, 160, 120],
        }))

        _, sidecar, geometry = self.compose_inputs([source], "--no-trim")

        detection = sidecar["panel_cleanup"][0]["label_detection"]
        self.assertEqual(detection["declared_placement"], "embedded")
        self.assertEqual(detection["placement"], "external-margin")
        self.assertEqual(detection["geometry_space"], "source-pt")
        self.assertEqual(sidecar["source_label_policy"], "crop-safe-margin")
        self.assertEqual(sidecar["native_label_values"], ["A"])
        self.assertEqual(geometry[next(iter(geometry))][0]["label"], "A")

    def test_verified_image_box_then_bounded_rim_cleanup_removes_combined_white_edge(self) -> None:
        source = self.directory / "verified_frame_and_rim_A.png"
        image = Image.new("RGB", (160, 120), (248, 248, 248))
        draw = ImageDraw.Draw(image)
        # Three-pixel PDF crop frame, two-pixel raster hairline, then image data.
        draw.rectangle((3, 3, 156, 89), fill=(250, 250, 250))
        draw.rectangle((3, 5, 156, 89), fill=(12, 12, 12))
        draw.text((9, 99), "A", fill=(0, 0, 0))
        image.save(source)
        Path(str(source) + ".postprocess.json").write_text(json.dumps({
            "source_panel_label": {
                "placement": "external-margin",
                "box_px": [6, 96, 28, 116],
                "image_box_px": [3, 3, 157, 90],
            },
        }))

        output, sidecar, _ = self.compose_inputs([source])

        cleanup = sidecar["panel_cleanup"][0]
        self.assertEqual(
            cleanup["label_margin_crop"]["verified_image_box_crop_px"],
            {"top": 3, "bottom": 30, "left": 3, "right": 3},
        )
        self.assertEqual(cleanup["edge_trim_px"]["top"], 2)
        self.assertLess(Image.open(output).convert("RGB").getpixel((0, 0))[0], 20)

    def test_verified_per_edge_trim_removes_five_pixel_exterior_band(self) -> None:
        source = self.directory / "verified_bottom_band_A.png"
        image = Image.new("RGB", (160, 140), (248, 248, 248))
        draw = ImageDraw.Draw(image)
        draw.rectangle((0, 20, 159, 139), fill=(12, 12, 12))
        draw.rectangle((0, 135, 159, 139), fill=(252, 252, 252))
        draw.text((9, 3), "A", fill=(0, 0, 0))
        image.save(source)
        Path(str(source) + ".postprocess.json").write_text(json.dumps({
            "source_panel_label": {
                "placement": "external-margin",
                "box_px": [6, 2, 28, 18],
                "image_box_px": [0, 20, 160, 140],
            },
            "verified_edge_trim_px": {"top": 0, "bottom": 5, "left": 0, "right": 0},
            "verified_edge_trim_reason": "verified-pdf-exterior-band",
        }))

        _, sidecar, _ = self.compose_inputs([source])

        cleanup = sidecar["panel_cleanup"][0]
        self.assertEqual(cleanup["verified_edge_trim_px"]["bottom"], 5)
        self.assertEqual(cleanup["edge_trim_px"]["bottom"], 0)
        self.assertEqual(cleanup["total_edge_trim_px"]["bottom"], 5)
        self.assertEqual(cleanup["residual_edge_review"]["status"], "clear")

    def test_source_geometry_overrides_stale_exterior_flag_for_embedded_label(self) -> None:
        source = self.directory / "stale_exterior_A.png"
        image = Image.new("RGB", (160, 120), (7, 7, 7))
        ImageDraw.Draw(image).text((9, 94), "A", fill=(255, 255, 255))
        image.save(source)
        Path(str(source) + ".postprocess.json").write_text(json.dumps({
            "source_panel_label": {
                "placement": "external-margin",
                "box_px": [6, 90, 28, 114],
                "image_box_px": [0, 0, 160, 120],
            },
        }))

        output, sidecar, geometry = self.compose_inputs([source], "--no-trim")

        detection = sidecar["panel_cleanup"][0]["label_detection"]
        self.assertEqual(detection["placement"], "embedded")
        self.assertEqual(sidecar["source_label_policy"], "preserve")
        self.assertEqual(sidecar["embedded_labels"], ["A"])
        self.assertEqual(geometry[output.stem], [])

    def test_mixed_figure_preserves_only_embedded_and_relabels_exterior_panel(self) -> None:
        embedded = self.directory / "mixed_A.png"
        Image.new("RGB", (160, 120), (7, 7, 7)).save(embedded)
        Path(str(embedded) + ".postprocess.json").write_text(json.dumps({
            "source_panel_label": {
                "placement": "embedded",
                "box_px": [6, 90, 28, 114],
                "image_box_px": [0, 0, 160, 120],
            },
        }))

        exterior = self.directory / "mixed_B.png"
        image = Image.new("RGB", (160, 120), (248, 248, 248))
        ImageDraw.Draw(image).rectangle((0, 0, 159, 89), fill=(10, 10, 10))
        ImageDraw.Draw(image).text((9, 99), "B", fill=(0, 0, 0))
        image.save(exterior)
        Path(str(exterior) + ".postprocess.json").write_text(json.dumps({
            "source_panel_label": {
                "placement": "external-margin",
                "box_px": [6, 96, 28, 116],
                "image_box_px": [0, 0, 160, 90],
            },
        }))

        _, sidecar, geometry = self.compose_inputs([embedded, exterior], "--no-trim")

        self.assertEqual(sidecar["source_label_policy"], "mixed")
        self.assertEqual(sidecar["embedded_labels"], ["A"])
        self.assertEqual(sidecar["native_label_values"], ["B"])
        self.assertEqual([entry["label_action"] for entry in sidecar["panel_cleanup"]],
                         ["preserved", "cropped-exterior-margin"])
        self.assertEqual([entry["label"] for entry in geometry[next(iter(geometry))]], ["B"])

    def test_exterior_margin_with_image_content_falls_back_to_preservation(self) -> None:
        source = self.directory / "unsafe_margin_A.png"
        image = Image.new("RGB", (160, 120), (8, 8, 8))
        ImageDraw.Draw(image).rectangle((50, 100, 130, 119), fill=(190, 190, 190))
        image.save(source)
        Path(str(source) + ".postprocess.json").write_text(json.dumps({
            "source_panel_label": {
                "placement": "external-margin",
                "box_px": [6, 96, 28, 120],
            },
        }))

        output, sidecar, geometry = self.compose_inputs([source], "--no-trim")

        self.assertEqual(sidecar["source_label_policy"], "preserve")
        self.assertFalse(sidecar["native_labels"])
        self.assertEqual(Image.open(output).getpixel((91, 126)), (190, 190, 190))
        self.assertEqual(geometry[output.stem], [])

    def test_banded_recomposer_allows_explicit_zero_safety_margin(self) -> None:
        source = self.directory / "zero_margin_panel.png"
        image = Image.new("RGB", (160, 120), (12, 12, 12))
        image.save(source)
        Path(str(source) + ".postprocess.json").write_text(json.dumps({
            "source_label_placement": "embedded",
            "embedded_label": "A",
        }))

        output, sidecar, geometry = self.compose_inputs(
            [source],
            "--source-label-policy",
            "preserve",
            "--no-trim",
            "--asset-type",
            "clinical-image",
            "--safety-margin-px",
            "0",
        )

        with Image.open(output) as composite:
            self.assertEqual(composite.size, image.size)
        self.assertEqual(sidecar["safety_margin_px"], 0)
        self.assertEqual(sidecar["asset_type"], "clinical-image")
        self.assertEqual(geometry[output.stem], [])

    def test_verified_absent_label_evidence_adds_one_native_label(self) -> None:
        source = self.directory / "absent-source.png"
        Image.new("RGB", (180, 120), (20, 50, 90)).save(source)
        panel = self.directory / "absent-panel.png"
        crop = subprocess.run([
            sys.executable, str(POSTPROCESS), "panel-crop", str(source), str(panel),
            "--box", "0", "0", "180", "120",
            "--label", "A", "--label-placement", "absent",
        ], capture_output=True, text=True)
        self.assertEqual(crop.returncode, 0, crop.stderr + crop.stdout)

        output, sidecar, geometry = self.compose_inputs(
            [panel], "--asset-type", "clinical-image", "--no-trim"
        )

        self.assertEqual(sidecar["verified_absent_labels"], ["A"])
        self.assertEqual(sidecar["native_label_values"], ["A"])
        self.assertEqual(sidecar["embedded_labels"], [])
        self.assertEqual(
            sidecar["panel_cleanup"][0]["label_action"],
            "native-from-verified-absence",
        )
        self.assertEqual([entry["label"] for entry in geometry[output.stem]], ["A"])

        panel_sidecar_path = Path(str(panel) + ".postprocess.json")
        panel_sidecar = json.loads(panel_sidecar_path.read_text())
        panel_sidecar["source_panel_label"]["absence_evidence"][
            "decoded_rgb_sha256"
        ] = "0" * 64
        panel_sidecar_path.write_text(json.dumps(panel_sidecar))
        failed = subprocess.run([
            sys.executable, str(SCRIPT), str(self.directory / "tampered.png"),
            "--inputs", str(panel), "--labels", "A",
            "--geometry", str(self.directory / "tampered-geometry.json"),
            "--asset-type", "clinical-image", "--no-trim",
        ], capture_output=True, text=True)
        self.assertNotEqual(failed.returncode, 0)
        self.assertIn("review evidence", failed.stderr + failed.stdout)

    def test_verified_absent_label_rejects_present_label_geometry(self) -> None:
        source = self.directory / "absent-with-box-source.png"
        Image.new("RGB", (180, 120), (20, 50, 90)).save(source)
        result = subprocess.run([
            sys.executable, str(POSTPROCESS), "panel-crop", str(source),
            str(self.directory / "absent-with-box.png"),
            "--box", "0", "0", "180", "120",
            "--label", "A", "--label-placement", "absent",
            "--label-box", "2", "2", "20", "20",
            "--image-box", "0", "0", "180", "120",
        ], capture_output=True, text=True)
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("cannot have", result.stderr + result.stdout)

    def test_left_span_2x2_template_rearranges_five_embedded_panels(self) -> None:
        panels = []
        sizes = [(320, 500), (420, 300), (430, 300), (420, 290), (430, 290)]
        for index, size in enumerate(sizes):
            path = self.directory / f"span-{index}.png"
            Image.new("RGB", size, (30 + index * 25, 55, 95)).save(path)
            Path(str(path) + ".postprocess.json").write_text(json.dumps({
                "source_panel_label": {
                    "label": chr(ord("A") + index),
                    "status": "present",
                    "placement": "embedded",
                    "box_px": [size[0] - 30, size[1] - 30, size[0] - 5, size[1] - 5],
                    "image_box_px": [0, 0, size[0], size[1]],
                },
            }))
            panels.append(path)

        output, sidecar, geometry = self.compose_inputs(
            panels,
            "--asset-type", "clinical-image",
            "--layout-template", "left-span-2x2",
            "--source-label-policy", "preserve",
            "--no-trim",
        )

        boxes = sidecar["panel_boxes_px"]
        self.assertEqual(sidecar["layout_mode"], "template")
        self.assertEqual(sidecar["layout_template"], "left-span-2x2")
        self.assertEqual(sidecar["safety_margin_px"], 0)
        self.assertEqual(boxes[0]["y"], 0)
        self.assertEqual(boxes[0]["h"], sidecar["padded_size_px"][1])
        self.assertEqual(boxes[1]["y"], boxes[2]["y"])
        self.assertEqual(boxes[3]["y"], boxes[4]["y"])
        self.assertGreater(boxes[3]["y"], boxes[1]["y"])
        self.assertEqual(geometry[output.stem], [])

    def test_right_span_2x2_preserves_cross_row_source_topology_and_aspect(self) -> None:
        panels = []
        sizes = [(320, 265), (290, 265), (425, 527), (320, 262), (290, 262)]
        for index, size in enumerate(sizes):
            path = self.directory / f"right-span-{index}.png"
            Image.new("RGB", size, (30 + index * 25, 60, 100)).save(path)
            Path(str(path) + ".postprocess.json").write_text(json.dumps({
                "source_panel_label": {
                    "label": chr(ord("A") + index),
                    "status": "present",
                    "placement": "embedded",
                    "box_px": [4, size[1] - 34, 30, size[1] - 4],
                    "image_box_px": [0, 0, size[0], size[1]],
                },
            }))
            panels.append(path)

        output, sidecar, geometry = self.compose_inputs(
            panels,
            "--asset-type", "clinical-image",
            "--layout-template", "right-span-2x2",
            "--source-label-policy", "preserve",
            "--no-trim",
        )

        boxes = sidecar["panel_boxes_px"]
        canvas_width, canvas_height = sidecar["padded_size_px"]
        self.assertEqual(sidecar["layout_mode"], "template")
        self.assertEqual(sidecar["layout_template"], "right-span-2x2")
        self.assertEqual((sidecar["rows"], sidecar["cols"]), (2, 3))
        self.assertEqual(boxes[0]["y"], boxes[1]["y"])
        self.assertEqual(boxes[3]["y"], boxes[4]["y"])
        self.assertGreater(boxes[3]["y"], boxes[0]["y"])
        self.assertEqual(boxes[2]["y"], 0)
        self.assertEqual(boxes[2]["h"], canvas_height)
        self.assertEqual(boxes[2]["x"] + boxes[2]["w"], canvas_width)
        self.assertGreater(boxes[2]["x"], boxes[1]["x"])
        for source_size, box in zip(sizes, boxes):
            source_ratio = source_size[0] / source_size[1]
            output_ratio = box["w"] / box["h"]
            self.assertLess(abs(output_ratio / source_ratio - 1.0), 0.006)
        self.assertEqual(geometry[output.stem], [])

        handled, failures = image_polarity._deterministic_helper_evidence(output, sidecar)
        self.assertTrue(handled)
        self.assertFalse(failures)

    def test_two_span_right_stack_preserves_aspect_and_relative_scale(self) -> None:
        panels = []
        sizes = [(316, 510), (322, 510), (393, 265), (393, 245)]
        for index, size in enumerate(sizes):
            path = self.directory / f"two-span-{index}.png"
            Image.new("RGB", size, (35 + index * 30, 65, 100)).save(path)
            Path(str(path) + ".postprocess.json").write_text(json.dumps({
                "source_panel_label": {
                    "label": chr(ord("A") + index),
                    "status": "present",
                    "placement": "embedded",
                    "box_px": [4, size[1] - 34, 30, size[1] - 4],
                    "image_box_px": [0, 0, size[0], size[1]],
                },
            }))
            panels.append(path)

        output, sidecar, geometry = self.compose_inputs(
            panels,
            "--asset-type", "clinical-image",
            "--layout-template", "two-span-right-stack",
            "--source-label-policy", "preserve",
            "--no-trim",
        )

        boxes = sidecar["panel_boxes_px"]
        canvas_width, canvas_height = sidecar["padded_size_px"]
        self.assertEqual(sidecar["layout_mode"], "template")
        self.assertEqual(sidecar["layout_template"], "two-span-right-stack")
        self.assertEqual((sidecar["rows"], sidecar["cols"]), (2, 3))
        self.assertEqual(boxes[0]["y"], 0)
        self.assertEqual(boxes[1]["y"], 0)
        self.assertEqual(boxes[0]["h"], canvas_height)
        self.assertEqual(boxes[1]["h"], canvas_height)
        self.assertEqual(boxes[2]["x"], boxes[3]["x"])
        self.assertEqual(boxes[2]["w"], boxes[3]["w"])
        self.assertEqual(boxes[3]["y"], boxes[2]["h"] + sidecar["gap"])
        self.assertEqual(boxes[3]["y"] + boxes[3]["h"], canvas_height)
        self.assertEqual(boxes[3]["x"] + boxes[3]["w"], canvas_width)
        for source_size, box in zip(sizes, boxes):
            source_ratio = source_size[0] / source_size[1]
            output_ratio = box["w"] / box["h"]
            self.assertLess(abs(output_ratio / source_ratio - 1.0), 0.006)
        source_relative_height = sizes[3][1] / sizes[0][1]
        output_relative_height = boxes[3]["h"] / boxes[0]["h"]
        self.assertLess(abs(output_relative_height - source_relative_height), 0.05)
        self.assertLess(output_relative_height, 0.60)
        self.assertEqual(geometry[output.stem], [])

        handled, failures = image_polarity._deterministic_helper_evidence(output, sidecar)
        self.assertTrue(handled)
        self.assertFalse(failures)

    def test_panel_with_solid_corner_mask_is_rejected_against_its_source_crop(self) -> None:
        original_path = self.directory / "verified_source.png"
        image = Image.new("RGB", (160, 120), (5, 5, 5))
        ImageDraw.Draw(image).rectangle((0, 90, 42, 119), fill=(210, 210, 210))
        image.save(original_path)
        masked_path = self.directory / "masked_panel.png"
        ImageDraw.Draw(image).rectangle((0, 90, 42, 119), fill=(0, 0, 0))
        image.save(masked_path)
        Path(str(masked_path) + ".postprocess.json").write_text(json.dumps({
            "source": str(original_path),
            "crop_box_px": [0, 0, 160, 120],
            "source_label_placement": "embedded",
        }))

        result = subprocess.run([
            sys.executable,
            str(SCRIPT),
            str(self.directory / "must_not_exist.png"),
            "--inputs",
            str(masked_path),
            "--labels",
            "A",
            "--geometry",
            str(self.directory / "geometry.json"),
        ], capture_output=True, text=True)

        self.assertNotEqual(result.returncode, 0)
        self.assertIn("solid-color corner overwrite", result.stderr)
        self.assertFalse((self.directory / "must_not_exist.png").exists())


if __name__ == "__main__":
    unittest.main()
