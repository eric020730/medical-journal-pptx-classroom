#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Render the enhanced dark-academic teaching-deck style in widescreen 16:9.

Visual style ("nice"): near-black-navy page background (#061428), a 1.0" mid-navy
header band (#102444) with a kicker line above the title, a bright light-blue
divider (#5BA9E6) under the header, a muted footer (citation left, zero-padded
page number right), the Dr. Leether logo at the upper-right of non-title/
non-thanks slides, full-bleed section dividers with a large accent number, and
figures/tables on optional white cards.

Slide types: title, section, content, figure, thanks (alias: thankyou).
Slide-visible text is English; speaker notes are Traditional Chinese with
**term** rendered as real bold runs.

Usage:
    python3 build_deck.py <spec.json> --out <output.pptx>

Pair with add_panel_labels.py to stamp native fixed-size A/B/C/D labels onto the
finished deck (see SKILL.md "Native Post-Build Panel Labels"). The figure image
box used by this builder is 12.13 x 4.95 in (pass these to
recompose_panels_banded.py as --slide-box-w-in/--slide-box-h-in).
"""
import json, os, re, argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

DEFAULT_LOGO_PATH = Path(__file__).resolve().parents[1] / "assets" / "dr_leether_logo.png"

BG      = RGBColor(0x06, 0x14, 0x28)
HEADER  = RGBColor(0x10, 0x24, 0x44)
DIVIDER = RGBColor(0x5B, 0xA9, 0xE6)
TITLE_C = RGBColor(0xEA, 0xF2, 0xFC)
BODY_C  = RGBColor(0xD8, 0xE4, 0xF2)
SECON_C = RGBColor(0x8F, 0xA8, 0xC8)
ACCENT  = RGBColor(0x5B, 0xA9, 0xE6)
CARD    = RGBColor(0xFF, 0xFF, 0xFF)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FIG_BOX = (Inches(0.6), Inches(1.25), Inches(12.13), Inches(4.95))


def rect(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(1, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def parse_runs(line):
    out = []
    m = re.match(r'^(\s*[•\-]\s*)([A-Za-z][^:]{1,40}:)(.*)$', line)
    if m:
        out.append((m.group(1), False)); out.append((m.group(2), True)); rest = m.group(3)
    else:
        rest = line
    for part in re.split(r'(\*\*[^*]+\*\*)', rest):
        if not part: continue
        if part.startswith("**") and part.endswith("**"):
            out.append((part[2:-2], True))
        else:
            out.append((part, False))
    return out or [(line, False)]


def parse_runs_notes(line):
    out = []
    for part in re.split(r'(\*\*[^*]+\*\*)', line):
        if not part: continue
        if part.startswith("**") and part.endswith("**"):
            out.append((part[2:-2], True))
        else:
            out.append((part, False))
    return out or [(line, False)]


def txt(slide, x, y, w, h, lines, size=18, color=BODY_C, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
        line_spacing=1.0, space_after=4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(lines, str): lines = [lines]
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = align; p.line_spacing = line_spacing
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        for seg, b in parse_runs(ln):
            r = p.add_run(); r.text = seg
            r.font.size = Pt(size); r.font.bold = b or bold
            r.font.name = font; r.font.color.rgb = color
    return tb


def set_notes(slide, notes):
    if not notes: return
    tf = slide.notes_slide.notes_text_frame; tf.clear()
    first = True
    for ln in notes.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        for seg, b in parse_runs_notes(ln):
            r = p.add_run(); r.text = seg; r.font.bold = b; r.font.size = Pt(12)


def _emf_aspect(path):
    """Read an EMF's true width/height aspect from its header rclBounds."""
    try:
        import struct
        with open(path, "rb") as f:
            head = f.read(40)
        x0, y0, x1, y1 = struct.unpack_from("<4i", head, 8)
        w, h = (x1 - x0), (y1 - y0)
        if w > 0 and h > 0:
            return w / h
    except Exception:
        pass
    return None


class Builder:
    def __init__(self, spec, base, logo):
        self.spec = spec; self.meta = spec.get("meta", {})
        self.base = base; self.logo = logo
        self.prs = Presentation()
        self.prs.slide_width = EMU_W; self.prs.slide_height = EMU_H
        self.blank = self.prs.slide_layouts[6]

    def bg(self, slide, color):
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

    def add_logo(self, slide):
        if not self.logo: return
        sz = Inches(0.62)
        slide.shapes.add_picture(self.logo, EMU_W - sz - Inches(0.35), Inches(0.19), sz, sz)

    def footer(self, slide, n):
        txt(slide, Inches(0.55), Inches(7.05), Inches(10.5), Inches(0.35),
            self.meta.get("footer_label", ""), size=10, color=SECON_C)
        txt(slide, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35),
            f"{n:02d}", size=10, color=SECON_C, align=PP_ALIGN.RIGHT)

    def header(self, slide, title, kicker=None):
        rect(slide, 0, 0, EMU_W, Inches(1.0), HEADER)
        rect(slide, 0, Inches(1.0), EMU_W, Pt(3), DIVIDER)
        if kicker:
            txt(slide, Inches(0.55), Inches(0.11), Inches(11), Inches(0.28),
                kicker, size=11.5, color=ACCENT, bold=True)
            txt(slide, Inches(0.55), Inches(0.40), Inches(11.0), Inches(0.56),
                title, size=23, color=TITLE_C, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        else:
            txt(slide, Inches(0.55), Inches(0.05), Inches(11.0), Inches(0.9),
                title, size=25, color=TITLE_C, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        self.add_logo(slide)

    def title(self, s):
        sl = self.prs.slides.add_slide(self.blank); self.bg(sl, BG)
        rect(sl, 0, Inches(2.05), EMU_W, Pt(3), DIVIDER)
        rect(sl, 0, Inches(5.0), EMU_W, Pt(3), DIVIDER)
        if s.get("kicker"):
            txt(sl, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.4),
                s["kicker"], size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, Inches(0.9), Inches(2.25), Inches(11.5), Inches(2.6),
            s["title"], size=34, color=TITLE_C, bold=True, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        txt(sl, Inches(0.9), Inches(5.25), Inches(11.5), Inches(0.8),
            s.get("authors", ""), size=15, color=BODY_C, align=PP_ALIGN.CENTER)
        txt(sl, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.6),
            s.get("citation", ""), size=13, color=SECON_C, align=PP_ALIGN.CENTER)
        set_notes(sl, s.get("notes"))

    def section(self, s, n):
        sl = self.prs.slides.add_slide(self.blank); self.bg(sl, HEADER)
        rect(sl, 0, Inches(3.05), EMU_W, Pt(3), DIVIDER)
        txt(sl, Inches(1.0), Inches(2.0), Inches(2.6), Inches(1.0),
            str(s.get("number", "")), size=46, color=ACCENT, bold=True)
        txt(sl, Inches(1.0), Inches(3.25), Inches(11.3), Inches(1.6),
            s["title"], size=32, color=TITLE_C, bold=True, line_spacing=1.05)
        if s.get("subtitle"):
            txt(sl, Inches(1.0), Inches(4.6), Inches(11.0), Inches(0.8),
                s["subtitle"], size=16, color=SECON_C)
        self.add_logo(sl); self.footer(sl, n); set_notes(sl, s.get("notes"))

    def content(self, s, n):
        sl = self.prs.slides.add_slide(self.blank); self.bg(sl, BG)
        self.header(sl, s["title"], s.get("kicker"))
        image = s.get("image")
        body_width = Inches(7.15) if image else Inches(12.0)
        txt(sl, Inches(0.7), Inches(1.35), body_width, Inches(5.4),
            s.get("body", s.get("items", s.get("bullets", []))),
            size=s.get("body_size", 18), color=BODY_C,
            line_spacing=1.12, space_after=7)
        if image:
            if not os.path.isabs(image):
                image = os.path.join(self.base, image)
            bx, by, bw, bh = (
                Inches(8.20), Inches(1.40), Inches(4.50), Inches(5.25)
            )
            is_vector = str(image).lower().endswith((".emf", ".wmf"))
            if is_vector:
                aspect = _emf_aspect(image) or float(s.get("image_aspect") or 1.5)
            else:
                with Image.open(image) as source:
                    aspect = source.width / source.height
            if bw / bh > aspect:
                draw_h = bh; draw_w = Emu(int(bh * aspect))
            else:
                draw_w = bw; draw_h = Emu(int(bw / aspect))
            draw_x = bx + Emu(int((bw - draw_w) / 2))
            draw_y = by + Emu(int((bh - draw_h) / 2))
            if s.get("card", False) or is_vector:
                pad = Inches(0.08)
                rect(sl, draw_x - pad, draw_y - pad,
                     draw_w + 2 * pad, draw_h + 2 * pad, CARD)
            sl.shapes.add_picture(image, draw_x, draw_y, draw_w, draw_h)
        self.footer(sl, n); set_notes(sl, s.get("notes"))

    def figure(self, s, n):
        sl = self.prs.slides.add_slide(self.blank); self.bg(sl, BG)
        self.header(sl, s["title"], s.get("kicker"))
        img = s["image"]
        if not os.path.isabs(img): img = os.path.join(self.base, img)
        panel_labels = s.get("panel_labels") or []
        x_fracs = s.get("panel_label_x_fracs") or []
        panel_boxes = s.get("panel_boxes") or []
        if len(panel_labels) > 1 and not s.get("panel_geometry_exception"):
            if len(x_fracs) < len(panel_labels) and len(panel_boxes) < len(panel_labels):
                raise RuntimeError(
                    "Figure slide %r has %d panel_labels but no panel geometry. "
                    "Provide 'panel_label_x_fracs' or 'panel_boxes', or set "
                    "'panel_geometry_exception': true with a documented reason."
                    % (s.get("title", ""), len(panel_labels))
                )
        bx, by, bw, bh = FIG_BOX
        if panel_labels:
            bh = Inches(4.80)
        bwi, bhi = Emu(int(bw)), Emu(int(bh))
        is_vector = str(img).lower().endswith((".emf", ".wmf"))
        if is_vector:
            # Vector tables (EMF/WMF): PIL cannot open them. Read the true aspect
            # from the EMF header; fall back to the spec's image_aspect.
            ar = _emf_aspect(img) or float(s.get("image_aspect") or 1.5)
        else:
            iw, ih = Image.open(img).size; ar = iw / ih
        if s.get("image_width_in"):
            # Fixed on-screen width: makes split-table parts (Table_1A/1B) render
            # at the SAME width regardless of pixel height. Clamp into the box.
            dw = Emu(int(Inches(float(s["image_width_in"]))))
            dh = Emu(int(dw / ar))
            if dh > bhi:
                dh = bhi; dw = Emu(int(bhi * ar))
        elif bwi / bhi > ar:
            dh = bhi; dw = Emu(int(bhi * ar))
        else:
            dw = bwi; dh = Emu(int(bwi / ar))
        px = bx + Emu(int((bwi - dw) / 2)); py = by + Emu(int((bhi - dh) / 2))
        # A vector table has a transparent/black-text background, so on the dark
        # slide it always needs a white card behind it (like the base builder).
        if s.get("card", False) or is_vector:
            pad = Inches(0.10) if is_vector else Inches(0.08)
            rect(sl, px - pad, py - pad, dw + 2 * pad, dh + 2 * pad, CARD)
        sl.shapes.add_picture(img, px, py, dw, dh)
        if panel_labels:
            label_y = min(py + dh + Inches(0.03), Inches(6.10))
            label_w, label_h = Inches(0.26), Inches(0.22)
            for index, label in enumerate(panel_labels):
                if index < len(x_fracs):
                    right_frac = float(x_fracs[index])
                elif index < len(panel_boxes) and isinstance(panel_boxes[index], dict):
                    box = panel_boxes[index]
                    right_frac = float(box.get("right_x_frac", (index + 1) / len(panel_labels)))
                else:
                    right_frac = (index + 1) / len(panel_labels)
                label_x = px + Emu(int(dw * right_frac)) - label_w
                txt(sl, label_x, label_y, label_w, label_h, str(label), size=12,
                    color=SECON_C, bold=True, align=PP_ALIGN.RIGHT,
                    anchor=MSO_ANCHOR.MIDDLE, space_after=0)
        if s.get("caption"):
            caption_y = Inches(6.40) if panel_labels else Inches(6.35)
            txt(sl, Inches(0.6), caption_y, Inches(12.1), Inches(0.55),
                s["caption"], size=11.5, color=SECON_C, align=PP_ALIGN.CENTER)
        self.footer(sl, n); set_notes(sl, s.get("notes"))

    def thanks(self, s):
        sl = self.prs.slides.add_slide(self.blank); self.bg(sl, HEADER)
        rect(sl, 0, Inches(3.5), EMU_W, Pt(3), DIVIDER)
        txt(sl, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.3),
            s.get("title", "Thank You"), size=44, color=TITLE_C, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(sl, Inches(0.9), Inches(3.8), Inches(11.5), Inches(1.2),
            s.get("subtitle", ""), size=16, color=SECON_C, align=PP_ALIGN.CENTER)
        if s.get("citation"):
            txt(sl, Inches(0.9), Inches(5.05), Inches(11.5), Inches(0.65),
                s["citation"], size=13, color=SECON_C, align=PP_ALIGN.CENTER)
        set_notes(sl, s.get("notes"))

    def build(self):
        n = 0
        for s in self.spec["slides"]:
            t = s["type"]
            if t == "title":
                self.title(s); n += 1
            elif t in ("section", "part"):
                n += 1; self.section(s, n)
            elif t in ("content", "outline", "references"):
                n += 1; self.content(s, n)
            elif t == "figure":
                n += 1; self.figure(s, n)
            elif t in ("thanks", "thankyou"):
                n += 1; self.thanks(s)
            else:
                raise ValueError(f"unknown slide type {t!r}")
        return self.prs


def resolve_logo(meta, spec_dir):
    lp = meta.get("logo_path")
    if lp:
        if not os.path.isabs(lp): lp = str((Path(spec_dir) / lp).resolve())
        if os.path.exists(lp): return lp
        print(f"WARNING: meta.logo_path {meta.get('logo_path')!r} not found; using bundled logo.")
    if DEFAULT_LOGO_PATH.exists(): return str(DEFAULT_LOGO_PATH)
    print("WARNING: no logo available; slides render without a logo.")
    return None


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("spec")
    ap.add_argument("--out", required=True)
    a = ap.parse_args()
    spec = json.load(open(a.spec, encoding="utf-8"))
    base = os.path.dirname(os.path.abspath(a.spec))
    logo = resolve_logo(spec.get("meta", {}), base)
    prs = Builder(spec, base, logo).build()
    prs.save(a.out)
    print(f"saved {a.out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
