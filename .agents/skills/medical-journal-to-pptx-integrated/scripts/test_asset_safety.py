#!/usr/bin/env python3
"""Negative regression tests for destructive raster and panel operations."""
from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from PIL import Image, ImageDraw

import postprocess_assets
import build_deck


SCRIPTS = Path(__file__).resolve().parent
POSTPROCESS = SCRIPTS / "postprocess_assets.py"
BANDED = SCRIPTS / "recompose_panels_banded.py"
ALIGNED = SCRIPTS / "recompose_panels_aligned.py"
ADD_LABELS = SCRIPTS / "add_panel_labels.py"


class AssetSafetyTests(unittest.TestCase):
    def setUp(self) -> None:
        self.tmp = tempfile.TemporaryDirectory()
        self.directory = Path(self.tmp.name)

    def tearDown(self) -> None:
        self.tmp.cleanup()

    def run_script(self, script: Path, *args: object) -> subprocess.CompletedProcess[str]:
        return subprocess.run(
            [sys.executable, str(script), *(str(arg) for arg in args)],
            capture_output=True,
            text=True,
        )

    def dark_clinical_image(self, name: str = "clinical.png") -> Path:
        path = self.directory / name
        image = Image.new("RGB", (200, 160), (8, 8, 8))
        draw = ImageDraw.Draw(image)
        draw.ellipse((35, 25, 165, 145), fill=(90, 90, 90))
        draw.rectangle((185, 40, 199, 120), fill=(210, 30, 30))
        image.save(path)
        return path

    def test_figure_auto_preserves_dark_clinical_canvas(self) -> None:
        source = self.dark_clinical_image()
        output = self.directory / "final.png"
        result = self.run_script(
            POSTPROCESS, "trim", source, output, "--asset-type", "figure"
        )
        self.assertEqual(result.returncode, 0, result.stderr)
        with Image.open(output) as image:
            self.assertEqual(image.size, (232, 192))
            self.assertEqual(image.crop((16, 16, 216, 176)).tobytes(), Image.open(source).tobytes())
        sidecar = json.loads(Path(str(output) + ".postprocess.json").read_text())
        self.assertFalse(sidecar["bg_aware_applied"])
        self.assertEqual(sidecar["safety_margin_px"], 16)

    def test_labels_requires_explicit_cut(self) -> None:
        source = self.dark_clinical_image()
        output = self.directory / "must_not_exist.png"
        result = self.run_script(
            POSTPROCESS, "labels", source, output, "--labels", "A"
        )
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("explicit", result.stderr + result.stdout)
        self.assertFalse(output.exists())

    def test_in_place_trim_is_rejected_without_changing_source(self) -> None:
        source = self.dark_clinical_image()
        before = source.read_bytes()
        result = self.run_script(POSTPROCESS, "trim", source, source)
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("in-place", result.stderr + result.stdout)
        self.assertEqual(source.read_bytes(), before)

    def test_split_table_rejects_out_of_range_final_margin(self) -> None:
        source = self.directory / "table.png"
        Image.new("RGB", (200, 200), "white").save(source)
        result = self.run_script(
            POSTPROCESS,
            "split-table",
            source,
            self.directory / "a.png",
            self.directory / "b.png",
            "--split-y", "120",
            "--repeat-header-y", "40",
            "--margin", "0",
        )
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("8-24", result.stderr + result.stdout)

    def test_final_sidecar_requires_typed_fields_and_real_pixel_canvas(self) -> None:
        image = self.directory / "fake-table.png"
        Image.new("RGB", (128, 128), "black").save(image)
        sidecar = {
            "command": "trim",
            "source": "source.png",
            "asset_type": "table",
            "margin": "0",
            "safety_margin_px": 16,
            "table_safety_margin_px": 16,
            "padding_background": "#FFFFFF",
            "unpadded_size_px": [96, 96],
            "padded_size_px": [128, 128],
            "intermediate": "true",
        }

        failures = postprocess_assets.validate_final_sidecar(image, sidecar)

        self.assertTrue(any("intermediate must be" in failure for failure in failures))
        self.assertTrue(any("margin must be an integer" in failure for failure in failures))
        self.assertTrue(any("actual outer pixels" in failure for failure in failures))

    def test_final_sidecar_accepts_a_physical_white_safety_canvas(self) -> None:
        image = self.directory / "safe-table.png"
        core = Image.new("RGB", (96, 96), (40, 60, 80))
        canvas = Image.new("RGB", (128, 128), "white")
        canvas.paste(core, (16, 16))
        canvas.save(image)
        sidecar = {
            "command": "trim",
            "source": "source.png",
            "asset_type": "table",
            "margin": 16,
            "safety_margin_px": 16,
            "table_safety_margin_px": 16,
            "padding_background": [255, 255, 255],
            "unpadded_size_px": [96, 96],
            "padded_size_px": [128, 128],
            "intermediate": False,
        }

        self.assertEqual(postprocess_assets.validate_final_sidecar(image, sidecar), [])

    def test_final_sidecar_rejects_malformed_dimensions_blank_core_and_alpha(self) -> None:
        base = {
            "command": "trim", "source": "source.png", "asset_type": "table",
            "margin": 16, "safety_margin_px": 16, "table_safety_margin_px": 16,
            "padding_background": "#FfFfFf", "unpadded_size_px": [96, 96],
            "padded_size_px": [128, 128], "intermediate": False,
        }

        malformed = self.directory / "malformed.png"
        Image.new("RGB", (128, 128), "white").save(malformed)
        malformed_sidecar = {**base, "unpadded_size_px": ["96", "96"]}
        malformed_failures = postprocess_assets.validate_final_sidecar(
            malformed, malformed_sidecar
        )
        self.assertTrue(any("two positive integers" in item for item in malformed_failures))
        self.assertTrue(any("no visible content" in item for item in malformed_failures))

        transparent = self.directory / "transparent.png"
        rgba = Image.new("RGBA", (128, 128), (255, 255, 255, 0))
        ImageDraw.Draw(rgba).rectangle((16, 16, 111, 111), fill=(40, 60, 80, 255))
        rgba.save(transparent)
        alpha_failures = postprocess_assets.validate_final_sidecar(transparent, base)
        self.assertTrue(any("actual outer pixels" in item for item in alpha_failures))

    def test_final_sidecar_accepts_lossy_white_safety_canvas_with_tolerance(self) -> None:
        image = self.directory / "safe-table.jpg"
        canvas = Image.new("RGB", (128, 128), "white")
        ImageDraw.Draw(canvas).rectangle((16, 16, 111, 111), fill=(35, 75, 120))
        canvas.save(image, quality=82)
        sidecar = {
            "command": "trim", "source": "source.png", "asset_type": "table",
            "margin": 16, "safety_margin_px": 16, "table_safety_margin_px": 16,
            "padding_background": "#fFfFfF", "unpadded_size_px": [96, 96],
            "padded_size_px": [128, 128], "intermediate": False,
        }

        self.assertEqual(postprocess_assets.validate_final_sidecar(image, sidecar), [])

    def test_notes_audit_rejects_simplified_and_repeated_character_padding(self) -> None:
        spec = self.directory / "notes.json"
        slides = [
            {"type": "title", "notes": "📌 " + "简" * 24}
            for _ in range(40)
        ]
        spec.write_text(json.dumps({"slides": slides}, ensure_ascii=False))

        result = self.run_script(POSTPROCESS, "notes-audit", "--spec", spec)

        self.assertNotEqual(result.returncode, 0)
        self.assertIn("Simplified", result.stdout + result.stderr)
        self.assertIn("repetitive", result.stdout + result.stderr)
        self.assertIsNotNone(
            postprocess_assets._NOTE_SIMPLIFIED_ONLY_RE.search(
                "💡个别患者关系风险管理资料"
            )
        )

    def test_notes_audit_rejects_non_string_notes_without_traceback(self) -> None:
        spec = self.directory / "malformed-notes.json"
        spec.write_text(json.dumps({
            "slides": [{"type": "title", "notes": ["not", "a", "string"]}]
        }))

        result = self.run_script(POSTPROCESS, "notes-audit", "--spec", spec)

        self.assertNotEqual(result.returncode, 0)
        self.assertIn("notes must be a string", result.stdout + result.stderr)
        self.assertNotIn("Traceback", result.stdout + result.stderr)

    def test_notes_audit_requires_every_visible_panel_to_be_described(self) -> None:
        spec = self.directory / "panels.json"
        good = "📌 這張投影片完整說明研究設計結果與臨床判斷重點，幫助教學討論。"
        slides = [{"type": "title", "notes": good} for _ in range(39)]
        slides.append({
            "type": "figure",
            "panel_labels": ["A", "B"],
            "notes": "🖼️ 【A: 研究組】完整說明研究結果與臨床判斷重點，但尚未解釋另一面板。",
        })
        spec.write_text(json.dumps({"slides": slides}, ensure_ascii=False))

        result = self.run_script(POSTPROCESS, "notes-audit", "--spec", spec)

        self.assertNotEqual(result.returncode, 0)
        self.assertIn("do not describe visible panel", result.stdout + result.stderr)

    def test_aligned_recomposer_rejects_label_panel_mismatch(self) -> None:
        a = self.dark_clinical_image("a.png")
        b = self.dark_clinical_image("b.png")
        result = self.run_script(
            ALIGNED,
            self.directory / "aligned.png",
            "--inputs", a, b,
            "--labels", "A",
        )
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("label count", result.stderr + result.stdout)

    def test_banded_recomposer_rejects_label_panel_mismatch(self) -> None:
        a = self.dark_clinical_image("a.png")
        b = self.dark_clinical_image("b.png")
        result = self.run_script(
            BANDED,
            self.directory / "banded.png",
            "--inputs", a, b,
            "--labels", "A",
        )
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("label count", result.stderr + result.stdout)

    def test_banded_auto_preserves_unknown_source_label_state(self) -> None:
        panel = self.dark_clinical_image("unknown.png")
        geometry = self.directory / "geometry.json"
        output = self.directory / "banded.png"
        result = self.run_script(
            BANDED,
            output,
            "--inputs", panel,
            "--labels", "A",
            "--geometry", geometry,
            "--no-trim",
        )
        self.assertEqual(result.returncode, 0, result.stderr)
        sidecar = json.loads(Path(str(output) + ".postprocess.json").read_text())
        self.assertEqual(sidecar["source_label_policy"], "preserve")
        self.assertFalse(sidecar["native_labels"])
        self.assertEqual(json.loads(geometry.read_text())[output.stem], [])


class NativePanelLabelTests(unittest.TestCase):
    def setUp(self) -> None:
        try:
            from pptx import Presentation  # noqa: F401
        except ImportError as error:
            self.skipTest(f"python-pptx unavailable: {error}")
        self.tmp = tempfile.TemporaryDirectory()
        self.directory = Path(self.tmp.name)

    def tearDown(self) -> None:
        if hasattr(self, "tmp"):
            self.tmp.cleanup()

    def make_deck(self) -> tuple[Path, Path, Path]:
        from pptx import Presentation
        from pptx.util import Inches

        image = self.directory / "figure.png"
        Image.new("RGB", (400, 240), (20, 20, 20)).save(image)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image), Inches(1), Inches(1), width=Inches(10))
        deck = self.directory / "deck.pptx"
        spec = self.directory / "spec.json"
        specification = {"slides": [{"type": "figure", "image": str(image)}]}
        spec.write_text(json.dumps(specification))
        prs.save(deck)
        normalized = Presentation(deck)
        manifest = build_deck.make_build_manifest(
            specification, spec, "standard", normalized
        )
        build_deck._set_core_manifest_properties(normalized, manifest)
        normalized.save(deck)
        build_deck._embed_manifest_part(deck, manifest)
        geometry = self.directory / "geometry.json"
        return deck, spec, geometry

    def run_labels(self, source: Path, output: Path, spec: Path, geometry: Path):
        return subprocess.run(
            [
                sys.executable, str(ADD_LABELS), str(source), str(output),
                "--spec", str(spec), "--geometry", str(geometry),
            ],
            capture_output=True,
            text=True,
        )

    def test_non_finite_geometry_is_rejected(self) -> None:
        deck, spec, geometry = self.make_deck()
        geometry.write_text('{"figure": [{"label": "A", "fx_right": NaN, "fy_center": 0.8}]}')
        result = self.run_labels(deck, self.directory / "out.pptx", spec, geometry)
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("finite", result.stderr + result.stdout)

    def test_out_of_range_geometry_is_rejected(self) -> None:
        deck, spec, geometry = self.make_deck()
        geometry.write_text(json.dumps({
            "figure": [{"label": "A", "fx_right": 1.2, "fy_center": 0.8}]
        }))
        result = self.run_labels(deck, self.directory / "out.pptx", spec, geometry)
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("[0,1]", result.stderr + result.stdout)

    def test_repeated_stamping_is_idempotent(self) -> None:
        from pptx import Presentation

        deck, spec, geometry = self.make_deck()
        geometry.write_text(json.dumps({
            "figure": [{"label": "A", "fx_right": 0.8, "fy_center": 0.8}]
        }))
        first = self.directory / "first.pptx"
        second = self.directory / "second.pptx"
        one = self.run_labels(deck, first, spec, geometry)
        self.assertEqual(one.returncode, 0, one.stderr)
        two = self.run_labels(first, second, spec, geometry)
        self.assertEqual(two.returncode, 0, two.stderr)
        prs = Presentation(second)
        labels = [
            shape for shape in prs.slides[0].shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip() == "A"
        ]
        self.assertEqual(len(labels), 1)
        self.assertTrue(labels[0].name.startswith("MJ_PANEL_LABEL_"))


if __name__ == "__main__":
    unittest.main()
