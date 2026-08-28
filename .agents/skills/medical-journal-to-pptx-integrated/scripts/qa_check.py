#!/usr/bin/env python3
"""Validate portable medical-journal specifications and finished PowerPoints."""

from __future__ import annotations

import argparse
import hashlib
import json
import math
import re
import sys
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

import build_deck
import image_polarity
import vector_table
from notes_quality import (
    CJK_RE,
    SIMPLIFIED_ONLY_RE,
    duplicate_note_failures,
    has_closing_takeaway,
    has_structural_emoji,
    note_diversity_failure,
)
from postprocess_assets import validate_final_sidecar


SKILL_ROOT = Path(__file__).resolve().parents[1]
CONFIG = {
    "modes": {
        "full": {"minimum_slides": 40, "maximum_slides": 55},
    }
}
LOGO_PATH = SKILL_ROOT / "assets" / "dr_leether_logo.png"
CHINESE_RE = CJK_RE
FIGURE_RE = re.compile(r"\bfig(?:ure)?\.?\s*(\d+)\b", re.IGNORECASE)
PANEL_FRAGMENT_RE = re.compile(r"(?:^|[_-])panel[_-]?[a-z0-9]+$", re.IGNORECASE)
RAW_STREAM_RE = re.compile(r"^image_p\d+_\d+$", re.IGNORECASE)
SPLIT_TABLE_RE = re.compile(r"^(table[_-]?\d+)[_-]?([a-z])$", re.IGNORECASE)
KEYCAP_RE = re.compile(r"[0-9]\ufe0f?\u20e3")
TABLE_MARGIN_MIN = 8
POSTPROCESS_SUFFIX = ".postprocess.json"
EMOJI_RANGES = ((0x1F300, 0x1FAFF), (0x2600, 0x27BF), (0x2B00, 0x2BFF))
ALLOWED_SLIDE_TYPES = {
    "title", "outline", "part", "content", "figure", "references", "thanks"
}
FOOTER_RE = re.compile(
    r"^\s*[^|]+\bet\s+al\b\s*[—-]\s*[^|]*\b(?:19|20)\d{2}\b\s*\|\s*\S.+$",
    re.IGNORECASE,
)
RASTER_SUFFIXES = build_deck.RASTER_SUFFIXES
MIN_NOTE_CJK = {
    "title": 12,
    "outline": 20,
    "part": 12,
    "section": 12,
    "content": 20,
    "figure": 20,
    "references": 16,
    "thanks": 12,
}


def has_emoji(text: str) -> bool:
    return has_structural_emoji(text)


def _cjk_count(text: str) -> int:
    return len(CHINESE_RE.findall(text))


def _note_diversity_failure(text: str) -> str | None:
    return note_diversity_failure(text)


def _hex_rgb(value: Any) -> tuple[int, int, int] | None:
    if not isinstance(value, str) or not re.fullmatch(r"#?[0-9A-Fa-f]{6}", value.strip()):
        return None
    raw = value.strip().lstrip("#")
    return tuple(int(raw[index:index + 2], 16) for index in (0, 2, 4))


def _contrast_ratio(foreground: tuple[int, int, int], background: tuple[int, int, int]) -> float:
    def luminance(rgb: tuple[int, int, int]) -> float:
        values: list[float] = []
        for channel in rgb:
            normalized = channel / 255.0
            values.append(
                normalized / 12.92
                if normalized <= 0.04045
                else ((normalized + 0.055) / 1.055) ** 2.4
            )
        return 0.2126 * values[0] + 0.7152 * values[1] + 0.0722 * values[2]

    bright, dark = sorted((luminance(foreground), luminance(background)), reverse=True)
    return (bright + 0.05) / (dark + 0.05)


def _nonempty_lines(value: Any) -> list[str]:
    if isinstance(value, list):
        return [str(item).strip() for item in value if str(item).strip()]
    if isinstance(value, str):
        return [line.strip() for line in value.splitlines() if line.strip()]
    return []


def _require_nonempty_string(
    slide: dict[str, Any], field: str, index: int, slide_type: str, failures: list[str]
) -> None:
    if not isinstance(slide.get(field), str) or not slide[field].strip():
        failures.append(f"Slide {index} ({slide_type}) requires a non-empty {field} string.")


def _validate_slide_schema(
    slide: dict[str, Any], index: int, failures: list[str]
) -> str:
    slide_type = slide.get("type")
    if not isinstance(slide_type, str) or slide_type not in ALLOWED_SLIDE_TYPES:
        failures.append(
            f"Slide {index} has unsupported type {slide_type!r}; allowed types are "
            + ", ".join(sorted(ALLOWED_SLIDE_TYPES))
            + "."
        )
        return str(slide_type or "missing")

    _require_nonempty_string(slide, "title", index, slide_type, failures)
    if not isinstance(slide.get("notes"), str):
        failures.append(f"Slide {index} ({slide_type}) speaker notes must be a string.")
    if slide_type == "title":
        for field in ("authors", "citation"):
            _require_nonempty_string(slide, field, index, slide_type, failures)
    elif slide_type == "outline":
        items = _nonempty_lines(slide.get("items"))
        if len(items) < 3:
            failures.append(f"Slide {index} (outline) requires at least 3 substantive items.")
    elif slide_type == "part":
        number = slide.get("number")
        if not isinstance(number, (int, str)) or not str(number).strip():
            failures.append(f"Slide {index} ({slide_type}) requires a non-empty section number.")
    elif slide_type == "content":
        body = _nonempty_lines(slide.get("body"))
        if len(body) < 3:
            failures.append(
                f"Slide {index} (content) requires at least 3 substantive body lines."
            )
        if sum(len(line) for line in body) < 45:
            failures.append(
                f"Slide {index} (content) body is too short to be a substantive teaching slide."
            )
    elif slide_type == "figure":
        _require_nonempty_string(slide, "image", index, slide_type, failures)
        _require_nonempty_string(slide, "caption", index, slide_type, failures)
        labels = slide.get("panel_labels", [])
        if not isinstance(labels, list):
            failures.append(f"Slide {index} (figure) panel_labels must be a list.")
            labels = []
        elif any(not isinstance(label, (str, int)) or not str(label).strip() for label in labels):
            failures.append(f"Slide {index} (figure) panel_labels must be non-empty strings/numbers.")
        normalized_labels = [str(label).strip().upper() for label in labels]
        if len(set(normalized_labels)) != len(normalized_labels):
            failures.append(f"Slide {index} (figure) panel_labels must be unique.")

        fractions = slide.get("panel_label_x_fracs")
        boxes = slide.get("panel_boxes")
        if fractions is not None and boxes is not None:
            failures.append(
                f"Slide {index} (figure) must use either panel_label_x_fracs or panel_boxes, not both."
            )
        anchors: list[float] = []
        if fractions is not None:
            if not isinstance(fractions, list) or len(fractions) != len(labels):
                failures.append(
                    f"Slide {index} (figure) panel_label_x_fracs must match panel_labels exactly."
                )
            else:
                for value in fractions:
                    try:
                        anchors.append(float(value))
                    except (TypeError, ValueError):
                        anchors.append(float("nan"))
        if boxes is not None:
            if not isinstance(boxes, list) or len(boxes) != len(labels):
                failures.append(
                    f"Slide {index} (figure) panel_boxes must match panel_labels exactly."
                )
            else:
                box_anchors: list[float] = []
                for box in boxes:
                    try:
                        if not isinstance(box, dict):
                            raise TypeError
                        box_anchors.append(float(box["right_x_frac"]))
                    except (KeyError, TypeError, ValueError):
                        box_anchors.append(float("nan"))
                if fractions is None:
                    anchors = box_anchors
        if anchors:
            if any(not math.isfinite(value) or not 0.0 <= value <= 1.0 for value in anchors):
                failures.append(
                    f"Slide {index} (figure) panel geometry must use finite values in [0,1]."
                )
            if any(right <= left for left, right in zip(anchors, anchors[1:])):
                failures.append(
                    f"Slide {index} (figure) panel geometry must increase in reading order."
                )
    elif slide_type == "references":
        items = _nonempty_lines(slide.get("items"))
        if not 5 <= len(items) <= 10:
            failures.append(
                f"Slide {index} (references) requires 5-10 reference items; found {len(items)}."
            )
    elif slide_type == "thanks":
        _require_nonempty_string(slide, "citation", index, slide_type, failures)
    return slide_type


def visible_texts(slide: Any) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
    return texts


def visible_spec_fields(slide: dict[str, Any]) -> list[tuple[str, str]]:
    fields: list[tuple[str, str]] = []
    for key in ("title", "caption", "authors", "citation", "subtitle", "kicker"):
        value = slide.get(key)
        if isinstance(value, str) and value.strip():
            fields.append((key, value))
    for key in ("bullets", "items", "body"):
        value = slide.get(key)
        if isinstance(value, str) and value.strip():
            fields.append((key, value))
        elif isinstance(value, list):
            fields.extend(
                (f"{key}[{index}]", item)
                for index, item in enumerate(value)
                if isinstance(item, str) and item.strip()
            )
    return fields


def referenced_panel_letters(notes: str) -> set[str]:
    letters = {
        match.group(1).upper()
        for match in re.finditer(r"【\s*([A-Za-z]|\d+)\s*(?:[圖图]|[:：])", notes)
    }
    letters.update(
        match.group(1).upper()
        for match in re.finditer(r"\bpanel\s*([A-Za-z]|\d+)\b", notes, re.IGNORECASE)
    )
    return letters


def resolve_asset(value: str | Path, spec_path: Path) -> Path:
    asset = Path(value).expanduser()
    return asset.resolve() if asset.is_absolute() else (spec_path.parent / asset).resolve()


def read_sidecar(image_path: Path) -> dict[str, Any] | None:
    sidecar = image_path.with_suffix(image_path.suffix + POSTPROCESS_SUFFIX)
    if not sidecar.is_file():
        return None
    try:
        value = json.loads(sidecar.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        return {}
    return value if isinstance(value, dict) else {}


def split_table_group(image_path: Path) -> str | None:
    match = SPLIT_TABLE_RE.match(image_path.stem)
    return match.group(1).replace("-", "_").lower() if match else None


def extraction_manifest(spec_path: Path, spec: dict[str, Any]) -> Path | None:
    meta = spec.get("meta", {})
    if isinstance(meta, dict):
        for key in ("extraction_manifest", "extracted_manifest"):
            value = meta.get(key)
            if isinstance(value, str) and value.strip():
                return resolve_asset(value, spec_path)
    for candidate in (
        spec_path.parent / "extracted" / "manifest.json",
        spec_path.parent.parent / "extracted" / "manifest.json",
    ):
        if candidate.is_file():
            return candidate.resolve()
    return None


def _check_split_tables(
    groups: dict[str, list[tuple[int, dict[str, Any], Path]]],
    failures: list[str],
    warnings: list[str],
) -> None:
    for group, entries in sorted(groups.items()):
        if len(entries) < 2:
            continue
        widths: dict[str, int] = {}
        on_screen: list[float | None] = []
        for _, entry, path in entries:
            try:
                with Image.open(path) as image:
                    widths[path.name] = image.width
            except (OSError, ValueError) as error:
                failures.append(f"Split table asset {path.name} cannot be opened: {error}")
                continue
            value = entry.get("image_width_in")
            try:
                on_screen.append(float(value) if value is not None else None)
            except (TypeError, ValueError):
                failures.append(f"Split table {group} has an invalid image_width_in: {value!r}.")

        if len(set(widths.values())) > 1:
            failures.append(f"Split table {group} has unequal image pixel widths: {widths}.")
        configured = {value for value in on_screen if value is not None}
        if len(configured) > 1:
            failures.append(
                f"Split table {group} has unequal on-screen image_width_in values: "
                f"{sorted(configured)}."
            )
        elif None in on_screen:
            warnings.append(
                f"Split table {group} should set the same image_width_in on every slide."
            )


def _check_content_structure(
    slides: list[dict[str, Any]], failures: list[str], warnings: list[str], mode: str
) -> None:
    content = [slide for slide in slides if slide.get("type") == "content"]
    flat: list[int] = []
    for index, slide in enumerate(slides, start=1):
        if slide.get("type") != "content":
            continue
        items: list[str] = []
        for key in ("bullets", "body", "items"):
            value = slide.get(key)
            if isinstance(value, list):
                items = [str(item) for item in value if str(item).strip()]
                break
            if isinstance(value, str):
                items = [line for line in value.splitlines() if line.strip()]
                break
        if len(items) < 3:
            continue
        joined = "\n".join(items)
        has_heading = any(re.search(r"[A-Za-z][\w /-]{2,}:", item) for item in items)
        has_logic = any(marker in joined for marker in ("→", "✅", "⚠"))
        if not has_heading and not has_logic:
            flat.append(index)

    if flat and mode == "full" and len(flat) > len(content) / 2:
        failures.append(
            f"{len(flat)}/{len(content)} teaching slides are flat bullet lists "
            f"without structured headings or clinical takeaways: {flat[:10]}."
        )
    elif flat:
        warnings.append(
            f"Teaching slide(s) {flat} would benefit from section headings or clinical takeaways."
        )


def validate_specification(
    spec_path: Path, *, mode: str = "full", audit_images: bool = True
) -> dict[str, Any]:
    spec_path = spec_path.expanduser().resolve()
    failures: list[str] = []
    warnings: list[str] = []
    if not spec_path.is_file():
        return {
            "ok": False,
            "spec": str(spec_path),
            "mode": mode,
            "failures": [f"Deck specification not found: {spec_path}"],
            "warnings": [],
        }
    try:
        spec = json.loads(spec_path.read_text(encoding="utf-8"))
        slides = spec.get("slides", [])
        if not isinstance(spec, dict) or not isinstance(slides, list):
            raise ValueError("The specification must contain a slides array.")
        if not all(isinstance(slide, dict) for slide in slides):
            raise ValueError("Each slide must be a JSON object.")
    except (json.JSONDecodeError, OSError, TypeError, ValueError, AttributeError) as error:
        return {
            "ok": False,
            "spec": str(spec_path),
            "mode": mode,
            "failures": [f"Deck specification is invalid: {error}"],
            "warnings": [],
        }

    if mode in CONFIG["modes"]:
        limits = CONFIG["modes"][mode]
        if not limits["minimum_slides"] <= len(slides) <= limits["maximum_slides"]:
            failures.append(
                f"{mode} mode expects {limits['minimum_slides']}-{limits['maximum_slides']} "
                f"slides but the specification contains {len(slides)}."
            )
    else:
        failures.append(f"Unknown QA mode: {mode}")

    meta = spec.get("meta")
    if not isinstance(meta, dict):
        failures.append("Top-level meta must be a JSON object.")
        meta = {}
    page_color = _hex_rgb(meta.get("bg_page"))
    text_color = _hex_rgb(meta.get("text_primary"))
    if page_color is not None and text_color is not None:
        ratio = _contrast_ratio(text_color, page_color)
        if ratio < 4.5:
            failures.append(
                f"Specification palette text_primary/bg_page contrast is {ratio:.2f}:1; "
                "at least 4.5:1 is required."
            )
    footer = meta.get("footer_label")
    if not isinstance(footer, str) or not footer.strip():
        failures.append("meta.footer_label is required and must be a non-empty English string.")
    elif CHINESE_RE.search(footer):
        failures.append("Visible meta.footer_label contains Chinese text.")
    elif not FOOTER_RE.match(footer):
        failures.append(
            "meta.footer_label must follow '<first author> et al — <journal> <year> | <English topic>'."
        )

    if slides:
        if slides[0].get("type") != "title":
            failures.append("Slide 1 must be the title slide.")
        if slides[-1].get("type") != "thanks":
            failures.append("The final slide must be the thanks slide.")

    counts: dict[str, int] = {}
    figures: dict[str, list[int]] = {}
    figure_asset_hashes: dict[str, list[tuple[int, str]]] = {}
    split_tables: dict[str, list[tuple[int, dict[str, Any], Path]]] = {}
    emoji_notes = 0
    note_entries: list[tuple[int, str]] = []
    image_assets = 0

    for index, slide in enumerate(slides, start=1):
        slide_type = _validate_slide_schema(slide, index, failures)
        counts[slide_type] = counts.get(slide_type, 0) + 1
        for field, value in visible_spec_fields(slide):
            if CHINESE_RE.search(value):
                failures.append(
                    f"Slide {index} ({slide_type}) visible field {field} contains Chinese text."
                )

        raw_notes = slide.get("notes")
        notes = raw_notes.strip() if isinstance(raw_notes, str) else ""
        if not notes:
            failures.append(f"Slide {index} has no speaker notes in its specification.")
        else:
            note_entries.append((index, notes))
            minimum_cjk = MIN_NOTE_CJK.get(slide_type, 12)
            cjk_count = _cjk_count(notes)
            if cjk_count < minimum_cjk:
                failures.append(
                    f"Slide {index} ({slide_type}) speaker notes are not substantive: "
                    f"{cjk_count} CJK characters; at least {minimum_cjk} are required."
                )
            if diversity_failure := _note_diversity_failure(notes):
                failures.append(f"Slide {index} ({slide_type}) {diversity_failure}.")
            if SIMPLIFIED_ONLY_RE.search(notes):
                sample = "".join(dict.fromkeys(SIMPLIFIED_ONLY_RE.findall(notes)))[:8]
                failures.append(
                    f"Slide {index} speaker notes appear to contain Simplified-Chinese "
                    f"forms ({sample}); use Traditional Chinese."
                )
            if not has_emoji(notes):
                failures.append(
                    f"Slide {index} speaker notes require at least one structural emoji marker."
                )
            else:
                emoji_notes += 1
            if slide_type == "content" and not has_closing_takeaway(notes):
                failures.append(
                    f"Slide {index} content speaker notes require a closing takeaway marker "
                    "(✅, 💡, or ⚠) in the latter half."
                )

        image_name = slide.get("image")
        if not isinstance(image_name, str) or not image_name.strip():
            if slide_type == "figure":
                # The schema validator has already emitted the field-specific error.
                pass
            continue
        image = resolve_asset(image_name, spec_path)
        if not image.is_file():
            failures.append(f"Slide {index} ({slide_type}) image is missing: {image}")
            continue
        image_assets += 1
        suffix = image.suffix.lower()
        if suffix not in RASTER_SUFFIXES | vector_table.VECTOR_SUFFIXES:
            failures.append(
                f"Slide {index} ({slide_type}) uses unsupported image format "
                f"{image.suffix!r}; use an authenticated raster or .emf vector table."
            )
        if PANEL_FRAGMENT_RE.search(image.stem):
            failures.append(f"Slide {index} references a panel fragment: {image.name}")
        if RAW_STREAM_RE.match(image.stem):
            failures.append(
                f"Slide {index} references raw PDF image {image.name}; use the "
                "color-decoded extracted/figures image instead."
            )

        sidecar_path = image.with_suffix(image.suffix + POSTPROCESS_SUFFIX)
        sidecar = read_sidecar(image)
        if suffix in RASTER_SUFFIXES | vector_table.VECTOR_SUFFIXES and sidecar is None:
            failures.append(f"Slide {index} is missing its image sidecar: {sidecar_path}")
        elif sidecar == {}:
            failures.append(f"Slide {index} has an invalid image sidecar: {sidecar_path}")
        elif suffix in RASTER_SUFFIXES and isinstance(sidecar, dict):
            failures.extend(
                f"Slide {index}: {failure}"
                for failure in validate_final_sidecar(image, sidecar)
            )
        elif suffix in vector_table.VECTOR_SUFFIXES and isinstance(sidecar, dict):
            failures.extend(
                f"Slide {index}: {failure}"
                for failure in vector_table.sidecar_structure_failures(sidecar)
            )

        if slide_type != "figure":
            continue

        labels = slide.get("panel_labels") or []
        if not isinstance(labels, list):
            failures.append(f"Figure slide {index} panel_labels must be a list.")
            labels = []
        metadata = sidecar or {}
        preserved = metadata.get("source_label_policy") == "preserve"
        if preserved and (metadata.get("native_labels") or labels):
            failures.append(
                f"Figure slide {index} preserves embedded panel labels but also requests "
                "native panel labels; this would duplicate the source letters."
            )
        cleanup_entries = metadata.get("panel_cleanup") or []
        edge_limit = metadata.get("max_edge_px", 4)
        boundary_limit = metadata.get("max_boundary_shift_px", 24)
        for panel_number, cleanup in enumerate(cleanup_entries, start=1):
            if not isinstance(cleanup, dict):
                failures.append(f"Figure slide {index} panel {panel_number} cleanup metadata is invalid.")
                continue
            overwritten = cleanup.get("label_overwritten_pixels", 0)
            if isinstance(overwritten, (int, float)) and overwritten > 0:
                failures.append(
                    f"Figure slide {index} panel {panel_number} overwrites {overwritten} "
                    "source-image pixels while handling a panel label."
                )
            trims = cleanup.get("edge_trim_px") or {}
            if any(
                not isinstance(value, int) or value < 0
                or not isinstance(edge_limit, int) or value > edge_limit
                for value in trims.values()
            ):
                failures.append(
                    f"Figure slide {index} panel {panel_number} exceeds its bounded "
                    f"{edge_limit}px edge-cleanup limit."
                )
            verified_limit = metadata.get("verified_edge_trim_max_px", 12)
            verified = cleanup.get("verified_edge_trim_px")
            verified_reason = cleanup.get("verified_edge_trim_reason", "")
            if verified is not None:
                valid_verified = (
                    isinstance(verified, dict)
                    and set(verified) == {"top", "bottom", "left", "right"}
                    and isinstance(verified_limit, int)
                    and not isinstance(verified_limit, bool)
                    and 0 <= verified_limit <= 12
                    and all(
                        isinstance(value, int)
                        and not isinstance(value, bool)
                        and 0 <= value <= verified_limit
                        for value in verified.values()
                    )
                )
                reason_allowed = verified_reason in {
                    "verified-pdf-exterior-band",
                    "verified-image-box-correction",
                    "manual-visual-review",
                }
                if not valid_verified or (any(verified.values()) and not reason_allowed):
                    failures.append(
                        f"Figure slide {index} panel {panel_number} has invalid verified "
                        "edge-trim evidence."
                    )
                total = cleanup.get("total_edge_trim_px")
                if valid_verified and total is not None and (
                    not isinstance(total, dict)
                    or set(total) != {"top", "bottom", "left", "right"}
                    or any(
                        total.get(side) != verified[side] + trims.get(side, 0)
                        for side in verified
                    )
                ):
                    failures.append(
                        f"Figure slide {index} panel {panel_number} has inconsistent "
                        "total edge-trim metadata."
                    )
            review = cleanup.get("residual_edge_review")
            if isinstance(review, dict) and review.get("status") == "needs-review":
                candidates = review.get("candidates") or {}
                warnings.append(
                    f"Figure slide {index} panel {panel_number} retains a narrow full-edge "
                    f"bright band requiring image-box or verified-trim review: {candidates}."
                )
            for adjustment in cleanup.get("boundary_adjustments") or []:
                if not isinstance(adjustment, dict):
                    failures.append(
                        f"Figure slide {index} panel {panel_number} boundary-adjustment metadata is invalid."
                    )
                    continue
                shift = adjustment.get("shift_px")
                if (
                    not isinstance(shift, int)
                    or shift <= 0
                    or not isinstance(boundary_limit, int)
                    or shift > boundary_limit
                    or adjustment.get("reason") != "preserve-complete-embedded-label-frame"
                ):
                    failures.append(
                        f"Figure slide {index} panel {panel_number} exceeds its bounded "
                        f"{boundary_limit}px label-safe boundary adjustment."
                    )
        if len(labels) > 1 and not slide.get("panel_geometry_exception"):
            fractions = slide.get("panel_label_x_fracs") or []
            boxes = slide.get("panel_boxes") or []
            sidecar_boxes = (sidecar or {}).get("panel_boxes") or []
            valid_geometry = (
                isinstance(fractions, list) and len(fractions) >= len(labels)
            ) or (
                isinstance(boxes, list) and len(boxes) >= len(labels)
            ) or (
                isinstance(sidecar_boxes, list) and len(sidecar_boxes) >= len(labels)
            ) or bool((sidecar or {}).get("native_labels"))
            if not valid_geometry:
                failures.append(
                    f"Figure slide {index} has {len(labels)} panel labels but no panel geometry."
                )
        semantic_labels: list[Any] = []
        for value in (
            labels,
            metadata.get("labels"),
            metadata.get("embedded_labels"),
            metadata.get("native_label_values"),
        ):
            if isinstance(value, list):
                semantic_labels.extend(value)
        label_set = {
            str(label).strip().upper()
            for label in semantic_labels
            if str(label).strip()
        }
        referenced = referenced_panel_letters(notes)
        unknown = referenced - label_set
        if unknown:
            failures.append(
                f"Figure slide {index} notes reference missing panel(s): {sorted(unknown)}."
            )
        missing_mentions = label_set - referenced
        if missing_mentions:
            failures.append(
                f"Figure slide {index} notes do not describe visible panel(s): "
                f"{sorted(missing_mentions)}."
            )

        caption = f"{slide.get('caption', '')} {slide.get('title', '')}"
        is_table = (
            "table" in image.name.lower()
            or str((sidecar or {}).get("asset_type", "")).lower() == "table"
            or bool(re.search(r"\btable\b", caption, re.IGNORECASE))
        )
        if is_table:
            margin = (sidecar or {}).get("margin")
            if isinstance(margin, int) and margin < TABLE_MARGIN_MIN:
                failures.append(
                    f"Figure slide {index} table {image.name} has an unsafe "
                    f"{margin}px edge margin; at least {TABLE_MARGIN_MIN}px is required."
                )
            group = split_table_group(image)
            if group:
                split_tables.setdefault(group, []).append((index, slide, image))
        else:
            if not slide.get("figure_slide_exception"):
                try:
                    digest = hashlib.sha256(image.read_bytes()).hexdigest()
                    figure_asset_hashes.setdefault(digest, []).append((index, image.name))
                except OSError as error:
                    failures.append(
                        f"Figure slide {index} image hash cannot be calculated: {error}."
                    )
            match = FIGURE_RE.search(caption)
            if match:
                figures.setdefault(match.group(1), []).append(index)

    failures.extend(duplicate_note_failures(note_entries))

    if counts.get("title", 0) != 1:
        failures.append(f"Expected exactly one title slide; found {counts.get('title', 0)}.")
    if counts.get("thanks", 0) != 1:
        failures.append(f"Expected exactly one closing slide; found {counts.get('thanks', 0)}.")
    if counts.get("content", 0) == 0:
        failures.append("Presentation contains no teaching-content slides.")
    if mode == "full":
        for required in ("outline", "references"):
            if counts.get(required, 0) != 1:
                failures.append(
                    f"Full teaching deck requires exactly one {required} slide; "
                    f"found {counts.get(required, 0)}."
                )
        if counts.get("part", 0) == 0:
            failures.append("Full teaching deck requires at least one numbered section divider.")

    title = next((slide for slide in slides if slide.get("type") == "title"), None)
    if title is not None:
        missing = [key for key in ("authors", "citation") if not str(title.get(key, "")).strip()]
        if missing:
            failures.append(f"Title slide is missing source-paper metadata: {', '.join(missing)}.")
    if slides and emoji_notes == 0:
        failures.append("Speaker notes do not contain any structural emoji markers.")

    outline = next((slide for slide in slides if slide.get("type") == "outline"), None)
    if outline is not None:
        items = outline.get("items") or []
        if isinstance(items, list) and items:
            if sum(bool(KEYCAP_RE.search(str(item))) for item in items) < len(items):
                warnings.append("Outline items should use numbered keycap markers.")
            if any(not re.search(r"slides?\s*\d+\s*[-–]\s*\d+", str(item), re.IGNORECASE)
                   for item in items):
                warnings.append("Outline items should indicate their corresponding slide ranges.")

    for number, indexes in sorted(figures.items()):
        if len(indexes) > 1:
            failures.append(
                f"Paper Figure {number} appears on multiple slides: "
                + ", ".join(str(index) for index in indexes)
            )
    for entries in figure_asset_hashes.values():
        if len(entries) > 1:
            failures.append(
                "The same figure image bytes are reused on multiple slides: "
                + ", ".join(f"slide {index} ({name})" for index, name in entries)
                + ". Each paper Figure must occupy exactly one slide."
            )
    _check_split_tables(split_tables, failures, warnings)
    _check_content_structure(slides, failures, warnings, mode)

    polarity_summary: dict[str, Any] | None = None
    manifest = extraction_manifest(spec_path, spec)
    if audit_images and manifest is not None:
        polarity = image_polarity.audit_extraction(manifest, persist=False)
        failures.extend(polarity.get("failures", []))
        warnings.extend(polarity.get("warnings", []))
        final_assets = image_polarity.audit_final_assets(spec_path, polarity)
        failures.extend(final_assets.get("failures", []))
        polarity_summary = {
            "checked_figures": polarity.get("checked_figures", 0),
            "unsafe_raw_streams": polarity.get("unsafe_raw_streams", 0),
            "checked_final_assets": final_assets.get("checked_assets", 0),
        }
    elif audit_images and image_assets:
        failures.append(
            "Every slide image requires a readable extraction manifest; original-PDF "
            "authenticity and complete provenance cannot be verified."
        )

    return {
        "ok": not failures,
        "stage": "specification",
        "spec": str(spec_path),
        "mode": mode,
        "slides": len(slides),
        "slides_with_emoji_notes": emoji_notes,
        "slide_types": counts,
        "image_polarity": polarity_summary,
        "failures": failures,
        "warnings": warnings,
    }


def expected_logo_path(spec: dict[str, Any] | None, spec_path: Path | None) -> Path:
    if spec is not None and spec_path is not None:
        meta = spec.get("meta", {})
        value = meta.get("logo_path") if isinstance(meta, dict) else None
        if isinstance(value, str) and value.strip():
            custom = resolve_asset(value, spec_path)
            if custom.is_file():
                return custom
    return LOGO_PATH


def logo_hash(path: Path) -> str | None:
    if not path.is_file():
        return None
    return hashlib.sha256(path.read_bytes()).hexdigest()


def slide_has_logo(slide: Any, expected_hash: str) -> bool:
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            if hashlib.sha256(shape.image.blob).hexdigest() == expected_hash:
                return True
        except (AttributeError, ValueError):
            continue
    return False


def _check_build_manifest(
    presentation: Any,
    pptx_path: Path,
    spec_path: Path | None,
    style: str | None,
    failures: list[str],
) -> dict[str, Any] | None:
    manifest, error = build_deck.read_build_manifest(pptx_path)
    if error or manifest is None:
        failures.append(
            "PowerPoint is missing a valid embedded build manifest; rebuild it with the "
            f"current canonical builder ({error or 'unknown manifest error'})."
        )
        return None
    if manifest.get("schema") != build_deck.MANIFEST_SCHEMA:
        failures.append(f"Unsupported build-manifest schema: {manifest.get('schema')!r}.")
    try:
        installed_version = (SKILL_ROOT / "VERSION").read_text(encoding="utf-8").strip()
    except OSError:
        installed_version = None
    if installed_version and manifest.get("skill_version") != installed_version:
        failures.append(
            f"PowerPoint was built with skill version {manifest.get('skill_version')!r}, "
            f"not installed version {installed_version!r}."
        )
    expected_digest = build_deck.manifest_sha256(manifest)
    if presentation.core_properties.identifier != expected_digest:
        failures.append(
            "PowerPoint core-property manifest digest does not match the embedded manifest."
        )
    keywords = presentation.core_properties.keywords or ""
    if "medical-journal-manifest/v1" not in keywords:
        failures.append("PowerPoint core properties do not identify the medical-journal manifest.")
    if style is not None and manifest.get("style") != style:
        failures.append(
            f"PowerPoint was built with style {manifest.get('style')!r}, not requested style {style!r}."
        )

    actual_size = {
        "width": int(presentation.slide_width),
        "height": int(presentation.slide_height),
    }
    if manifest.get("presentation_size_emu") != actual_size:
        failures.append(
            "PowerPoint canvas dimensions do not match the embedded build manifest."
        )
    try:
        actual_package = build_deck.make_package_binding(pptx_path)
    except (OSError, ValueError, zipfile.BadZipFile, ET.ParseError) as error:
        failures.append(f"PowerPoint package fingerprint cannot be calculated: {error}")
        actual_package = None
    if actual_package is not None and manifest.get("package_parts") != actual_package:
        failures.append(
            "PowerPoint package XML, relationships, media, masters, layouts, themes, "
            "or embedded objects changed after the canonical build."
        )

    if spec_path is not None and spec_path.is_file():
        try:
            specification = json.loads(spec_path.read_text(encoding="utf-8"))
            binding = build_deck.make_spec_binding(specification, spec_path)
        except (OSError, json.JSONDecodeError, TypeError, AttributeError) as error:
            failures.append(f"Cannot calculate the current specification binding: {error}")
        else:
            if manifest.get("spec_sha256") != binding["spec_sha256"]:
                failures.append(
                    "PowerPoint was not built from this deck specification (canonical spec SHA-256 mismatch)."
                )
            if manifest.get("slides") != binding["slides"]:
                failures.append(
                    "PowerPoint source binding does not match the current per-slide content/image SHA-256 inventory."
                )
            if manifest.get("external_bindings") != binding.get("external_bindings"):
                failures.append(
                    "PowerPoint article asset-map binding does not match the current caption/source map."
                )
            if style is not None:
                try:
                    baseline = build_deck.render_specification(
                        specification, spec_path.parent, style=style
                    )
                    expected_size = {
                        "width": int(baseline.slide_width),
                        "height": int(baseline.slide_height),
                    }
                    expected_visual = build_deck.make_serialized_rendered_binding(
                        baseline, exclude_managed_labels=True
                    )
                    expected_package = build_deck.make_serialized_package_binding(
                        baseline, exclude_managed_labels=True
                    )
                    actual_visual = build_deck.make_rendered_binding(
                        presentation, exclude_managed_labels=True
                    )
                    actual_canonical_package = build_deck.make_package_binding(
                        pptx_path, exclude_managed_labels=True
                    )
                except Exception as error:
                    failures.append(f"Cannot rebuild the canonical visual baseline: {error}")
                else:
                    if actual_size != expected_size:
                        failures.append(
                            "PowerPoint canvas dimensions differ from the fresh canonical "
                            "spec/style rebuild."
                        )
                    if actual_visual != expected_visual:
                        failures.append(
                            "PowerPoint visual state differs from a fresh canonical rebuild of "
                            "the supplied spec/style."
                        )
                    if actual_canonical_package != expected_package:
                        failures.append(
                            "PowerPoint package differs from the fresh canonical spec/style "
                            "rebuild, including slide-level XML, masters, layouts, themes, "
                            "relationships, or media."
                        )

    actual_rendered = build_deck.make_rendered_binding(presentation)
    if manifest.get("rendered_slides") != actual_rendered:
        failures.append(
            "PowerPoint slide text, notes, image bytes, geometry, fill, line, crop, or typography "
            "changed after the canonical build."
        )
    source_slides = manifest.get("slides")
    rendered_slides = manifest.get("rendered_slides")
    if isinstance(source_slides, list) and isinstance(rendered_slides, list):
        for index, (source, rendered) in enumerate(
            zip(source_slides, rendered_slides), start=1
        ):
            if not isinstance(source, dict) or not isinstance(rendered, dict):
                failures.append(f"Build manifest slide {index} has invalid binding records.")
                continue
            image_sha = source.get("image_sha256")
            picture_hashes = rendered.get("picture_sha256") or []
            if image_sha is not None and image_sha not in picture_hashes:
                failures.append(
                    f"Slide {index} does not embed the image bytes bound by its specification SHA-256."
                )
    return manifest


def _picture_hash(shape: Any) -> str | None:
    try:
        return hashlib.sha256(shape.image.blob).hexdigest()
    except (AttributeError, ValueError):
        return None


def _pptx_rgb(color: Any) -> tuple[int, int, int] | None:
    try:
        value = color.rgb
    except (AttributeError, ValueError):
        return None
    if value is None:
        return None
    raw = str(value)
    if not re.fullmatch(r"[0-9A-Fa-f]{6}", raw):
        return None
    return tuple(int(raw[index:index + 2], 16) for index in (0, 2, 4))


def _shape_fill_rgb(shape: Any) -> tuple[int, int, int] | None:
    try:
        return _pptx_rgb(shape.fill.fore_color)
    except (AttributeError, ValueError, TypeError):
        return None


def _check_text_shape_visibility(
    slide: Any, index: int, slide_w: int, slide_h: int, failures: list[str]
) -> None:
    """Catch fully off-slide text and explicit foreground/background collisions."""
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.text_frame.text.strip():
            continue
        left, top = int(shape.left), int(shape.top)
        width, height = int(shape.width), int(shape.height)
        right, bottom = left + width, top + height
        if width <= 0 or height <= 0 or right <= 0 or bottom <= 0 or left >= slide_w or top >= slide_h:
            failures.append(
                f"Slide {index} text shape {shape.name!r} lies completely outside the visible slide."
            )
        background = _shape_fill_rgb(shape)
        if background is None:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if not run.text.strip():
                    continue
                try:
                    foreground = _pptx_rgb(run.font.color)
                except (AttributeError, ValueError):
                    foreground = None
                if foreground is not None and _contrast_ratio(foreground, background) < 3.0:
                    failures.append(
                        f"Slide {index} text shape {shape.name!r} has explicit foreground/fill "
                        f"contrast below 3.0:1."
                    )
                    break


def _managed_panel_marker(image_stem: str, label: str) -> str:
    safe_figure = re.sub(r"[^A-Za-z0-9_.-]+", "_", image_stem)
    safe_label = re.sub(r"[^A-Za-z0-9_.-]+", "_", label)
    return f"MJ_PANEL_LABEL_{safe_figure}_{safe_label}"


def _check_managed_panel_shapes(
    slide: Any,
    specification: dict[str, Any] | None,
    spec_path: Path | None,
    index: int,
    failures: list[str],
) -> None:
    """Reject name-only manifest exclusions and constrain approved native labels."""
    expected: dict[str, str] = {}
    if isinstance(specification, dict):
        image_value = specification.get("image")
        labels = specification.get("panel_labels") or []
        if isinstance(image_value, str) and isinstance(labels, list):
            image_path = resolve_asset(image_value, spec_path) if spec_path is not None else Path(image_value)
            sidecar = read_sidecar(image_path) or {}
            native = sidecar.get("native_label_values") if sidecar.get("native_labels") else None
            if isinstance(native, list):
                labels = native
            for value in labels:
                label = str(value).strip()
                if label:
                    expected[_managed_panel_marker(image_path.stem, label)] = label

    seen: set[str] = set()
    for shape in slide.shapes:
        name = str(getattr(shape, "name", ""))
        if not name.startswith("MJ_PANEL_LABEL_"):
            continue
        if name not in expected:
            failures.append(
                f"Slide {index} contains an unrecognized managed panel-label shape {name!r}."
            )
            continue
        if name in seen:
            failures.append(f"Slide {index} contains duplicate managed panel-label marker {name!r}.")
        seen.add(name)
        text = shape.text_frame.text.strip() if getattr(shape, "has_text_frame", False) else ""
        if text != expected[name]:
            failures.append(
                f"Slide {index} managed panel-label shape {name!r} has unexpected text {text!r}."
            )
        if int(shape.width) > int(Inches(1.0)) or int(shape.height) > int(Inches(0.75)):
            failures.append(
                f"Slide {index} managed panel-label shape {name!r} exceeds the allowed label box."
            )


def _check_panel_labels(
    slide: Any,
    specification: dict[str, Any],
    index: int,
    expected_labels: list[Any],
    expected_logo_hash: str | None,
    slide_w: int,
    slide_h: int,
    failures: list[str],
) -> None:
    expected = [str(label).strip().upper() for label in expected_labels if str(label).strip()]
    if not expected:
        return
    label_shapes = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text_frame.text.strip().upper() in set(expected)
    ]
    actual = [shape.text_frame.text.strip().upper() for shape in label_shapes]
    if len(actual) != len(expected) or sorted(actual) != sorted(expected):
        failures.append(
            f"Figure slide {index} requires exactly panel labels {expected}; found {actual}."
        )

    pictures = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        and (expected_logo_hash is None or _picture_hash(shape) != expected_logo_hash)
    ]
    if not pictures:
        failures.append(f"Figure slide {index} has no non-logo picture for panel-label placement.")
        return
    picture = max(pictures, key=lambda shape: int(shape.width) * int(shape.height))
    expansion = int(Inches(0.65))
    x_fracs = specification.get("panel_label_x_fracs") or []
    boxes = specification.get("panel_boxes") or []
    for label_index, label in enumerate(expected):
        matches = [shape for shape in label_shapes if shape.text_frame.text.strip().upper() == label]
        if len(matches) != 1:
            continue
        shape = matches[0]
        left, top, right, bottom = (
            int(shape.left), int(shape.top), int(shape.left + shape.width), int(shape.top + shape.height)
        )
        if left < 0 or top < 0 or right > slide_w or bottom > slide_h:
            failures.append(f"Figure slide {index} panel label {label} lies outside the slide.")
        if (
            right < int(picture.left) - expansion
            or left > int(picture.left + picture.width) + expansion
            or bottom < int(picture.top) - expansion
            or top > int(picture.top + picture.height) + expansion
        ):
            failures.append(
                f"Figure slide {index} panel label {label} is not anchored near the figure image."
            )
        expected_frac: float | None = None
        try:
            if label_index < len(x_fracs):
                expected_frac = float(x_fracs[label_index])
            elif label_index < len(boxes) and isinstance(boxes[label_index], dict):
                expected_frac = float(boxes[label_index].get("right_x_frac"))
        except (TypeError, ValueError):
            expected_frac = None
        if expected_frac is not None and picture.width:
            actual_frac = (right - int(picture.left)) / int(picture.width)
            if (
                not math.isfinite(expected_frac)
                or not 0 <= expected_frac <= 1
                or abs(actual_frac - expected_frac) > 0.12
            ):
                failures.append(
                    f"Figure slide {index} panel label {label} is at x={actual_frac:.3f}; "
                    f"expected right_x_frac={expected_frac:.3f}."
                )


def validate_presentation(
    pptx_path: Path,
    *,
    spec_path: Path | None = None,
    mode: str = "full",
    style: str | None = None,
) -> dict[str, Any]:
    pptx_path = pptx_path.expanduser().resolve()
    spec_path = spec_path.expanduser().resolve() if spec_path is not None else None
    failures: list[str] = []
    warnings: list[str] = []
    if not pptx_path.is_file():
        return {"ok": False, "failures": [f"Presentation not found: {pptx_path}"], "warnings": []}
    try:
        presentation = Presentation(str(pptx_path))
    except Exception as error:
        return {"ok": False, "failures": [f"PowerPoint cannot be opened: {error}"], "warnings": []}

    slides = list(presentation.slides)
    count = len(slides)
    if not slides:
        failures.append("Presentation does not contain any slides.")
    if mode in CONFIG["modes"]:
        limits = CONFIG["modes"][mode]
        if not limits["minimum_slides"] <= count <= limits["maximum_slides"]:
            failures.append(
                f"{mode} mode expects {limits['minimum_slides']}-{limits['maximum_slides']} "
                f"slides but the presentation contains {count}."
            )
    else:
        failures.append(f"Unknown QA mode: {mode}")

    width, height = int(presentation.slide_width), int(presentation.slide_height)
    if height <= 0 or abs(width / height - 16 / 9) > 0.03:
        failures.append(f"Presentation is not widescreen 16:9 ({width} × {height}).")

    spec_slides: list[dict[str, Any]] = []
    specification_data: dict[str, Any] | None = None
    spec_report: dict[str, Any] | None = None
    if spec_path is not None:
        spec_report = validate_specification(spec_path, mode=mode)
        failures.extend(spec_report["failures"])
        warnings.extend(spec_report["warnings"])
        if spec_path.is_file():
            try:
                loaded = json.loads(spec_path.read_text(encoding="utf-8"))
                if isinstance(loaded, dict):
                    specification_data = loaded
                    spec_slides = list(loaded.get("slides", []))
            except (json.JSONDecodeError, OSError, TypeError, AttributeError):
                spec_slides = []
            if spec_slides and len(spec_slides) != count:
                failures.append(
                    f"Deck specification contains {len(spec_slides)} slides, "
                    f"but the PowerPoint contains {count}."
                )

    _check_build_manifest(presentation, pptx_path, spec_path, style, failures)

    logo_path = expected_logo_path(specification_data, spec_path)
    expected_logo = logo_hash(logo_path)
    if expected_logo is None:
        failures.append(f"Expected presentation logo is missing: {logo_path}")

    note_count = 0
    emoji_note_count = 0
    note_entries: list[tuple[int, str]] = []
    logo_count = 0
    visible_cjk_slides: list[int] = []
    for index, slide in enumerate(slides, start=1):
        if str(slide._element.get("show", "1")).strip().lower() in {"0", "false", "off"}:
            failures.append(f"Slide {index} is hidden and would be omitted from presentation delivery.")
        _check_text_shape_visibility(slide, index, width, height, failures)
        current_spec = spec_slides[index - 1] if index <= len(spec_slides) else None
        _check_managed_panel_shapes(
            slide, current_spec, spec_path, index, failures
        )
        if any(CHINESE_RE.search(text) for text in visible_texts(slide)):
            visible_cjk_slides.append(index)
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except (AttributeError, ValueError):
            notes = ""
        if not notes:
            failures.append(f"Slide {index} has no speaker notes.")
        else:
            note_count += 1
            note_entries.append((index, notes))
            slide_type = (
                str(spec_slides[index - 1].get("type", "missing"))
                if index <= len(spec_slides) else "missing"
            )
            cjk_count = _cjk_count(notes)
            minimum_cjk = MIN_NOTE_CJK.get(slide_type, 12)
            if cjk_count < minimum_cjk:
                failures.append(
                    f"Slide {index} speaker notes are not substantive: {cjk_count} CJK "
                    f"characters; at least {minimum_cjk} are required."
                )
            if diversity_failure := _note_diversity_failure(notes):
                failures.append(f"Slide {index} {diversity_failure}.")
            if SIMPLIFIED_ONLY_RE.search(notes):
                failures.append(f"Slide {index} speaker notes appear to use Simplified Chinese.")
            if "**" in notes:
                failures.append(f"Slide {index} speaker notes contain unrendered ** markup.")
            if has_emoji(notes):
                emoji_note_count += 1
            else:
                failures.append(f"Slide {index} speaker notes contain no structural emoji marker.")
            if slide_type == "content" and not has_closing_takeaway(notes):
                failures.append(
                    f"Slide {index} content speaker notes require a closing takeaway marker "
                    "(✅, 💡, or ⚠) in the latter half."
                )

        present_logo = bool(expected_logo and slide_has_logo(slide, expected_logo))
        logo_count += int(present_logo)
        if index <= len(spec_slides):
            specification = spec_slides[index - 1]
            slide_type = specification.get("type")
            if slide_type not in {"title", "thanks"} and not present_logo:
                failures.append(
                    f"Slide {index} ({slide_type}) does not contain the expected logo {logo_path.name}."
                )
            if slide_type == "figure":
                expected_labels = specification.get("panel_labels") or []
                image = specification.get("image")
                if not expected_labels and isinstance(image, str) and spec_path is not None:
                    sidecar = read_sidecar(resolve_asset(image, spec_path)) or {}
                    if sidecar.get("native_labels"):
                        expected_labels = (
                            sidecar.get("native_label_values")
                            or sidecar.get("labels")
                            or []
                        )
                if isinstance(expected_labels, list):
                    _check_panel_labels(
                        slide,
                        specification,
                        index,
                        expected_labels,
                        expected_logo,
                        width,
                        height,
                        failures,
                    )

    failures.extend(duplicate_note_failures(note_entries))

    if visible_cjk_slides:
        failures.append(
            "English-only slide text contains Chinese on slide(s): "
            + ", ".join(str(index) for index in visible_cjk_slides)
        )
    if slides and emoji_note_count == 0:
        failures.append("Speaker notes do not contain any structural emoji markers.")
    if spec_path is None and count > 2 and logo_count == 0:
        failures.append("No embedded presentation logo was found.")

    return {
        "ok": not failures,
        "stage": "presentation",
        "pptx": str(pptx_path),
        "spec": str(spec_path) if spec_path else None,
        "mode": mode,
        "style": style,
        "slides": count,
        "slides_with_notes": note_count,
        "slides_with_emoji_notes": emoji_note_count,
        "slides_with_logo": logo_count,
        "image_polarity": spec_report.get("image_polarity") if spec_report else None,
        "failures": list(dict.fromkeys(failures)),
        "warnings": list(dict.fromkeys(warnings)),
    }


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("target", type=Path, help="PowerPoint or, with --spec-only, deck JSON")
    parser.add_argument("--spec", type=Path)
    parser.add_argument("--spec-only", action="store_true", help="Validate a deck before building")
    parser.add_argument("--mode", choices=tuple(CONFIG["modes"]), default="full")
    parser.add_argument("--style", choices=("standard", "nice"), default="standard")
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args(argv)

    target = args.target.expanduser().resolve()
    if args.spec_only:
        report = validate_specification(target, mode=args.mode)
    else:
        report = validate_presentation(
            target,
            spec_path=args.spec.expanduser().resolve() if args.spec else None,
            mode=args.mode,
            style=args.style,
        )

    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        if "slides" in report:
            if args.spec_only:
                print(f"Specification slides: {report['slides']} | Mode: {report['mode']}")
            else:
                print(
                    f"Slides: {report['slides']} | Notes: {report['slides_with_notes']} | "
                    f"Logo: {report['slides_with_logo']} | Mode: {report['mode']}"
                )
        polarity = report.get("image_polarity")
        if polarity:
            print(
                f"PDF image polarity: {polarity['checked_figures']} checked; "
                f"{polarity['unsafe_raw_streams']} unsafe raw stream(s) isolated."
            )
        for failure in report["failures"]:
            print(f"[FAIL] {failure}", file=sys.stderr)
        for warning in report["warnings"]:
            print(f"[WARN] {warning}", file=sys.stderr)
        stage = "Specification" if args.spec_only else "Presentation"
        print(f"{stage} QA {'passed' if report['ok'] else 'failed'}.")

    return 0 if report["ok"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
