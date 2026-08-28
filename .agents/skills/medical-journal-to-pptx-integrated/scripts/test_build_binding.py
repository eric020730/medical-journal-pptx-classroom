#!/usr/bin/env python3
"""Regression tests for strict deck schema and PPTX/spec/style binding."""

from __future__ import annotations

import json
import re
import subprocess
import sys
import tempfile
import unittest
import zipfile
from unittest import mock
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Pt

import build_deck
import notes_quality
import qa_check


NOTE = "📌 這張投影片提供完整教學重點，協助理解研究內容與臨床判斷，並連結後續討論。"


def valid_spec(*, logo_path: str | None = None) -> dict:
    meta = {"footer_label": "Author et al — Journal 2026 | Synthetic topic"}
    if logo_path:
        meta["logo_path"] = logo_path
    slides = [
        {
            "type": "title",
            "title": "Synthetic Medical Journal Study",
            "authors": "Author A, Author B, et al.",
            "citation": "Journal 2026;1:1-10",
            "notes": NOTE + " 本頁介紹研究標題、作者資訊與期刊來源。",
        },
        {
            "type": "outline",
            "title": "Outline",
            "items": [
                "1️⃣ Background — Slides 3-10",
                "2️⃣ Methods — Slides 11-20",
                "3️⃣ Results — Slides 21-38",
            ],
            "notes": NOTE + " 本頁依序預告背景、方法與結果三個教學區段。",
        },
        {
            "type": "part", "number": 1, "title": "Background",
            "notes": NOTE + " 本頁開啟研究背景段落並提示接下來的判讀方向。",
        },
    ]
    teaching_topics = [
        "研究問題界定", "疾病負擔評估", "既有治療缺口", "研究假說形成",
        "試驗設計選擇", "收案族群界定", "納入條件判讀", "排除條件判讀",
        "介入措施細節", "對照策略合理性", "隨機分派流程", "盲法執行品質",
        "主要終點定義", "次要終點定義", "樣本數估算", "統計模型選擇",
        "缺失資料處理", "敏感度分析", "基線平衡判讀", "主要結果解讀",
        "次要結果解讀", "效果量臨床意義", "信賴區間判讀", "亞組分析限制",
        "不良事件比較", "依從性與交叉", "追蹤完整程度", "偏差風險評估",
        "外部效度判讀", "證據確定性", "研究優勢整理", "研究限制整理",
        "臨床決策轉譯", "共享決策應用", "未來研究方向",
    ]
    for number, topic in enumerate(teaching_topics, start=1):
        slides.append(
            {
                "type": "content",
                "title": f"Teaching Point {number}",
                "body": [
                    "Study context:",
                    "• Synthetic evidence supports this teaching example",
                    "✅ Clinical meaning remains explicit and actionable",
                ],
                "notes": (
                    NOTE + f" 本頁聚焦{topic}，補充頁面專屬的證據判讀。"
                    " ✅ 臨床結論應結合研究限制審慎運用。"
                ),
            }
        )
    slides.extend(
        [
            {
                "type": "references",
                "title": "Key References",
                "items": [f"{number}. Author et al. Journal 202{number};1:1-10." for number in range(1, 6)],
                "notes": NOTE + " 本頁整理五筆參考來源與查核文獻的方法。",
            },
            {
                "type": "thanks",
                "title": "Thank You",
                "citation": "Author et al — Journal 2026",
                "notes": NOTE + " 本頁總結教學內容並邀請聽眾提出問題。",
            },
        ]
    )
    return {"meta": meta, "slides": slides}


class BuildBindingTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temp = tempfile.TemporaryDirectory()
        self.root = Path(self.temp.name)

    def tearDown(self) -> None:
        self.temp.cleanup()

    def write_spec(self, value: dict, name: str = "deck.json") -> Path:
        path = self.root / name
        path.write_text(json.dumps(value, ensure_ascii=False), encoding="utf-8")
        return path

    def test_exact_spec_and_style_pass_but_wrong_binding_fails(self) -> None:
        spec_path = self.write_spec(valid_spec())
        pptx_path = self.root / "deck.pptx"
        build_deck.build(spec_path, pptx_path, style="standard")

        report = qa_check.validate_presentation(
            pptx_path, spec_path=spec_path, style="standard"
        )
        self.assertTrue(report["ok"], report["failures"])

        wrong_style = qa_check.validate_presentation(
            pptx_path, spec_path=spec_path, style="nice"
        )
        self.assertFalse(wrong_style["ok"])
        self.assertTrue(any("not requested style" in item for item in wrong_style["failures"]))

        changed = valid_spec()
        changed["slides"][3]["title"] = "Different Teaching Point"
        changed_path = self.write_spec(changed, "changed.json")
        wrong_spec = qa_check.validate_presentation(
            pptx_path, spec_path=changed_path, style="standard"
        )
        self.assertFalse(wrong_spec["ok"])
        self.assertTrue(any("canonical spec SHA-256 mismatch" in item for item in wrong_spec["failures"]))

    def test_post_build_visible_edit_is_detected(self) -> None:
        spec_path = self.write_spec(valid_spec())
        pptx_path = self.root / "deck.pptx"
        build_deck.build(spec_path, pptx_path, style="standard")
        prs = Presentation(pptx_path)
        changed = False
        for shape in prs.slides[3].shapes:
            if getattr(shape, "has_text_frame", False) and "Teaching Point" in shape.text:
                shape.text_frame.text = "Tampered Title"
                changed = True
                break
        self.assertTrue(changed)
        edited = self.root / "edited.pptx"
        prs.save(edited)
        report = qa_check.validate_presentation(
            edited, spec_path=spec_path, style="standard"
        )
        self.assertFalse(report["ok"])
        self.assertTrue(
            any(
                "changed after the canonical build" in item
                or "missing a valid embedded build manifest" in item
                for item in report["failures"]
            )
        )

    def test_post_build_geometry_and_typography_edits_are_detected(self) -> None:
        spec_path = self.write_spec(valid_spec())
        pptx_path = self.root / "deck.pptx"
        build_deck.build(spec_path, pptx_path, style="standard")
        prs = Presentation(pptx_path)
        target = next(
            shape for shape in prs.slides[3].shapes
            if getattr(shape, "has_text_frame", False) and "Teaching Point" in shape.text
        )
        target.left = -target.width - 1
        target.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
        edited = self.root / "geometry-edited.pptx"
        prs.save(edited)

        report = qa_check.validate_presentation(
            edited, spec_path=spec_path, style="standard"
        )

        self.assertFalse(report["ok"])
        self.assertTrue(any("geometry" in failure for failure in report["failures"]))
        self.assertTrue(any("outside the visible slide" in failure for failure in report["failures"]))

    def test_forged_manifest_cannot_hide_managed_overlay_or_hidden_slide(self) -> None:
        spec_path = self.write_spec(valid_spec())
        built = self.root / "built.pptx"
        altered = self.root / "forged-overlay.pptx"
        build_deck.build(spec_path, built, style="standard")
        presentation = Presentation(built)
        slide = presentation.slides[3]
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            presentation.slide_width,
            presentation.slide_height,
        )
        overlay._element.nvSpPr.cNvPr.set("name", "MJ_PANEL_LABEL_ATTACK")
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        slide._element.set("show", "0")
        presentation.save(altered)

        # A locally rewritten digest is not authoritative: the supplied spec
        # and style are rendered again to establish the canonical baseline.
        manifest, error = build_deck.read_build_manifest(altered)
        self.assertIsNone(error)
        self.assertIsNotNone(manifest)
        forged = Presentation(altered)
        manifest["rendered_slides"] = build_deck.make_rendered_binding(forged)
        build_deck._set_core_manifest_properties(forged, manifest)
        forged.save(altered)
        build_deck._embed_manifest_part(altered, manifest)

        report = qa_check.validate_presentation(
            altered, spec_path=spec_path, style="standard"
        )

        self.assertFalse(report["ok"])
        self.assertTrue(any("fresh canonical rebuild" in item for item in report["failures"]))
        self.assertTrue(any("unrecognized managed panel-label" in item for item in report["failures"]))
        self.assertTrue(any("is hidden" in item for item in report["failures"]))

    def test_proportional_canvas_size_tamper_is_detected(self) -> None:
        spec_path = self.write_spec(valid_spec())
        built = self.root / "built.pptx"
        altered = self.root / "canvas-altered.pptx"
        build_deck.build(spec_path, built, style="standard")

        with zipfile.ZipFile(built, "r") as source, zipfile.ZipFile(
            altered, "w", compression=zipfile.ZIP_DEFLATED
        ) as destination:
            for member in source.infolist():
                payload = source.read(member.filename)
                if member.filename == "ppt/presentation.xml":
                    matches = 0

                    def double_dimensions(match: re.Match[bytes]) -> bytes:
                        nonlocal matches
                        matches += 1
                        tag = match.group(0)
                        for attribute in (b"cx", b"cy"):
                            tag = re.sub(
                                attribute + rb'="(\d+)"',
                                lambda value: attribute
                                + b'="'
                                + str(int(value.group(1)) * 2).encode("ascii")
                                + b'"',
                                tag,
                                count=1,
                            )
                        return tag

                    payload = re.sub(
                        rb"<p:sldSz\b[^>]*/>", double_dimensions, payload, count=1
                    )
                    self.assertEqual(matches, 1)
                destination.writestr(member, payload)

        # Forge every locally controllable binding so the regression proves
        # the independent fresh spec/style baseline, not only manifest mismatch.
        manifest, error = build_deck.read_build_manifest(altered)
        self.assertIsNone(error)
        self.assertIsNotNone(manifest)
        forged = Presentation(altered)
        manifest["presentation_size_emu"] = {
            "width": int(forged.slide_width),
            "height": int(forged.slide_height),
        }
        manifest["rendered_slides"] = build_deck.make_rendered_binding(forged)
        forged.save(altered)
        manifest["package_parts"] = build_deck.make_package_binding(altered)
        build_deck._set_core_manifest_properties(forged, manifest)
        forged.save(altered)
        build_deck._embed_manifest_part(altered, manifest)

        report = qa_check.validate_presentation(
            altered, spec_path=spec_path, style="standard"
        )

        self.assertFalse(report["ok"])
        self.assertTrue(
            any("canvas dimensions" in item for item in report["failures"]),
            report["failures"],
        )
        self.assertTrue(
            any("fresh canonical" in item for item in report["failures"]),
            report["failures"],
        )

    def test_slide_transition_timing_xml_tamper_is_detected(self) -> None:
        spec_path = self.write_spec(valid_spec())
        built = self.root / "built.pptx"
        altered = self.root / "transition-altered.pptx"
        build_deck.build(spec_path, built, style="standard")
        transition = (
            b'<p:transition advClick="0" advTm="0"><p:cut/></p:transition>'
        )
        changed = 0
        with zipfile.ZipFile(built, "r") as source, zipfile.ZipFile(
            altered, "w", compression=zipfile.ZIP_DEFLATED
        ) as destination:
            for member in source.infolist():
                payload = source.read(member.filename)
                if re.fullmatch(r"ppt/slides/slide\d+\.xml", member.filename):
                    self.assertIn(b"</p:sld>", payload)
                    payload = payload.replace(b"</p:sld>", transition + b"</p:sld>", 1)
                    changed += 1
                destination.writestr(member, payload)
        self.assertEqual(changed, 40)

        report = qa_check.validate_presentation(
            altered, spec_path=spec_path, style="standard"
        )

        self.assertFalse(report["ok"])
        self.assertTrue(
            any("package" in item.lower() for item in report["failures"]),
            report["failures"],
        )

    def test_slide_master_visual_tamper_is_detected(self) -> None:
        spec_path = self.write_spec(valid_spec())
        built = self.root / "built.pptx"
        altered = self.root / "master-altered.pptx"
        build_deck.build(spec_path, built, style="nice")
        injected = (
            b'<p:sp><p:nvSpPr><p:cNvPr id="9999" name="INJECTED FULL-SLIDE RED MASTER"/>'
            b'<p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="0" y="0"/>'
            b'<a:ext cx="12192000" cy="6858000"/></a:xfrm><a:prstGeom prst="rect">'
            b'<a:avLst/></a:prstGeom><a:solidFill><a:srgbClr val="FF0000"/></a:solidFill>'
            b'<a:ln><a:noFill/></a:ln></p:spPr></p:sp>'
        )
        changed = 0
        with zipfile.ZipFile(built, "r") as source, zipfile.ZipFile(
            altered, "w", compression=zipfile.ZIP_DEFLATED
        ) as destination:
            for member in source.infolist():
                payload = source.read(member.filename)
                if member.filename == "ppt/slideMasters/slideMaster1.xml":
                    self.assertIn(b"</p:spTree>", payload)
                    payload = payload.replace(
                        b"</p:spTree>", injected + b"</p:spTree>", 1
                    )
                    changed += 1
                destination.writestr(member, payload)
        self.assertEqual(changed, 1)

        report = qa_check.validate_presentation(
            altered, spec_path=spec_path, style="nice"
        )

        self.assertFalse(report["ok"])
        self.assertTrue(
            any("package" in item.lower() for item in report["failures"]),
            report["failures"],
        )

    def test_manifest_package_wiring_tamper_is_detected(self) -> None:
        spec_path = self.write_spec(valid_spec())
        built = self.root / "built.pptx"
        altered = self.root / "wiring-altered.pptx"
        build_deck.build(spec_path, built, style="standard")
        changed = 0
        with zipfile.ZipFile(built, "r") as source, zipfile.ZipFile(
            altered, "w", compression=zipfile.ZIP_DEFLATED
        ) as destination:
            for member in source.infolist():
                payload = source.read(member.filename)
                if member.filename == "_rels/.rels":
                    root = ET.fromstring(payload)
                    for relation in root:
                        if relation.get("Type") == build_deck.MANIFEST_REL_TYPE:
                            relation.set("Target", "docProps/missing-custom.xml")
                            changed += 1
                    ET.register_namespace("", build_deck.RELATIONSHIP_NS)
                    payload = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                destination.writestr(member, payload)
        self.assertEqual(changed, 1)

        report = qa_check.validate_presentation(
            altered, spec_path=spec_path, style="standard"
        )

        self.assertFalse(report["ok"])
        self.assertTrue(
            any("wrong target" in item for item in report["failures"]),
            report["failures"],
        )

    def test_explicit_white_on_white_text_fill_is_rejected(self) -> None:
        spec_path = self.write_spec(valid_spec())
        pptx_path = self.root / "deck.pptx"
        build_deck.build(spec_path, pptx_path, style="standard")
        prs = Presentation(pptx_path)
        target = next(
            shape for shape in prs.slides[3].shapes
            if getattr(shape, "has_text_frame", False)
            and shape.text_frame.paragraphs[0].runs
        )
        target.fill.solid()
        target.fill.fore_color.rgb = RGBColor(255, 255, 255)
        for paragraph in target.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
        edited = self.root / "contrast-edited.pptx"
        prs.save(edited)

        report = qa_check.validate_presentation(
            edited, spec_path=spec_path, style="standard"
        )

        self.assertFalse(report["ok"])
        self.assertTrue(any("contrast below" in failure for failure in report["failures"]))

    def test_spec_rejects_white_on_white_palette_and_repeated_note_padding(self) -> None:
        value = valid_spec()
        value["meta"]["bg_page"] = "#FFFFFF"
        value["meta"]["text_primary"] = "#FFFFFF"
        for slide in value["slides"]:
            slide["notes"] = "📌 " + "中" * 24
        spec_path = self.write_spec(value)

        report = qa_check.validate_specification(spec_path, audit_images=False)

        self.assertFalse(report["ok"])
        self.assertTrue(any("contrast" in failure for failure in report["failures"]))
        self.assertTrue(any("repetitive padding" in failure for failure in report["failures"]))

    def test_notes_reject_short_phrase_loops_and_cross_slide_boilerplate(self) -> None:
        self.assertFalse(qa_check.has_emoji("這段文字沒有開頭掃描標記" * 8 + "💡"))
        self.assertIn(
            "repeated phrase",
            qa_check._note_diversity_failure("💡" + "研究結果" * 5) or "",
        )
        value = valid_spec()
        repeated = "📌 本頁完整說明研究背景結果限制與臨床判斷，協助教學討論與證據應用。"
        for slide in value["slides"]:
            slide["notes"] = repeated
        spec_path = self.write_spec(value)

        report = qa_check.validate_specification(spec_path, audit_images=False)

        self.assertFalse(report["ok"])
        self.assertTrue(any(
            "reuse the same normalized speaker notes" in failure
            for failure in report["failures"]
        ))

    def test_notes_reject_boilerplate_differing_only_by_page_number(self) -> None:
        value = valid_spec()
        for index, slide in enumerate(value["slides"], start=1):
            slide["notes"] = (
                "📌 本頁整理研究背景、方法、結果與臨床判讀，提供完整而清楚的教學說明，"
                f"幫助聽眾掌握重點與限制。第{index}頁。"
                " ✅ 臨床結論應依證據品質審慎解讀。"
            )
        spec_path = self.write_spec(value)

        report = qa_check.validate_specification(spec_path, audit_images=False)

        self.assertFalse(report["ok"])
        self.assertTrue(
            any(
                "reuse the same normalized speaker notes" in failure
                for failure in report["failures"]
            ),
            report["failures"],
        )

    def test_notes_reject_boilerplate_differing_only_by_teaching_ordinal(self) -> None:
        value = valid_spec()
        for index, slide in enumerate(value["slides"], start=1):
            slide["notes"] = (
                "📌 本頁整理研究背景、方法、結果與臨床判讀，提供完整教學說明。"
                f"這是第{index}個教學重點。"
                " ✅ 臨床結論應依證據品質審慎解讀。"
            )
        spec_path = self.write_spec(value)

        report = qa_check.validate_specification(spec_path, audit_images=False)

        self.assertFalse(report["ok"])
        self.assertTrue(
            any(
                "reuse the same normalized speaker notes" in failure
                for failure in report["failures"]
            ),
            report["failures"],
        )

    def test_note_signature_preserves_clinical_numeric_values(self) -> None:
        first = "📌 治療組的 hazard ratio 為 0.72，樣本數為 415 人。"
        second = "📌 治療組的 hazard ratio 為 0.83，樣本數為 512 人。"
        self.assertNotEqual(
            notes_quality.normalized_note_signature(first),
            notes_quality.normalized_note_signature(second),
        )
        self.assertNotEqual(
            notes_quality.normalized_note_signature("第1張影像顯示主要終點"),
            notes_quality.normalized_note_signature("第2張影像顯示主要終點"),
        )
        numeric_notes = [
            (
                index,
                "📌 本頁完整整理研究方法、主要結果、限制與臨床判讀，"
                f"治療效果 HR {hr}，95% CI {ci}，樣本數 {sample} 人。"
                " ✅ 臨床結論應依證據品質與病人特徵審慎解讀。",
            )
            for index, hr, ci, sample in (
                (1, "0.72", "0.60–0.86", 415),
                (2, "0.83", "0.71–0.98", 512),
                (3, "0.91", "0.79–1.05", 638),
            )
        ]
        self.assertEqual(notes_quality.duplicate_note_failures(numeric_notes), [])

    def test_notes_reject_near_duplicate_boilerplate_with_one_unique_character(self) -> None:
        value = valid_spec()
        unique_characters = (
            "甲乙丙丁戊己庚辛壬癸天地玄黃宇宙洪荒日月盈昃辰宿列張寒來暑往"
            "秋收冬藏閏餘成歲律呂調陽雲騰致雨露結為霜"
        )
        self.assertGreaterEqual(len(unique_characters), len(value["slides"]))
        for character, slide in zip(unique_characters, value["slides"]):
            slide["notes"] = (
                "📌 本頁完整整理研究背景、方法、主要結果、研究限制與臨床判讀，"
                "提供清楚且可直接使用的教學說明。"
                f"{character} ✅ 臨床結論應依證據品質、病人特徵與治療情境審慎解讀。"
            )
        spec_path = self.write_spec(value)

        report = qa_check.validate_specification(spec_path, audit_images=False)

        self.assertFalse(report["ok"])
        self.assertTrue(
            any("near-identical speaker-note boilerplate" in item for item in report["failures"]),
            report["failures"],
        )

    def test_non_string_notes_are_rejected_without_coercion(self) -> None:
        value = valid_spec()
        value["slides"][3]["notes"] = ["not", "a", "string"]
        spec_path = self.write_spec(value)

        report = qa_check.validate_specification(spec_path, audit_images=False)

        self.assertFalse(report["ok"])
        self.assertTrue(any("speaker notes must be a string" in item for item in report["failures"]))

    def test_custom_logo_is_verified_from_spec(self) -> None:
        logo = self.root / "custom-logo.png"
        Image.new("RGB", (80, 80), (12, 90, 180)).save(logo)
        spec_path = self.write_spec(valid_spec(logo_path=logo.name))
        pptx_path = self.root / "custom-logo.pptx"
        build_deck.build(spec_path, pptx_path, style="standard")
        report = qa_check.validate_presentation(
            pptx_path, spec_path=spec_path, style="standard"
        )
        self.assertTrue(report["ok"], report["failures"])

    def test_empty_content_and_one_character_notes_are_rejected(self) -> None:
        bad = valid_spec()
        for slide in bad["slides"]:
            slide["notes"] = "中"
            if slide["type"] == "content":
                slide["body"] = []
        spec_path = self.write_spec(bad)
        report = qa_check.validate_specification(spec_path, audit_images=False)
        self.assertFalse(report["ok"])
        self.assertTrue(any("not substantive" in item for item in report["failures"]))
        self.assertTrue(any("requires at least 3" in item for item in report["failures"]))

    def test_nice_content_image_and_thanks_citation_are_rendered(self) -> None:
        image = self.root / "content-image.png"
        Image.new("RGB", (320, 180), (220, 230, 240)).save(image)
        value = valid_spec()
        value["slides"][3]["image"] = image.name
        spec_path = self.write_spec(value)
        pptx_path = self.root / "nice.pptx"
        build_deck.build(
            spec_path,
            pptx_path,
            style="nice",
            allow_unprocessed_assets=True,
        )
        prs = Presentation(pptx_path)
        pictures = [
            shape for shape in prs.slides[3].shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        self.assertGreaterEqual(len(pictures), 2)  # content image + corner logo
        thanks_text = "\n".join(
            shape.text for shape in prs.slides[-1].shapes
            if getattr(shape, "has_text_frame", False)
        )
        self.assertIn("Author et al — Journal 2026", thanks_text)

    def test_standard_panel_boxes_drive_label_x_positions(self) -> None:
        image = self.root / "figure.png"
        Image.new("RGB", (600, 300), (230, 230, 230)).save(image)
        value = valid_spec()
        value["slides"][3] = {
            "type": "figure",
            "title": "Synthetic Figure",
            "image": image.name,
            "caption": "Figure 1. Synthetic panels.",
            "panel_labels": ["A", "B"],
            "panel_boxes": [
                {"right_x_frac": 0.25},
                {"right_x_frac": 0.80},
            ],
            "notes": "🖼️ 這張合成圖完整展示兩個研究面板，協助理解主要結果與臨床意義。",
        }
        spec_path = self.write_spec(value)
        pptx_path = self.root / "panels.pptx"
        build_deck.build(
            spec_path,
            pptx_path,
            style="standard",
            allow_unprocessed_assets=True,
        )
        prs = Presentation(pptx_path)
        slide = prs.slides[3]
        pictures = [
            shape for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        figure = max(pictures, key=lambda shape: int(shape.width) * int(shape.height))
        labels = {
            shape.text.strip(): shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and shape.text.strip() in {"A", "B"}
        }
        self.assertEqual(set(labels), {"A", "B"})
        for label, expected in (("A", 0.25), ("B", 0.80)):
            shape = labels[label]
            actual = (
                int(shape.left + shape.width) - int(figure.left)
            ) / int(figure.width)
            self.assertAlmostEqual(actual, expected, delta=0.01)

    def test_qa_spec_rejects_incomplete_final_raster_sidecar(self) -> None:
        image = self.root / "unsafe-content.png"
        Image.new("RGB", (320, 180), (220, 230, 240)).save(image)
        Path(str(image) + ".postprocess.json").write_text(
            json.dumps({"source": "rendered.png", "asset_type": "figure"}),
            encoding="utf-8",
        )
        value = valid_spec()
        value["slides"][3]["image"] = image.name
        spec_path = self.write_spec(value)

        report = qa_check.validate_specification(spec_path, audit_images=False)

        self.assertFalse(report["ok"])
        self.assertTrue(
            any("safety_margin_px" in failure for failure in report["failures"]),
            report["failures"],
        )
        self.assertTrue(
            any("padded_size_px" in failure for failure in report["failures"]),
            report["failures"],
        )

    def test_qa_warns_on_narrow_residual_full_edge_bright_band(self) -> None:
        panel = self.root / "residual-panel.png"
        image = Image.new("RGB", (160, 120), (10, 10, 10))
        ImageDraw.Draw(image).rectangle((0, 115, 159, 119), fill=(252, 252, 252))
        image.save(panel)
        final = self.root / "residual-composite.png"
        subprocess.run([
            sys.executable,
            str(Path(__file__).with_name("recompose_panels_banded.py")),
            str(final),
            "--inputs", str(panel),
            "--geometry", str(self.root / "residual-geometry.json"),
        ], check=True, capture_output=True, text=True)
        value = valid_spec()
        value["slides"][3] = {
            "type": "figure",
            "title": "Figure 1. Synthetic edge review",
            "image": final.name,
            "caption": "Figure 1. Synthetic edge-review example.",
            "notes": "🖼️ 本頁說明合成影像的邊緣檢查結果與保守裁切限制，協助確認來源像素完整性。",
        }
        spec_path = self.write_spec(value)

        report = qa_check.validate_specification(spec_path, audit_images=False)

        self.assertTrue(report["ok"], report["failures"])
        self.assertTrue(
            any("retains a narrow full-edge bright band" in warning for warning in report["warnings"]),
            report["warnings"],
        )

    def test_qa_rejects_bounded_bright_frame_when_cleanup_is_disabled(self) -> None:
        panel = self.root / "disabled-cleanup-panel.png"
        image = Image.new("RGB", (160, 120), (10, 10, 10))
        ImageDraw.Draw(image).rectangle((0, 0, 1, 119), fill=(252, 252, 252))
        image.save(panel)
        final = self.root / "disabled-cleanup-composite.png"
        subprocess.run([
            sys.executable,
            str(Path(__file__).with_name("recompose_panels_banded.py")),
            str(final),
            "--inputs", str(panel),
            "--geometry", str(self.root / "disabled-cleanup-geometry.json"),
            "--no-trim",
        ], check=True, capture_output=True, text=True)
        value = valid_spec()
        value["slides"][3] = {
            "type": "figure",
            "title": "Figure 1. Synthetic disabled cleanup",
            "image": final.name,
            "caption": "Figure 1. Synthetic disabled-cleanup example.",
            "notes": "🖼️ 本頁驗證停用自動修邊時，薄白框必須被品質檢查攔截，避免殘留於輸出的臨床影像。",
        }
        spec_path = self.write_spec(value)

        report = qa_check.validate_specification(spec_path, audit_images=False)

        self.assertFalse(report["ok"])
        self.assertTrue(
            any("Bounded cleanup was disabled" in failure for failure in report["failures"]),
            report["failures"],
        )

    def test_vector_only_slide_requires_extraction_manifest_before_qa_or_build(self) -> None:
        vector = self.root / "table.emf"
        vector.write_bytes(b"synthetic-emf")
        source_pdf = self.root / "paper.pdf"
        source_pdf.write_bytes(b"%PDF-1.7\nsynthetic")
        sidecar = vector.with_suffix(vector.suffix + ".postprocess.json")
        sidecar.write_text(json.dumps({
            "schema": "medical-journal-vector-table-sidecar/v1",
            "command": "vector-table",
            "asset_type": "table",
            "intermediate": False,
            "source": str(source_pdf),
            "source_pdf_sha256": "0" * 64,
            "page": 1,
            "requested_bbox_pt": [10.0, 10.0, 100.0, 80.0],
            "padding_pt": {"x": 15.0, "top": 6.0, "bottom": 10.0},
            "effective_bbox_pt": [0.0, 4.0, 115.0, 90.0],
            "background_rgb": "#FFFFFF",
            "canonical_svg_sha256": "1" * 64,
            "output_sha256": "2" * 64,
            "output_size_bytes": len(vector.read_bytes()),
            "image_aspect": 1.25,
            "converter": {"name": "LibreOffice", "version": "synthetic"},
        }), encoding="utf-8")
        value = valid_spec()
        value["slides"][3] = {
            "type": "figure",
            "title": "Table 1. Synthetic outcome",
            "image": vector.name,
            "caption": "Table 1. Synthetic outcome.",
            "notes": "📊 本頁說明研究表格中的主要結果、限制與臨床判讀方式，協助審慎應用證據。",
        }
        spec_path = self.write_spec(value)

        report = qa_check.validate_specification(spec_path)

        self.assertFalse(report["ok"])
        self.assertTrue(
            any("Every slide image requires" in failure for failure in report["failures"]),
            report["failures"],
        )
        with self.assertRaisesRegex(RuntimeError, "extraction manifest"):
            build_deck.build(spec_path, self.root / "vector.pptx", style="standard")

    def test_failed_overwrite_preserves_existing_deck_bytes(self) -> None:
        class FailingPresentation:
            def __init__(self) -> None:
                presentation = Presentation()
                presentation.slides.add_slide(presentation.slide_layouts[6])
                self.slides = presentation.slides
                self.core_properties = presentation.core_properties

            def save(self, _path: str) -> None:
                raise RuntimeError("injected save failure")

        spec_path = self.write_spec(valid_spec())
        destination = self.root / "existing.pptx"
        destination.write_bytes(b"previous verified deck")

        with mock.patch.object(
            build_deck.build_deck_standard,
            "build",
            return_value=FailingPresentation(),
        ):
            with self.assertRaisesRegex(RuntimeError, "injected save failure"):
                build_deck.build(
                    spec_path,
                    destination,
                    style="standard",
                    allow_unprocessed_assets=True,
                    overwrite=True,
                )

        self.assertEqual(destination.read_bytes(), b"previous verified deck")

    def test_manifest_replace_waits_until_source_archive_is_closed(self) -> None:
        deck = self.root / "manifest-source.pptx"
        Presentation().save(deck)
        manifest = {"schema": build_deck.MANIFEST_SCHEMA, "test": "windows-lock"}
        real_zip_file = zipfile.ZipFile
        real_replace = build_deck.os.replace
        source_archives: list[zipfile.ZipFile] = []

        def tracked_zip_file(
            file: object, mode: str = "r", *args: object, **kwargs: object
        ) -> zipfile.ZipFile:
            archive = real_zip_file(file, mode, *args, **kwargs)
            if Path(file) == deck and mode == "r":
                source_archives.append(archive)
            return archive

        def checked_replace(source: object, destination: object) -> None:
            self.assertEqual(len(source_archives), 1)
            self.assertIsNone(
                source_archives[0].fp,
                "Windows cannot replace a PPTX while its source ZipFile is open",
            )
            real_replace(source, destination)

        with mock.patch.object(
            build_deck.zipfile, "ZipFile", side_effect=tracked_zip_file
        ), mock.patch.object(build_deck.os, "replace", side_effect=checked_replace):
            build_deck._embed_manifest_part(deck, manifest)

        self.assertIsNone(build_deck.validate_manifest_wiring(deck))

    def test_native_label_stamping_preserves_build_manifest_binding(self) -> None:
        image = self.root / "manifest-panels.png"
        Image.new("RGB", (600, 300), (30, 30, 30)).save(image)
        value = valid_spec()
        value["slides"][3] = {
            "type": "figure",
            "title": "Synthetic Figure",
            "image": image.name,
            "caption": "Figure 1. Synthetic panel image.",
            "notes": "🖼️ 這張合成圖完整展示研究面板，協助理解主要結果與臨床判斷價值。",
        }
        spec_path = self.write_spec(value)
        geometry = self.root / "geometry.json"
        geometry.write_text(json.dumps({
            image.stem: [{"label": "A", "fx_right": 0.80, "fy_center": 0.80}]
        }))
        built = self.root / "built.pptx"
        stamped = self.root / "stamped.pptx"
        build_deck.build(
            spec_path, built, style="standard", allow_unprocessed_assets=True
        )

        result = subprocess.run(
            [
                sys.executable,
                str(Path(__file__).resolve().parent / "add_panel_labels.py"),
                str(built), str(stamped), "--spec", str(spec_path),
                "--geometry", str(geometry),
            ],
            capture_output=True,
            text=True,
        )
        self.assertEqual(result.returncode, 0, result.stderr)
        presentation = Presentation(stamped)
        failures: list[str] = []
        qa_check._check_build_manifest(
            presentation, stamped, spec_path, "standard", failures
        )
        self.assertEqual(failures, [])


if __name__ == "__main__":
    unittest.main()
