#!/usr/bin/env python3
"""Build one canonical medical-journal deck in standard or nice visual style."""

from __future__ import annotations

import argparse
import base64
import contextlib
import hashlib
import json
import os
import re
import sys
import tempfile
import zipfile
from xml.etree import ElementTree as ET
from pathlib import Path
from typing import Any

from pptx import Presentation

import build_deck_nice
import build_deck_standard


STYLE_IMAGE_BOXES = {
    "standard": {"width_in": 12.10, "height_in": 4.85},
    "nice": {"width_in": 12.13, "height_in": 4.95},
}

MANIFEST_SCHEMA = "medical-journal-build-manifest/v1"
MANIFEST_PART = "docProps/custom.xml"
MANIFEST_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/custom-properties"
)
MANIFEST_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.custom-properties+xml"
)
CUSTOM_PROPERTY_NS = (
    "http://schemas.openxmlformats.org/officeDocument/2006/custom-properties"
)
VT_NS = "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"
CUSTOM_PROPERTY_FMTID = "{D5CDD505-2E9C-101B-9397-08002B2CF9AE}"
RASTER_SUFFIXES = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".webp", ".bmp"}
PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
RELATIONSHIP_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"


def _canonical_bytes(value: Any) -> bytes:
    return json.dumps(
        value, ensure_ascii=False, sort_keys=True, separators=(",", ":")
    ).encode("utf-8")


def _sha256_bytes(value: bytes) -> str:
    return hashlib.sha256(value).hexdigest()


def _sha256_file(path: Path) -> str | None:
    try:
        return _sha256_bytes(path.read_bytes())
    except OSError:
        return None


def _canonical_package_xml(
    name: str, payload: bytes, *, exclude_managed_labels: bool
) -> bytes:
    """Canonicalize one OOXML part and remove only explicitly managed state."""
    root = ET.fromstring(payload)
    if name == "_rels/.rels":
        for relation in list(root):
            if relation.get("Type") == MANIFEST_REL_TYPE:
                root.remove(relation)
    elif name == "[Content_Types].xml":
        for override in list(root):
            if override.get("PartName") == f"/{MANIFEST_PART}":
                root.remove(override)
    elif exclude_managed_labels and name.startswith("ppt/slides/slide"):
        shape_tree_tag = f"{{{PRESENTATION_NS}}}spTree"
        nonvisual_tag = f"{{{PRESENTATION_NS}}}cNvPr"
        for shape_tree in root.iter(shape_tree_tag):
            for shape in list(shape_tree):
                nonvisual = shape.find(f".//{nonvisual_tag}")
                if (
                    nonvisual is not None
                    and str(nonvisual.get("name", "")).startswith("MJ_PANEL_LABEL_")
                ):
                    shape_tree.remove(shape)
    serialized = ET.tostring(root, encoding="unicode")
    return ET.canonicalize(
        xml_data=serialized, rewrite_prefixes=True
    ).encode("utf-8")


def make_package_binding(
    pptx_path: Path, *, exclude_managed_labels: bool = False
) -> list[dict[str, str]]:
    """Fingerprint every canonical package part that can affect the deck.

    Non-visual document properties (including the self-referential manifest and
    its core digest) are intentionally excluded. The manifest relationship and
    content-type entries are removed before XML canonicalization. All remaining XML, relationships,
    media, themes, masters, layouts, transitions, and embedded objects are bound.
    """
    pptx_path = Path(pptx_path).expanduser().resolve()
    with zipfile.ZipFile(pptx_path, "r") as package:
        members = [member for member in package.infolist() if not member.is_dir()]
        names = [member.filename for member in members]
        duplicates = sorted({name for name in names if names.count(name) > 1})
        if duplicates:
            raise ValueError(
                "PowerPoint package contains duplicate members: " + ", ".join(duplicates)
            )
        binding: list[dict[str, str]] = []
        for member in sorted(members, key=lambda value: value.filename):
            name = member.filename
            # Document properties and thumbnail metadata do not affect slide
            # rendering and may be updated by a serializer after fingerprinting.
            if name.startswith("docProps/"):
                continue
            payload = package.read(member)
            if name.endswith((".xml", ".rels")) or name == "[Content_Types].xml":
                payload = _canonical_package_xml(
                    name, payload, exclude_managed_labels=exclude_managed_labels
                )
            binding.append({"name": name, "sha256": _sha256_bytes(payload)})
    return binding


def make_serialized_package_binding(
    presentation: Any, *, exclude_managed_labels: bool = False
) -> list[dict[str, str]]:
    """Serialize twice before computing a stable full-package binding."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as handle:
        temporary = Path(handle.name)
    try:
        presentation.save(os.fspath(temporary))
        normalized = Presentation(os.fspath(temporary))
        make_rendered_binding(
            normalized, exclude_managed_labels=exclude_managed_labels
        )
        normalized.save(os.fspath(temporary))
        return make_package_binding(
            temporary, exclude_managed_labels=exclude_managed_labels
        )
    finally:
        temporary.unlink(missing_ok=True)


def _resolve_spec_asset(value: Any, spec_dir: Path) -> Path | None:
    if not isinstance(value, str) or not value.strip():
        return None
    candidate = Path(value).expanduser()
    return candidate.resolve() if candidate.is_absolute() else (spec_dir / candidate).resolve()


def _slide_spec_manifest(slide: dict[str, Any], index: int, spec_dir: Path) -> dict[str, Any]:
    content = {key: value for key, value in slide.items() if key != "image"}
    image = _resolve_spec_asset(slide.get("image"), spec_dir)
    return {
        "index": index,
        "type": slide.get("type"),
        "content_sha256": _sha256_bytes(_canonical_bytes(content)),
        "image_sha256": _sha256_file(image) if image is not None else None,
    }


def _visible_slide_text(slide: Any) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            # Native panel labels may be added after the initial build. They are
            # checked geometrically by QA and intentionally excluded here so the
            # package binding remains stable across that supported post-process.
            if text and not re.fullmatch(r"[A-Z]", text, re.IGNORECASE):
                texts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
    return texts


def _enum_text(value: Any) -> str | int | float | bool | None:
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    try:
        return int(value)
    except (TypeError, ValueError):
        return str(value)


def _color_payload(color: Any) -> dict[str, Any] | None:
    if color is None:
        return None
    payload: dict[str, Any] = {"type": _enum_text(getattr(color, "type", None))}
    for key in ("rgb", "theme_color", "brightness"):
        try:
            value = getattr(color, key)
        except (AttributeError, ValueError):
            continue
        if value is not None:
            payload[key] = str(value) if key == "rgb" else _enum_text(value)
    return payload


def _fill_payload(fill: Any) -> dict[str, Any] | None:
    if fill is None:
        return None
    payload: dict[str, Any] = {"type": _enum_text(getattr(fill, "type", None))}
    try:
        payload["fore_color"] = _color_payload(fill.fore_color)
    except (AttributeError, ValueError, TypeError):
        pass
    try:
        payload["back_color"] = _color_payload(fill.back_color)
    except (AttributeError, ValueError, TypeError):
        pass
    try:
        payload["transparency"] = _enum_text(fill.transparency)
    except (AttributeError, ValueError, TypeError):
        pass
    return payload


def _font_payload(font: Any) -> dict[str, Any]:
    payload: dict[str, Any] = {}
    for key in ("name", "size", "bold", "italic", "underline"):
        try:
            value = getattr(font, key)
        except (AttributeError, ValueError):
            continue
        if value is not None:
            payload[key] = int(value) if key == "size" else _enum_text(value)
    try:
        payload["color"] = _color_payload(font.color)
    except (AttributeError, ValueError):
        pass
    return payload


def _text_frame_payload(text_frame: Any) -> dict[str, Any]:
    payload: dict[str, Any] = {"text": text_frame.text}
    for key in (
        "margin_left", "margin_right", "margin_top", "margin_bottom",
        "vertical_anchor", "word_wrap", "auto_size",
    ):
        try:
            value = getattr(text_frame, key)
        except (AttributeError, ValueError):
            continue
        if value is not None:
            payload[key] = int(value) if key.startswith("margin_") else _enum_text(value)
    paragraphs: list[dict[str, Any]] = []
    for paragraph in text_frame.paragraphs:
        item: dict[str, Any] = {
            "text": paragraph.text,
            "level": paragraph.level,
            "alignment": _enum_text(paragraph.alignment),
            "runs": [
                {"text": run.text, "font": _font_payload(run.font)}
                for run in paragraph.runs
            ],
        }
        for key in ("line_spacing", "space_before", "space_after"):
            try:
                value = getattr(paragraph, key)
            except (AttributeError, ValueError):
                continue
            if value is not None:
                item[key] = _enum_text(value)
        paragraphs.append(item)
    payload["paragraphs"] = paragraphs
    return payload


def _shape_payload(
    shape: Any, z_order: int, *, exclude_managed_labels: bool = False
) -> dict[str, Any] | None:
    """Capture deterministic visual state for one shape."""
    name = str(getattr(shape, "name", ""))
    if exclude_managed_labels and name.startswith("MJ_PANEL_LABEL_"):
        return None
    payload: dict[str, Any] = {
        "z_order": z_order,
        "name": name,
        "shape_type": _enum_text(getattr(shape, "shape_type", None)),
        "left": int(getattr(shape, "left", 0)),
        "top": int(getattr(shape, "top", 0)),
        "width": int(getattr(shape, "width", 0)),
        "height": int(getattr(shape, "height", 0)),
        "rotation": float(getattr(shape, "rotation", 0) or 0),
    }
    try:
        payload["fill"] = _fill_payload(shape.fill)
    except (AttributeError, ValueError):
        pass
    try:
        payload["line"] = {
            "fill": _fill_payload(shape.line.fill),
            "width": int(shape.line.width) if shape.line.width is not None else None,
            "dash_style": _enum_text(shape.line.dash_style),
        }
    except (AttributeError, ValueError, TypeError):
        pass
    if getattr(shape, "has_text_frame", False):
        payload["text_frame"] = _text_frame_payload(shape.text_frame)
    if getattr(shape, "has_table", False):
        payload["table"] = [
            [
                {
                    "text_frame": _text_frame_payload(cell.text_frame),
                    "fill": _fill_payload(cell.fill),
                }
                for cell in row.cells
            ]
            for row in shape.table.rows
        ]
    try:
        payload["picture_crop"] = {
            key: round(float(getattr(shape, key)), 8)
            for key in ("crop_left", "crop_right", "crop_top", "crop_bottom")
        }
    except (AttributeError, ValueError, TypeError):
        pass
    try:
        payload["image_sha256"] = _sha256_bytes(shape.image.blob)
    except (AttributeError, ValueError, TypeError):
        pass
    try:
        children = [
            child_payload
            for index, child in enumerate(shape.shapes)
            if (
                child_payload := _shape_payload(
                    child, index, exclude_managed_labels=exclude_managed_labels
                )
            ) is not None
        ]
        payload["children"] = children
    except (AttributeError, ValueError):
        pass
    # Several python-pptx property getters lazily materialize default OOXML.
    # Hash only after every parsed visual field (and nested child) has been
    # inspected so the first and subsequent fingerprints are identical.
    try:
        payload["ooxml_sha256"] = _sha256_bytes(shape._element.xml.encode("utf-8"))
    except (AttributeError, ValueError, TypeError):
        pass
    return payload


def _rendered_slide_manifest(
    slide: Any, index: int, *, exclude_managed_labels: bool = False
) -> dict[str, Any]:
    try:
        notes = slide.notes_slide.notes_text_frame.text.strip()
    except (AttributeError, ValueError):
        notes = ""
    pictures: list[str] = []
    for shape in slide.shapes:
        try:
            blob = shape.image.blob
        except (AttributeError, ValueError):
            continue
        pictures.append(_sha256_bytes(blob))
    shapes = [
        shape_payload
        for z_order, shape in enumerate(slide.shapes)
        if (
            shape_payload := _shape_payload(
                shape, z_order, exclude_managed_labels=exclude_managed_labels
            )
        ) is not None
    ]
    try:
        background = _fill_payload(slide.background.fill)
    except (AttributeError, ValueError):
        background = None
    payload = {
        "texts": _visible_slide_text(slide),
        "notes": notes,
        "pictures": sorted(pictures),
        "background": background,
        "shapes": shapes,
        "slide_flags": {
            key: slide._element.get(key)
            for key in ("show", "showMasterSp", "showMasterPh")
        },
    }
    return {
        "index": index,
        "rendered_sha256": _sha256_bytes(_canonical_bytes(payload)),
        "picture_sha256": sorted(pictures),
    }


def make_rendered_binding(
    presentation: Any, *, exclude_managed_labels: bool = False
) -> list[dict[str, Any]]:
    return [
        _rendered_slide_manifest(
            slide, index, exclude_managed_labels=exclude_managed_labels
        )
        for index, slide in enumerate(presentation.slides, start=1)
    ]


def make_serialized_rendered_binding(
    presentation: Any, *, exclude_managed_labels: bool = False
) -> list[dict[str, Any]]:
    """Fingerprint the stable OOXML form produced by python-pptx.

    Newly-created shapes can acquire relationship ids and normalized XML only
    when the package is first saved. Canonical baselines therefore cross that
    serialization boundary before their raw shape XML is hashed.
    """
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as handle:
        temporary = Path(handle.name)
    try:
        presentation.save(os.fspath(temporary))
        normalized = Presentation(os.fspath(temporary))
        normalized.save(os.fspath(temporary))
        normalized = Presentation(os.fspath(temporary))
        return make_rendered_binding(
            normalized, exclude_managed_labels=exclude_managed_labels
        )
    finally:
        temporary.unlink(missing_ok=True)


def make_spec_binding(specification: dict[str, Any], spec_path: Path) -> dict[str, Any]:
    """Return the canonical source-spec and per-slide asset SHA-256 inventory."""
    return {
        "spec_sha256": _sha256_bytes(_canonical_bytes(specification)),
        "slides": [
            _slide_spec_manifest(slide, index, spec_path.parent)
            for index, slide in enumerate(specification.get("slides", []), start=1)
        ],
    }


def make_build_manifest(
    specification: dict[str, Any],
    spec_path: Path,
    style: str,
    presentation: Any,
    package_path: Path | None = None,
) -> dict[str, Any]:
    """Create a deterministic binding between source spec/assets and built slides."""
    version_file = Path(__file__).resolve().parents[1] / "VERSION"
    try:
        skill_version = version_file.read_text(encoding="utf-8").strip()
    except OSError:
        skill_version = "unknown"
    binding = make_spec_binding(specification, spec_path)
    rendered_slides = make_rendered_binding(presentation)
    package_binding = (
        make_package_binding(package_path)
        if package_path is not None
        else make_serialized_package_binding(presentation)
    )
    return {
        "schema": MANIFEST_SCHEMA,
        "skill_name": "medical-journal-to-pptx-integrated",
        "skill_version": skill_version,
        "style": style,
        "presentation_size_emu": {
            "width": int(presentation.slide_width),
            "height": int(presentation.slide_height),
        },
        **binding,
        "rendered_slides": rendered_slides,
        "package_parts": package_binding,
    }


def _manifest_xml(manifest: dict[str, Any]) -> tuple[bytes, str]:
    payload = _canonical_bytes(manifest)
    digest = _sha256_bytes(payload)
    encoded = base64.b64encode(payload).decode("ascii")
    ET.register_namespace("", CUSTOM_PROPERTY_NS)
    ET.register_namespace("vt", VT_NS)
    root = ET.Element(f"{{{CUSTOM_PROPERTY_NS}}}Properties")
    for pid, name, value in (
        (2, "MedicalJournalBuildManifest", encoded),
        (3, "MedicalJournalBuildManifestSha256", digest),
    ):
        prop = ET.SubElement(
            root,
            f"{{{CUSTOM_PROPERTY_NS}}}property",
            {"fmtid": CUSTOM_PROPERTY_FMTID, "pid": str(pid), "name": name},
        )
        ET.SubElement(prop, f"{{{VT_NS}}}lpwstr").text = value
    xml = ET.tostring(root, encoding="utf-8", xml_declaration=True)
    return xml, digest


def manifest_sha256(manifest: dict[str, Any]) -> str:
    return _sha256_bytes(_canonical_bytes(manifest))


def _set_core_manifest_properties(presentation: Any, manifest: dict[str, Any]) -> str:
    _, digest = _manifest_xml(manifest)
    core = presentation.core_properties
    core.keywords = (
        f"medical-journal-manifest/v1;style={manifest['style']};"
        f"version={manifest['skill_version']};spec={manifest['spec_sha256']}"
    )
    core.identifier = digest
    core.category = "medical-journal-to-pptx-integrated"
    return digest


def _embed_manifest_part(pptx_path: Path, manifest: dict[str, Any]) -> None:
    """Embed the full manifest as a standards-compatible custom property part.

    Core properties carry the manifest digest and source/style summary; the
    custom-properties part carries the full per-slide/package SHA-256 inventory.
    This standard package relationship remains loadable by PowerPoint and
    LibreOffice, unlike a customXml relationship attached at the package root.
    """
    manifest_xml, _ = _manifest_xml(manifest)
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    with zipfile.ZipFile(pptx_path, "r") as source:
        rels = ET.fromstring(source.read("_rels/.rels"))
        for relation in list(rels):
            if relation.get("Type") == MANIFEST_REL_TYPE:
                rels.remove(relation)
        used_ids = {relation.get("Id") for relation in rels}
        number = 1
        while f"rIdMJ{number}" in used_ids:
            number += 1
        ET.SubElement(
            rels,
            f"{{{rel_ns}}}Relationship",
            {
                "Id": f"rIdMJ{number}",
                "Type": MANIFEST_REL_TYPE,
                "Target": MANIFEST_PART,
            },
        )

        content_types = ET.fromstring(source.read("[Content_Types].xml"))
        for override in list(content_types):
            if override.get("PartName") == f"/{MANIFEST_PART}":
                content_types.remove(override)
        ET.SubElement(
            content_types,
            f"{{{ct_ns}}}Override",
            {
                "PartName": f"/{MANIFEST_PART}",
                "ContentType": MANIFEST_CONTENT_TYPE,
            },
        )
        # LibreOffice requires the conventional default namespaces in these
        # two package-root parts; semantically equivalent ``ns0:`` output is
        # rejected as an unloadable source file.
        ET.register_namespace("", rel_ns)
        relationships_xml = ET.tostring(
            rels, encoding="utf-8", xml_declaration=True
        )
        ET.register_namespace("", ct_ns)
        content_types_xml = ET.tostring(
            content_types, encoding="utf-8", xml_declaration=True
        )

        with tempfile.NamedTemporaryFile(
            prefix=f".{pptx_path.name}.", suffix=".tmp", dir=pptx_path.parent, delete=False
        ) as handle:
            temporary = Path(handle.name)
        try:
            with zipfile.ZipFile(temporary, "w", zipfile.ZIP_DEFLATED) as target:
                for item in source.infolist():
                    if item.filename in {"_rels/.rels", "[Content_Types].xml", MANIFEST_PART}:
                        continue
                    target.writestr(item, source.read(item.filename))
                target.writestr(
                    "_rels/.rels",
                    relationships_xml,
                )
                target.writestr(
                    "[Content_Types].xml",
                    content_types_xml,
                )
                target.writestr(MANIFEST_PART, manifest_xml)
            os.replace(temporary, pptx_path)
        finally:
            if temporary.exists():
                temporary.unlink()


def validate_manifest_wiring(pptx_path: Path) -> str | None:
    """Return an error unless the manifest uses one exact standard OOXML wire."""
    try:
        with zipfile.ZipFile(pptx_path, "r") as package:
            names = [member.filename for member in package.infolist()]
            for required in (MANIFEST_PART, "_rels/.rels", "[Content_Types].xml"):
                if names.count(required) != 1:
                    return f"package must contain exactly one {required!r} member"
            relationships = ET.fromstring(package.read("_rels/.rels"))
            content_types = ET.fromstring(package.read("[Content_Types].xml"))
        if relationships.tag != f"{{{RELATIONSHIP_NS}}}Relationships":
            return "package relationships use an unsupported root namespace"
        all_relationships = list(relationships)
        identifiers = [relation.get("Id") for relation in all_relationships]
        if any(not identifier for identifier in identifiers) or len(set(identifiers)) != len(
            identifiers
        ):
            return "package relationships contain missing or duplicate Id values"
        manifest_relationships = [
            relation
            for relation in all_relationships
            if relation.get("Type") == MANIFEST_REL_TYPE
        ]
        if len(manifest_relationships) != 1:
            return "package must contain exactly one custom-properties manifest relationship"
        relation = manifest_relationships[0]
        if relation.get("Target") != MANIFEST_PART or relation.get("TargetMode") is not None:
            return "custom-properties manifest relationship has the wrong target or mode"
        conflicting_targets = [
            item
            for item in all_relationships
            if item is not relation and item.get("Target") == MANIFEST_PART
        ]
        if conflicting_targets:
            return "manifest custom-properties part has a conflicting relationship"

        if content_types.tag != f"{{{CONTENT_TYPES_NS}}}Types":
            return "package content types use an unsupported root namespace"
        manifest_overrides = [
            override
            for override in list(content_types)
            if override.get("PartName") == f"/{MANIFEST_PART}"
        ]
        if len(manifest_overrides) != 1:
            return "package must contain exactly one manifest content-type override"
        if manifest_overrides[0].get("ContentType") != MANIFEST_CONTENT_TYPE:
            return "manifest custom-properties part has the wrong content type"
    except (OSError, KeyError, zipfile.BadZipFile, ET.ParseError) as error:
        return f"manifest package wiring cannot be read: {error}"
    return None


def read_build_manifest(pptx_path: Path) -> tuple[dict[str, Any] | None, str | None]:
    """Return ``(manifest, error)`` from a built PPTX package."""
    try:
        if wiring_error := validate_manifest_wiring(pptx_path):
            return None, wiring_error
        with zipfile.ZipFile(pptx_path, "r") as package:
            raw = package.read(MANIFEST_PART)
        root = ET.fromstring(raw)
        if root.tag != f"{{{CUSTOM_PROPERTY_NS}}}Properties":
            return None, "embedded manifest custom-properties root is invalid"
        values: dict[str, str] = {}
        seen_names: set[str] = set()
        for prop in root.findall(f"{{{CUSTOM_PROPERTY_NS}}}property"):
            name = prop.get("name")
            value = prop.find(f"{{{VT_NS}}}lpwstr")
            if isinstance(name, str) and value is not None and value.text is not None:
                if name in seen_names:
                    return None, f"embedded custom property {name!r} is duplicated"
                seen_names.add(name)
                values[name] = value.text
        encoded = values.get("MedicalJournalBuildManifest")
        expected_digest = values.get("MedicalJournalBuildManifestSha256")
        if not encoded or not expected_digest:
            return None, "embedded custom properties do not contain the build manifest"
        payload = base64.b64decode(encoded, validate=True)
        digest = _sha256_bytes(payload)
        if expected_digest != digest:
            return None, "embedded build manifest digest does not match its payload"
        value = json.loads(payload.decode("utf-8"))
        if not isinstance(value, dict):
            return None, "embedded build manifest is not a JSON object"
        return value, None
    except (OSError, KeyError, zipfile.BadZipFile, ET.ParseError, ValueError, json.JSONDecodeError) as error:
        return None, f"embedded build manifest cannot be read: {error}"


def normalize_specification(spec: dict[str, Any], *, style: str) -> dict[str, Any]:
    """Preserve canonical section types while adapting aliases for each builder."""
    if style == "standard":
        slides = []
        for slide in spec.get("slides", []):
            normalized = dict(slide)
            if normalized.get("type") == "section":
                normalized["type"] = "part"
            slides.append(normalized)
        return {**spec, "slides": slides}

    if style != "nice":
        raise ValueError(f"Unknown visual style: {style}")
    slides = []
    for slide in spec.get("slides", []):
        normalized = dict(slide)
        slide_type = normalized.get("type")
        if slide_type == "part":
            normalized["type"] = "section"
            number = normalized.get("number")
            if isinstance(number, int):
                normalized["number"] = f"{number:02d}"
        elif slide_type in {"outline", "references"}:
            normalized["type"] = "content"
            normalized.setdefault("body", normalized.get("items", []))
            normalized.setdefault(
                "kicker", "LEARNING OUTLINE" if slide_type == "outline" else "KEY REFERENCES"
            )
        slides.append(normalized)
    return {**spec, "slides": slides}


def render_specification(
    specification: dict[str, Any], spec_dir: Path, *, style: str
) -> Any:
    """Render a deterministic in-memory deck used by both build and independent QA."""
    normalized = normalize_specification(specification, style=style)
    if style == "standard":
        with contextlib.redirect_stdout(sys.stderr):
            return build_deck_standard.build(
                normalized, spec_dir, require_processed_assets=False
            )
    if style == "nice":
        logo = build_deck_nice.resolve_logo(normalized.get("meta", {}), spec_dir)
        if logo is None:
            raise RuntimeError("The bundled presentation logo is missing.")
        return build_deck_nice.Builder(normalized, os.fspath(spec_dir), logo).build()
    raise ValueError(f"Unknown visual style: {style}")


def build(
    spec_path: Path,
    destination: Path,
    *,
    style: str = "standard",
    allow_unprocessed_assets: bool = False,
    overwrite: bool = False,
) -> dict[str, Any]:
    spec_path = spec_path.expanduser().resolve()
    destination = destination.expanduser().resolve()
    if destination.exists() and not overwrite:
        raise FileExistsError(
            f"Presentation already exists: {destination}. Choose a new path or use --overwrite."
        )
    specification = json.loads(spec_path.read_text(encoding="utf-8"))
    if not allow_unprocessed_assets:
        build_deck_standard.require_postprocessed_figure_assets(specification, spec_path.parent)

    presentation = render_specification(specification, spec_path.parent, style=style)

    destination.parent.mkdir(parents=True, exist_ok=True)
    # Save and validate in the destination directory, then atomically replace
    # the requested path. An interrupted overwrite must never destroy a prior
    # verified deck.
    with tempfile.NamedTemporaryFile(
        prefix=f".{destination.name}.", suffix=".pptx", dir=destination.parent, delete=False
    ) as handle:
        temporary = Path(handle.name)
    try:
        # The first save assigns/normalizes OOXML relationship ids. Build the
        # manifest only from that serialized form so its raw shape XML hashes
        # remain stable when the package is opened for QA.
        presentation.save(os.fspath(temporary))
        normalized = Presentation(os.fspath(temporary))
        # Access every lazy python-pptx visual property before the package-level
        # fingerprint is read from disk, then serialize those defaults once.
        make_rendered_binding(normalized)
        normalized.save(os.fspath(temporary))
        normalized = Presentation(os.fspath(temporary))
        manifest = make_build_manifest(
            specification, spec_path, style, normalized, package_path=temporary
        )
        manifest_digest = _set_core_manifest_properties(normalized, manifest)
        normalized.save(os.fspath(temporary))

        verified = Presentation(os.fspath(temporary))
        if make_rendered_binding(verified) != manifest["rendered_slides"]:
            raise RuntimeError(
                "Built PowerPoint visual fingerprint changed after OOXML normalization."
            )
        if make_package_binding(temporary) != manifest["package_parts"]:
            raise RuntimeError(
                "Built PowerPoint package fingerprint changed after OOXML normalization."
            )
        _embed_manifest_part(temporary, manifest)
        with zipfile.ZipFile(temporary, "r") as package:
            corrupt_member = package.testzip()
            if corrupt_member is not None:
                raise RuntimeError(f"Built PowerPoint contains a corrupt ZIP member: {corrupt_member}")
        embedded, manifest_error = read_build_manifest(temporary)
        if manifest_error or embedded != manifest:
            raise RuntimeError(
                "Built PowerPoint manifest verification failed: "
                f"{manifest_error or 'payload mismatch'}"
            )
        if make_package_binding(temporary) != manifest["package_parts"]:
            raise RuntimeError(
                "Embedded manifest changed the canonical PowerPoint package fingerprint."
            )
        os.replace(temporary, destination)
    finally:
        if temporary.exists():
            temporary.unlink()
    return {
        "pptx": str(destination),
        "style": style,
        "slides": len(normalized.slides),
        "image_box": STYLE_IMAGE_BOXES[style],
        "manifest_sha256": manifest_digest,
    }


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("spec", type=Path)
    parser.add_argument("--out", type=Path, required=True)
    parser.add_argument("--style", choices=("standard", "nice"), default="standard")
    parser.add_argument("--allow-unprocessed-assets", action="store_true")
    parser.add_argument("--overwrite", action="store_true")
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args(argv)
    payload = build(
        args.spec,
        args.out,
        style=args.style,
        allow_unprocessed_assets=args.allow_unprocessed_assets,
        overwrite=args.overwrite,
    )
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2))
    else:
        print(f"saved {payload['pptx']} ({payload['slides']} slides; {args.style} style)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
