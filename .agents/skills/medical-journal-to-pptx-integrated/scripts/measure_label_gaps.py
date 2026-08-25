#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""measure_label_gaps.py — QA / calibration helper. Renders one figure slide to
PNG (via the already-exported PDF) and measures, in inches:
  * gap_above : label glyph TOP  -> its own panel bottom
  * glyphH    : rendered label glyph height
  * gap_below : label glyph BOTTOM -> next row's panel top

Use it to verify the two gaps match your target, or to re-derive
--glyph-ratio / --center-offset-in for a different font, point size, or renderer.

Usage
-----
  # 1) export the deck to PDF first, e.g.
  #    soffice --headless --convert-to pdf --outdir qa OUT.pptx
  python measure_label_gaps.py OUT.pptx qa/OUT.pdf \
      --spec deck_spec.json --geometry panel_geometry.json \
      --figure FIGURE_KEY --dpi 200

Requirements: python-pptx, pdftoppm (poppler), numpy, Pillow.
"""
import argparse, json, subprocess, tempfile, os
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


def biggest_picture(slide):
    best = None
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            ar = int(sh.width) * int(sh.height)
            if best is None or ar > best[0]:
                best = (ar, sh)
    return best[1] if best else None


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("pdf")
    ap.add_argument("--spec", required=True)
    ap.add_argument("--geometry", required=True)
    ap.add_argument("--figure", required=True, help="geometry key for the requested figure")
    ap.add_argument("--dpi", type=int, default=200)
    a = ap.parse_args()

    spec = json.load(open(a.spec)); geom = json.load(open(a.geometry))
    prs = Presentation(a.pptx); slides = list(prs.slides)
    sslides = spec.get("slides", spec)
    idx = None; pic_rect = None
    for i, (sl, sp) in enumerate(zip(slides, sslides)):
        if sp.get("type") == "figure" and sp.get("image", "").endswith(a.figure + ".png"):
            pic = biggest_picture(sl)
            pic_rect = (Emu(pic.left).inches, Emu(pic.top).inches,
                        Emu(pic.width).inches, Emu(pic.height).inches)
            idx = i + 1; break
    if idx is None:
        raise SystemExit("figure not found in spec")
    pl, pt, pw, ph = pic_rect

    tmp = tempfile.mkdtemp()
    subprocess.run(["pdftoppm", "-png", "-r", str(a.dpi), "-f", str(idx), "-l", str(idx),
                    a.pdf, os.path.join(tmp, "p")], check=True)
    png = [f for f in os.listdir(tmp) if f.endswith(".png")][0]
    A = np.asarray(Image.open(os.path.join(tmp, png)).convert("RGB")).astype(int)
    H, W, _ = A.shape; DPI = a.dpi

    g0 = geom[a.figure][0]
    rx = pl + g0["fx_right"] * pw - 0.05
    box_l = int((rx - 0.55) * DPI); box_r = int(rx * DPI)
    ytop = int(pt * DPI); ybot = int((pt + ph) * DPI)
    strip = A[ytop:ybot, box_l:box_r, :]
    ismax = strip.max(axis=2); ismin = strip.min(axis=2)
    labelish = (strip[:, :, 2] > 120) & (strip[:, :, 2] < 235) & (ismax - ismin > 25) & (ismax < 235)
    rc = labelish.sum(axis=1); rows = [y for y, c in enumerate(rc) if c >= 2]
    runs = []; prev = -9; st = None
    for y in rows:
        if y - prev > 3:
            if st is not None: runs.append((st, prev))
            st = y
        prev = y
    if st is not None: runs.append((st, prev))
    runs = [(s + ytop, e + ytop) for s, e in runs if e - s >= 6]

    col = int((pl + 0.22 * pw) * DPI); bright = A[:, col, :].sum(axis=1) > 200
    pr = []; prev = False; s0 = 0
    for y in range(ytop, ybot + 1):
        b = bright[y] if y < H else False
        if b and not prev: s0 = y
        if (not b) and prev: pr.append((s0, y - 1))
        prev = b
    pr = [r for r in pr if r[1] - r[0] > 40]
    if not runs or len(pr) < 2:
        print("could not measure (single-row figure or detection failed)"); return
    a_bot = pr[0][1]; c_top = pr[1][0]
    cand = [r for r in runs if a_bot < (r[0] + r[1]) / 2 < c_top]
    g = cand[0] if cand else runs[0]
    print(f"{a.figure}: gap_above={(g[0]-a_bot)/DPI:.3f}in  "
          f"glyphH={(g[1]-g[0])/DPI:.3f}in  gap_below={(c_top-g[1])/DPI:.3f}in")


if __name__ == "__main__":
    main()
