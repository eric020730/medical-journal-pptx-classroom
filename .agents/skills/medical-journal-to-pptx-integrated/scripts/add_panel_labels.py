#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""add_panel_labels.py — stamp A/B/C/D panel labels onto an ALREADY-BUILT .pptx
as native text boxes at a FIXED point size. Because the text is added after the
deck is built (not burned into the image), every label has identical actual size
on every slide, regardless of how each figure was scaled to fit its image box.

Positions come from the geometry JSON produced by recompose_panels_banded.py.
Figure slides are matched to geometry entries via the deck-spec JSON: each
slide of type "figure" whose `image` basename matches a geometry key gets its
labels. The figure picture on the slide is found as the largest picture shape.

Usage
-----
  python add_panel_labels.py IN.pptx OUT.pptx \
      --spec deck_spec.json --geometry panel_geometry.json \
      --label-pt 18 --color 8FA8C8

Requirements: python-pptx.
"""
import argparse, json, math, re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE

import build_deck


def biggest_picture(slide):
    best = None
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            area = int(sh.width) * int(sh.height)
            if best is None or area > best[0]:
                best = (area, sh)
    return best[1] if best else None


def _shape_text(shape):
    return shape.text.strip() if getattr(shape, "has_text_frame", False) else ""


def _marker(figure_name, label):
    safe_figure = re.sub(r"[^A-Za-z0-9_.-]+", "_", figure_name)
    safe_label = re.sub(r"[^A-Za-z0-9_.-]+", "_", str(label))
    return f"MJ_PANEL_LABEL_{safe_figure}_{safe_label}"


def _set_shape_name(shape, name):
    """Set the OOXML non-visual shape name across python-pptx versions."""
    shape._element.nvSpPr.cNvPr.set("name", name)


def _existing_at(slide, marker, label, left, top, width, height):
    matches = []
    target_cx = int(left) + int(width) // 2
    target_cy = int(top) + int(height) // 2
    tolerance = int(Inches(0.06))
    for shape in slide.shapes:
        if getattr(shape, "name", "") == marker:
            matches.append(shape)
            continue
        if _shape_text(shape) != str(label):
            continue
        cx = int(shape.left) + int(shape.width) // 2
        cy = int(shape.top) + int(shape.height) // 2
        if abs(cx - target_cx) <= tolerance and abs(cy - target_cy) <= tolerance:
            matches.append(shape)
    return matches


def _validate_geometry(geometry):
    if not isinstance(geometry, dict):
        raise SystemExit("geometry JSON must be an object keyed by image stem")
    for name, entries in geometry.items():
        if not isinstance(entries, list):
            raise SystemExit(f"geometry[{name!r}] must be a list")
        for entry in entries:
            if not isinstance(entry, dict):
                raise SystemExit(f"geometry[{name!r}] entries must be objects")
            if not entry.get("label"):
                continue
            try:
                fx_right = float(entry["fx_right"])
                fy_center = float(entry["fy_center"])
            except (KeyError, TypeError, ValueError) as error:
                raise SystemExit(f"invalid panel geometry for {name}: {entry!r}") from error
            if not (
                math.isfinite(fx_right)
                and math.isfinite(fy_center)
                and 0.0 <= fx_right <= 1.0
                and 0.0 <= fy_center <= 1.0
            ):
                raise SystemExit(
                    f"panel geometry fractions must be finite values in [0,1]: {entry!r}"
                )


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("in_pptx")
    ap.add_argument("out_pptx")
    ap.add_argument("--spec", required=True, help="deck spec JSON (slides with type/image)")
    ap.add_argument("--geometry", required=True, help="geometry JSON from recompose_panels_banded.py")
    ap.add_argument("--label-pt", type=float, default=18.0)
    ap.add_argument("--color", default="8FA8C8", help="hex label color, no '#'")
    ap.add_argument("--box-w-in", type=float, default=0.55)
    ap.add_argument("--box-h-in", type=float, default=0.30)
    ap.add_argument("--right-pad-in", type=float, default=0.05)
    ap.add_argument("--bold", action="store_true")
    ap.add_argument("--font", default="Calibri")
    a = ap.parse_args()

    in_path = Path(a.in_pptx).expanduser().resolve(strict=True)
    out_path = Path(a.out_pptx).expanduser().resolve(strict=False)
    if in_path == out_path:
        ap.error("in_pptx and out_pptx must be different paths")
    if a.box_w_in <= 0 or a.box_h_in <= 0 or a.right_pad_in < 0:
        ap.error("label box dimensions must be positive and right padding non-negative")
    if not math.isfinite(a.label_pt) or a.label_pt <= 0:
        ap.error("--label-pt must be a finite positive number")

    geom = json.loads(Path(a.geometry).read_text(encoding="utf-8"))
    spec = json.loads(Path(a.spec).read_text(encoding="utf-8"))
    _validate_geometry(geom)
    color = RGBColor(int(a.color[0:2], 16), int(a.color[2:4], 16), int(a.color[4:6], 16))
    BOX_W, BOX_H = Inches(a.box_w_in), Inches(a.box_h_in)
    RIGHT_PAD = Inches(a.right_pad_in)

    prs = Presentation(str(in_path))
    manifest, manifest_error = build_deck.read_build_manifest(in_path)
    if manifest_error or manifest is None:
        raise SystemExit(
            "input deck is not a valid canonical build: "
            f"{manifest_error or 'missing build manifest'}"
        )
    if manifest.get("rendered_slides") != build_deck.make_rendered_binding(prs):
        raise SystemExit("input deck visual state does not match its embedded build manifest")
    if manifest.get("package_parts") != build_deck.make_package_binding(in_path):
        raise SystemExit("input deck package state does not match its embedded build manifest")
    slides = list(prs.slides)
    sslides = spec.get("slides", spec) if isinstance(spec, dict) else spec
    if len(slides) != len(sslides):
        raise SystemExit(f"slide count mismatch: pptx={len(slides)} spec={len(sslides)}")

    added = 0
    for slide, sp in zip(slides, sslides):
        if sp.get("type") != "figure":
            continue
        img = sp.get("image", "")
        name = Path(img).stem
        if name not in geom:
            continue
        if not isinstance(geom[name], list):
            raise SystemExit(f"geometry[{name!r}] must be a list")
        image_path = Path(img).expanduser()
        if not image_path.is_absolute():
            image_path = Path(a.spec).expanduser().resolve().parent / image_path
        sidecar_path = Path(str(image_path) + ".postprocess.json")
        if sidecar_path.is_file():
            sidecar = json.loads(sidecar_path.read_text(encoding="utf-8"))
            if sidecar.get("source_label_policy") == "preserve" or (
                sidecar.get("embedded_labels") and not sidecar.get("native_labels")
            ):
                continue
        pic = biggest_picture(slide)
        if pic is None:
            continue
        L, T, W, H = int(pic.left), int(pic.top), int(pic.width), int(pic.height)
        for p in geom[name]:
            if not isinstance(p, dict):
                raise SystemExit(f"geometry[{name!r}] entries must be objects")
            if not p.get("label"):
                continue
            try:
                fx_right = float(p["fx_right"])
                fy_center = float(p["fy_center"])
            except (KeyError, TypeError, ValueError) as error:
                raise SystemExit(f"invalid panel geometry for {name}: {p!r}") from error
            if not (
                math.isfinite(fx_right)
                and math.isfinite(fy_center)
                and 0.0 <= fx_right <= 1.0
                and 0.0 <= fy_center <= 1.0
            ):
                raise SystemExit(
                    f"panel geometry fractions must be finite values in [0,1]: {p!r}"
                )
            rx = L + int(fx_right * W) - int(RIGHT_PAD)
            cy = T + int(fy_center * H)
            left = rx - int(BOX_W)
            top = cy - int(BOX_H) // 2
            if (
                left < 0 or top < 0
                or left + int(BOX_W) > int(prs.slide_width)
                or top + int(BOX_H) > int(prs.slide_height)
            ):
                raise SystemExit(
                    f"panel label {p['label']!r} for {name} would fall outside the slide"
                )
            marker = _marker(name, p["label"])
            existing = _existing_at(slide, marker, p["label"], left, top, BOX_W, BOX_H)
            if len(existing) > 1:
                raise SystemExit(
                    f"duplicate native panel label {p['label']!r} already exists for {name}"
                )
            if existing:
                _set_shape_name(existing[0], marker)
                continue
            tb = slide.shapes.add_textbox(left, top, BOX_W, BOX_H)
            _set_shape_name(tb, marker)
            tf = tb.text_frame
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.RIGHT
            run = para.add_run(); run.text = p["label"]
            run.font.size = Pt(a.label_pt); run.font.bold = a.bold
            run.font.name = a.font; run.font.color.rgb = color
            added += 1

    out_path.parent.mkdir(parents=True, exist_ok=True)
    # Save once before hashing: newly-added shape XML receives its stable
    # relationship ids only at the serialization boundary.
    prs.save(str(out_path))
    normalized = Presentation(str(out_path))
    build_deck.make_rendered_binding(normalized)
    normalized.save(str(out_path))
    normalized = Presentation(str(out_path))
    manifest["rendered_slides"] = build_deck.make_rendered_binding(normalized)
    manifest["package_parts"] = build_deck.make_package_binding(out_path)
    build_deck._set_core_manifest_properties(normalized, manifest)
    normalized.save(str(out_path))
    verified = Presentation(str(out_path))
    if build_deck.make_rendered_binding(verified) != manifest["rendered_slides"]:
        raise SystemExit("panel-labelled deck fingerprint changed after OOXML normalization")
    if build_deck.make_package_binding(out_path) != manifest["package_parts"]:
        raise SystemExit("panel-labelled package fingerprint changed after OOXML normalization")
    build_deck._embed_manifest_part(out_path, manifest)
    embedded, embedded_error = build_deck.read_build_manifest(out_path)
    if embedded_error or embedded != manifest:
        raise SystemExit(
            "failed to refresh panel-labelled deck manifest: "
            f"{embedded_error or 'payload mismatch'}"
        )
    if build_deck.make_package_binding(out_path) != manifest["package_parts"]:
        raise SystemExit("embedded manifest changed the canonical package fingerprint")
    print(f"added {added} native panel labels -> {out_path}")


if __name__ == "__main__":
    main()
