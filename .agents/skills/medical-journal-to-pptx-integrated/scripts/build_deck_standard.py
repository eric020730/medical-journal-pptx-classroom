#!/usr/bin/env python3
"""Build a styled medical-journal .pptx from a deck-spec JSON file.

Usage:
    python build_deck.py <spec.json> --out <output.pptx>

The spec schema lives in ../references/deck_spec_schema.md.
Visual style constants live in ../references/visual_style.md.

Style: dark academic / near-black-navy — deep #061428 page background,
mid-navy #0F2847 header band, light-blue accents. Content slides show
a chapter progress bar under the header and a footer that splits
citation (left) and an understated zero-padded page number (right).
"""
from __future__ import annotations
import argparse
import sys
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

import vector_table

# ---------- Style constants (match references/visual_style.md) ----------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

HEADER_H = Inches(1.00)
PROGRESS_Y = Inches(1.00)
PROGRESS_H = Inches(0.10)
BODY_Y = Inches(1.30)
BODY_H = Inches(5.70)
DEFAULT_LOGO_PATH = Path(__file__).resolve().parents[1] / "assets" / "dr_leether_logo.png"
CORNER_LOGO_H = Inches(0.55)
CORNER_LOGO_RIGHT = Inches(0.25)
CORNER_LOGO_TOP = int((HEADER_H - CORNER_LOGO_H) / 2)

DEFAULT_PALETTE = {
    # Dark academic / near-black-navy palette.
    "bg_page":        "#061428",   # near-black-navy page background
    "header_navy":    "#0F2847",   # mid-navy header band
    "accent_bright":  "#A8C5E8",   # light-blue accent: titles, progress fill, stripes
    "accent_dim":     "#2A4566",   # muted mid-blue: progress bar bg, hairlines
    "text_primary":   "#FFFFFF",   # body copy on dark bg
    "text_secondary": "#8FA8C8",   # footer, page number, citation (muted blue-gray)
    "text_on_header": "#FFFFFF",   # text sitting on navy header band
    "caption":        "#8FA8C8",   # figure caption
    "hairline":       "#2A4566",   # thin divider above footer

    # Backward-compatibility aliases so older specs still work:
    # map old keys onto the closest role in the new palette.
    "bg_dark":        "#061428",
    "header_bg":      "#0F2847",
    "accent_blue":    "#A8C5E8",
    "accent_teal":    "#2A4566",
    "accent_light":   "#A8C5E8",
    "accent_pale":    "#2A4566",
    "body_text":      "#FFFFFF",
    "accent_dark":    "#061428",
    "white":          "#FFFFFF",
}

POSTPROCESS_SUFFIX = ".postprocess.json"
RASTER_SUFFIXES = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".webp", ".bmp"}


def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def require_postprocessed_figure_assets(spec: dict, spec_dir: Path) -> None:
    missing = []
    invalid = []
    image_count = 0
    for idx, slide in enumerate(spec.get("slides", []), start=1):
        if not slide.get("image"):
            continue
        image_path = Path(slide["image"])
        if not image_path.is_absolute():
            image_path = (spec_dir / image_path).resolve()
        image_count += 1
        suffix = image_path.suffix.lower()
        if suffix not in RASTER_SUFFIXES | vector_table.VECTOR_SUFFIXES:
            invalid.append(f"slide {idx}: unsupported image format {image_path.suffix!r}")
            continue
        sidecar = image_path.with_suffix(image_path.suffix + POSTPROCESS_SUFFIX)
        if not sidecar.exists():
            missing.append(f"slide {idx}: {image_path}")
            continue
        if suffix in vector_table.VECTOR_SUFFIXES:
            try:
                metadata = json.loads(sidecar.read_text(encoding="utf-8"))
            except (OSError, json.JSONDecodeError) as error:
                invalid.append(f"slide {idx}: unreadable vector sidecar ({error})")
                continue
            if not isinstance(metadata, dict):
                invalid.append(f"slide {idx}: vector sidecar must be a JSON object")
                continue
            invalid.extend(
                f"slide {idx}: {failure}"
                for failure in vector_table.sidecar_structure_failures(metadata)
            )

    manifest: Path | None = None
    meta = spec.get("meta")
    if isinstance(meta, dict):
        for key in ("extraction_manifest", "extracted_manifest"):
            value = meta.get(key)
            if isinstance(value, str) and value.strip():
                candidate = Path(value).expanduser()
                manifest = candidate.resolve() if candidate.is_absolute() else (spec_dir / candidate).resolve()
                break
    if manifest is None:
        for candidate in (
            spec_dir / "extracted" / "manifest.json",
            spec_dir.parent / "extracted" / "manifest.json",
        ):
            if candidate.is_file():
                manifest = candidate.resolve()
                break
    if image_count and (manifest is None or not manifest.is_file()):
        invalid.append(
            "all slide images require a readable extraction manifest before canonical build"
        )

    if missing or invalid:
        detail = "\n".join(f"- {item}" for item in [*missing, *invalid])
        raise RuntimeError(
            "All slide images must pass authenticated postprocessing before deck build:\n"
            + detail
        )


# ---------- Shape helpers ----------

def add_rect(slide, x, y, w, h, fill_hex):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size_pt, color_hex, bold=False,
             italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             lines=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)

    if lines is None:
        lines = text.split("\n") if isinstance(text, str) else list(text)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = hex_to_rgb(color_hex)
    return tb


def add_image_fit(slide, path, x, y, w, h):
    """Add an image into the box (x,y,w,h), preserving aspect ratio, centered."""
    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        draw_w = w
        draw_h = int(w / img_ratio)
    else:
        draw_h = h
        draw_w = int(h * img_ratio)
    draw_x = x + (w - draw_w) // 2
    draw_y = y + (h - draw_h) // 2
    pic = slide.shapes.add_picture(path, draw_x, draw_y, draw_w, draw_h)
    return pic, (draw_x, draw_y, draw_w, draw_h)


def add_image_fixed_width(slide, path, x, y, w, h, fixed_w):
    """Add an image at a fixed width inside the box, centered vertically."""
    with Image.open(path) as im:
        iw, ih = im.size
    draw_w = min(int(fixed_w), int(w))
    draw_h = int(draw_w * ih / max(1, iw))
    if draw_h > h:
        draw_h = int(h)
        draw_w = int(draw_h * iw / max(1, ih))
    draw_x = x + (w - draw_w) // 2
    draw_y = y + (h - draw_h) // 2
    pic = slide.shapes.add_picture(path, draw_x, draw_y, draw_w, draw_h)
    return pic, (draw_x, draw_y, draw_w, draw_h)


def set_notes(slide, notes_text):
    if not notes_text:
        return
    tf = slide.notes_slide.notes_text_frame
    tf.text = ""
    for i, line in enumerate(notes_text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        pos = 0
        while pos < len(line):
            start = line.find("**", pos)
            if start == -1:
                run = p.add_run()
                run.text = line[pos:]
                pos = len(line)
                break
            if start > pos:
                run = p.add_run()
                run.text = line[pos:start]
            end = line.find("**", start + 2)
            if end == -1:
                run = p.add_run()
                run.text = line[start:]
                pos = len(line)
                break
            run = p.add_run()
            run.text = line[start + 2:end]
            run.font.bold = True
            pos = end + 2


# ---------- Common frame pieces ----------

def _draw_progress_bar(slide, pal, current_part, total_parts, y, h):
    """Draw a single bright-blue divider bar under the header.

    Render this as a clean accent line rather than a two-tone progress meter.
    Keep the signature arguments so existing slide builders do not change.
    """
    add_rect(slide, 0, y, SLIDE_W, h, pal["accent_bright"])


def _draw_footer(slide, pal, footer_label, slide_number, total_slides):
    """Footer = hairline + left citation + zero-padded current page."""
    add_rect(slide, 0, Inches(7.05), SLIDE_W, Inches(0.01), pal["hairline"])
    if footer_label:
        add_text(slide, Inches(0.40), Inches(7.12), Inches(10.50), Inches(0.32),
                 footer_label, 11, pal["text_secondary"],
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    if slide_number:
        add_text(slide, Inches(11.40), Inches(7.12), Inches(1.53), Inches(0.32),
                 f"{slide_number:02d}", 11, pal["text_secondary"],
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _add_corner_logo(slide, logo_path):
    """Add the Dr. Leether logo at top-right on non-title/non-closing slides."""
    if not logo_path:
        return
    try:
        with Image.open(logo_path) as im:
            iw, ih = im.size
        logo_w = int(CORNER_LOGO_H * (iw / ih))
        x = SLIDE_W - CORNER_LOGO_RIGHT - logo_w
        slide.shapes.add_picture(
            logo_path, x, CORNER_LOGO_TOP,
            height=CORNER_LOGO_H
        )
    except Exception:
        pass


def _content_frame(slide, pal, logo_path, footer_label, title,
                   current_part, total_parts, slide_number, total_slides):
    """Shared frame: near-black bg, tall header, progress bar, footer."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, pal["bg_page"])
    add_rect(slide, 0, 0, SLIDE_W, HEADER_H, pal["header_navy"])
    _draw_progress_bar(slide, pal, current_part, total_parts, PROGRESS_Y, PROGRESS_H)

    title_x = Inches(0.55)
    title_w = Inches(12.25)
    if logo_path:
        _add_corner_logo(slide, logo_path)
        title_w = Inches(11.55)

    add_text(slide, title_x, Inches(0.12), title_w, Inches(0.78),
             title, 28, pal["text_on_header"], bold=True,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    _draw_footer(slide, pal, footer_label, slide_number, total_slides)


# ---------- Slide builders ----------

def build_title(slide, data, pal, logo_path):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, pal["bg_page"])
    # Top header band + bright accent stripe
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.10), pal["header_navy"])
    add_rect(slide, 0, Inches(1.10), SLIDE_W, Inches(0.06), pal["accent_bright"])
    # Bottom bright accent stripe + header band
    add_rect(slide, 0, Inches(6.34), SLIDE_W, Inches(0.06), pal["accent_bright"])
    add_rect(slide, 0, Inches(6.40), SLIDE_W, Inches(1.10), pal["header_navy"])

    title = (data.get("title") or "").replace("\\n", "\n")
    add_text(slide, Inches(0.80), Inches(2.00), Inches(11.70), Inches(2.20),
             title, 40, pal["accent_bright"], bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    authors = data.get("authors") or ""
    if authors:
        add_text(slide, Inches(0.50), Inches(4.40), Inches(12.30), Inches(0.60),
                 authors, 20, pal["text_primary"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    citation = data.get("citation") or ""
    if citation:
        add_text(slide, Inches(0.50), Inches(5.05), Inches(12.30), Inches(0.60),
                 citation, 16, pal["text_secondary"], italic=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    set_notes(slide, data.get("notes"))


def build_part(slide, data, pal, logo_path,
               current_part, total_parts, slide_number, total_slides):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, pal["bg_page"])
    # Top header band + progress stripe
    add_rect(slide, 0, 0, SLIDE_W, HEADER_H, pal["header_navy"])
    _draw_progress_bar(slide, pal, current_part, total_parts,
                       PROGRESS_Y, PROGRESS_H)
    _add_corner_logo(slide, logo_path)

    # Big "PART N"
    add_text(slide, Inches(0.50), Inches(1.60), Inches(12.30), Inches(2.20),
             f"PART {data.get('number', '')}", 88, pal["accent_bright"],
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Section subtitle band (mid-navy on deep-navy bg reads as a card)
    add_rect(slide, 0, Inches(4.40), SLIDE_W, Inches(1.10), pal["header_navy"])
    add_text(slide, Inches(0.50), Inches(4.55), Inches(12.30), Inches(0.80),
             data.get("title", ""), 30, pal["text_primary"],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom hairline
    add_rect(slide, 0, Inches(7.48), SLIDE_W, Inches(0.02), pal["accent_bright"])

    set_notes(slide, data.get("notes"))


def build_content(slide, data, pal, logo_path, footer_label,
                  current_part, total_parts, slide_number, total_slides):
    _content_frame(slide, pal, logo_path, footer_label, data.get("title", ""),
                   current_part, total_parts, slide_number, total_slides)

    has_image = bool(data.get("image"))
    body_w = Inches(7.60) if has_image else Inches(12.25)

    lines = data.get("body") or []
    add_text(slide, Inches(0.55), BODY_Y, body_w, BODY_H,
             "", 22, pal["text_primary"], lines=lines)

    if has_image:
        add_image_fit(slide, data["image"],
                      Inches(8.35), Inches(1.45),
                      Inches(4.70), Inches(5.30))

    set_notes(slide, data.get("notes"))


def _split_outline_item(raw):
    """Parse 'N\ufe0f\u20e3  Topic Title \u2014 Slides 3-12' into (title, range_text).
    Tolerant: strips any leading emoji/number markers, splits on the last em or
    en dash, and normalises a trailing 'Slides X-Y' into 'Slides X\u2013Y'."""
    import re
    txt = str(raw).strip()
    title, rng = txt, ""
    for sep in (" \u2014 ", " \u2013 ", "  \u2014 ", " - Slides", " \u2014Slides"):
        if sep in txt:
            left, right = txt.rsplit(sep, 1)
            title, rng = left, right
            break
    # strip leading emoji / numerals / punctuation before the first letter
    title = re.sub(r"^[^A-Za-z]+", "", title).strip()
    rng = rng.strip()
    m = re.search(r"[Ss]lides?\s*([0-9]+)\s*[-\u2013]\s*([0-9]+)", txt)
    if m:
        rng = "Slides %s\u2013%s" % (m.group(1), m.group(2))
    elif rng and not rng.lower().startswith("slide"):
        rng = ""
    return title, rng


def build_outline(slide, data, pal, logo_path, footer_label,
                  current_part, total_parts, slide_number, total_slides):
    _content_frame(slide, pal, logo_path, footer_label,
                   data.get("title", "Outline"),
                   current_part, total_parts, slide_number, total_slides)
    items = data.get("items") or []
    n = len(items) or 1

    # vertical layout: evenly spaced cards in the body region
    top = Inches(1.45)
    bottom = Inches(6.78)
    avail = int(bottom) - int(top)
    gap = Inches(0.16)
    row_h = (avail - int(gap) * (n - 1)) // n
    row_h = min(row_h, int(Inches(0.92)))
    block_h = row_h * n + int(gap) * (n - 1)
    y0 = int(top) + (avail - block_h) // 2

    band_x = Inches(0.62)
    band_w = Inches(12.10)
    badge_d = min(row_h - int(Inches(0.14)), int(Inches(0.62)))
    pill_w = Inches(2.05)
    # Adaptive sizing: scale text with row height so a long outline (many rows)
    # stays readable and uncramped while a short one keeps generous type.
    row_h_in = row_h / float(Inches(1))
    title_pt = max(11, min(20, int(round(row_h_in * 24))))
    num_pt = max(12, min(24, int(round(row_h_in * 30))))
    pill_pt = max(9, min(14, int(round(row_h_in * 17))))
    pill_h = min(int(Inches(0.46)), int(row_h * 0.62))

    for idx, raw in enumerate(items):
        title, rng = _split_outline_item(raw)
        ry = y0 + idx * (row_h + int(gap))

        # row band (rounded) for a card feel
        band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      band_x, ry, band_w, row_h)
        band.fill.solid(); band.fill.fore_color.rgb = hex_to_rgb(pal["header_navy"])
        band.line.fill.background(); band.shadow.inherit = False
        try:
            band.adjustments[0] = 0.12
        except Exception:
            pass

        # number badge (filled accent circle)
        bx = band_x + Inches(0.22)
        by = ry + (row_h - badge_d) // 2
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx, by, badge_d, badge_d)
        badge.fill.solid(); badge.fill.fore_color.rgb = hex_to_rgb(pal["accent_bright"])
        badge.line.fill.background(); badge.shadow.inherit = False
        add_text(slide, bx, by, badge_d, badge_d, str(idx + 1), num_pt,
                 pal["bg_page"], bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

        # topic title
        tx = bx + badge_d + Inches(0.32)
        tw = band_x + band_w - Inches(0.24) - pill_w - Inches(0.20) - tx
        add_text(slide, tx, ry, tw, row_h, title, title_pt, pal["text_primary"],
                 bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

        # slide-range pill
        if rng:
            px = band_x + band_w - Inches(0.24) - pill_w
            py = ry + (row_h - int(pill_h)) // 2
            pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          px, py, pill_w, pill_h)
            pill.fill.solid(); pill.fill.fore_color.rgb = hex_to_rgb(pal["accent_dim"])
            pill.line.fill.background(); pill.shadow.inherit = False
            try:
                pill.adjustments[0] = 0.5
            except Exception:
                pass
            add_text(slide, px, py, pill_w, pill_h, rng, pill_pt,
                     pal["accent_bright"], bold=True, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)

    set_notes(slide, data.get("notes"))


def _emf_aspect(path):
    """Read an EMF's true width/height aspect from its header rclBounds."""
    try:
        import struct
        with open(path, "rb") as f:
            head = f.read(40)
        x0, y0, x1, y1 = struct.unpack_from("<4i", head, 8)
        w, h = (x1 - x0), (y1 - y0)
        if w > 0 and h > 0:
            return w / h
    except Exception:
        pass
    return None


def _place_vector_with_card(slide, path, bx, by, bw, bh, aspect, card_pad):
    """Place a vector image (EMF/WMF) fitted into a box, on a white card.

    Vector tables are resolution-independent (no quality loss at any zoom) but
    have a transparent/black-text background, so on the dark slide they need a
    white card behind them. The true aspect is read from the EMF header so the
    image is never distorted; `aspect` is only a fallback.
    """
    if str(path).lower().endswith(".emf"):
        true_aspect = _emf_aspect(path)
        if true_aspect:
            aspect = true_aspect
    box_aspect = bw / bh
    if box_aspect > aspect:
        draw_h = int(bh); draw_w = int(bh * aspect)
    else:
        draw_w = int(bw); draw_h = int(bw / aspect)
    draw_x = int(bx + (bw - draw_w) // 2)
    draw_y = int(by + (bh - draw_h) // 2)
    cx = max(0, draw_x - card_pad); cy = max(0, draw_y - card_pad)
    add_rect(slide, cx, cy, draw_w + 2 * card_pad, draw_h + 2 * card_pad, "#FFFFFF")
    slide.shapes.add_picture(path, draw_x, draw_y, width=draw_w, height=draw_h)
    return draw_x, draw_y, draw_w, draw_h


def build_figure(slide, data, pal, logo_path, footer_label,
                 current_part, total_parts, slide_number, total_slides):
    _content_frame(slide, pal, logo_path, footer_label, data.get("title", ""),
                   current_part, total_parts, slide_number, total_slides)
    image = data.get("image")
    caption = data.get("caption") or ""
    panel_labels = data.get("panel_labels") or []

    # Panel-geometry gate: a labeled multi-panel figure rendered as a single
    # composite image cannot have its per-panel label positions inferred.
    # Require explicit geometry (panel_label_x_fracs or panel_boxes) instead of
    # silently distributing labels evenly across the image.
    x_fracs = data.get("panel_label_x_fracs") or []
    panel_boxes = data.get("panel_boxes") or []
    if len(panel_labels) > 1 and not data.get("panel_geometry_exception"):
        has_geometry = (
            len(x_fracs) >= len(panel_labels)
            or len(panel_boxes) >= len(panel_labels)
        )
        if not has_geometry:
            raise RuntimeError(
                "Figure slide %r has %d panel_labels but no panel geometry. "
                "Provide 'panel_label_x_fracs' (one per label) or 'panel_boxes', "
                "or set 'panel_geometry_exception': true with a documented reason. "
                "Silently distributing labels across a single composite image is "
                "not allowed by the panel-geometry quality gate."
                % (data.get("title", ""), len(panel_labels))
            )

    image_box_h = Inches(4.70) if panel_labels else Inches(4.85)
    if image:
        box = (Inches(0.60), Inches(1.40), Inches(12.10), image_box_h)
        if str(image).lower().endswith((".emf", ".wmf")):
            draw_x, draw_y, draw_w, draw_h = _place_vector_with_card(
                slide, image, *box,
                aspect=float(data.get("image_aspect") or 1.5),
                card_pad=Inches(0.10),
            )
        elif data.get("image_width_in"):
            _, (draw_x, draw_y, draw_w, draw_h) = add_image_fixed_width(
                slide, image, *box, Inches(float(data["image_width_in"]))
            )
        else:
            _, (draw_x, draw_y, draw_w, draw_h) = add_image_fit(
                slide, image, *box
            )
        if panel_labels:
            label_y = min(draw_y + draw_h + Inches(0.03), Inches(6.10))
            label_h = Inches(0.22)
            label_w = Inches(0.26)
            n = len(panel_labels)
            for idx, label in enumerate(panel_labels):
                if idx < len(x_fracs):
                    right_frac = float(x_fracs[idx])
                elif idx < len(panel_boxes) and isinstance(panel_boxes[idx], dict):
                    right_frac = float(
                        panel_boxes[idx].get("right_x_frac", (idx + 1) / max(1, n))
                    )
                else:
                    right_frac = (idx + 1) / max(1, n)
                label_x = draw_x + int(draw_w * right_frac) - label_w
                add_text(slide, label_x, label_y, label_w, label_h,
                         str(label), 12, pal["caption"], bold=True,
                         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    if caption:
        caption_y = Inches(6.40) if panel_labels else Inches(6.35)
        add_text(slide, Inches(0.60), caption_y, Inches(12.10), Inches(0.55),
                 caption, 13, pal["caption"], italic=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    set_notes(slide, data.get("notes"))


def build_references(slide, data, pal, logo_path, footer_label,
                     current_part, total_parts, slide_number, total_slides):
    _content_frame(slide, pal, logo_path, footer_label,
                   data.get("title", "Key References"),
                   current_part, total_parts, slide_number, total_slides)
    items = data.get("items") or []
    add_text(slide, Inches(0.55), Inches(1.25), Inches(12.25), Inches(5.70),
             "", 15, pal["text_primary"], lines=items)
    set_notes(slide, data.get("notes"))


def build_thanks(slide, data, pal, logo_path):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, pal["bg_page"])
    # Same bands as title
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.10), pal["header_navy"])
    add_rect(slide, 0, Inches(1.10), SLIDE_W, Inches(0.06), pal["accent_bright"])
    add_rect(slide, 0, Inches(6.34), SLIDE_W, Inches(0.06), pal["accent_bright"])
    add_rect(slide, 0, Inches(6.40), SLIDE_W, Inches(1.10), pal["header_navy"])

    add_text(slide, Inches(0.80), Inches(2.00), Inches(11.70), Inches(1.80),
             data.get("title", "Thank You"), 60, pal["accent_bright"],
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    subtitle = data.get("subtitle") or ""
    if subtitle:
        add_text(slide, Inches(0.50), Inches(4.00), Inches(12.30), Inches(0.70),
                 subtitle, 22, pal["text_primary"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    citation = data.get("citation") or ""
    if citation:
        add_text(slide, Inches(0.50), Inches(4.90), Inches(12.30), Inches(0.60),
                 citation, 15, pal["text_secondary"], italic=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    set_notes(slide, data.get("notes"))


# ---------- Driver ----------

def build(spec: dict, spec_dir: Path, require_processed_assets: bool = True) -> Presentation:
    if require_processed_assets:
        require_postprocessed_figure_assets(spec, spec_dir)

    meta = spec.get("meta", {})
    pal = dict(DEFAULT_PALETTE)
    for k in ("bg_page", "header_navy", "accent_bright", "accent_dim",
              "text_primary", "text_secondary", "text_on_header",
              "caption", "hairline",
              "accent_dark", "accent_blue", "accent_teal", "accent_light",
              "accent_pale"):
        if k in meta:
            pal[k] = meta[k]

    logo_path = meta.get("logo_path")
    if logo_path:
        if not Path(logo_path).is_absolute():
            logo_path = str((spec_dir / logo_path).resolve())
        if not Path(logo_path).exists():
            # Do NOT silently disable the logo. Warn and fall back to the
            # bundled default so non-title/non-thanks slides still brand.
            print(
                "WARNING: meta.logo_path %r does not exist; falling back to "
                "bundled default logo." % logo_path,
                file=sys.stderr,
            )
            logo_path = str(DEFAULT_LOGO_PATH) if DEFAULT_LOGO_PATH.exists() else None
    elif DEFAULT_LOGO_PATH.exists():
        logo_path = str(DEFAULT_LOGO_PATH)

    if logo_path:
        print("logo: using %s" % logo_path, file=sys.stderr)
    else:
        print(
            "WARNING: no logo available (meta.logo_path missing/invalid and "
            "bundled default not found); slides will render without a logo.",
            file=sys.stderr,
        )

    footer_label = meta.get("footer_label", "")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slides = spec.get("slides", [])
    total_slides = len(slides)
    total_parts = sum(1 for s in slides if s.get("type") == "part")
    current_part = 0

    for i, s in enumerate(slides, start=1):
        t = s.get("type")
        if t == "part":
            current_part += 1

        for key in ("image",):
            if key in s and s[key] and not Path(s[key]).is_absolute():
                s[key] = str((spec_dir / s[key]).resolve())

        slide = prs.slides.add_slide(blank)

        if t == "title":
            build_title(slide, s, pal, logo_path)
        elif t == "part":
            build_part(slide, s, pal, logo_path,
                       current_part, total_parts, i, total_slides)
        elif t == "content":
            build_content(slide, s, pal, logo_path, footer_label,
                          current_part, total_parts, i, total_slides)
        elif t == "outline":
            build_outline(slide, s, pal, logo_path, footer_label,
                          current_part, total_parts, i, total_slides)
        elif t == "figure":
            build_figure(slide, s, pal, logo_path, footer_label,
                         current_part, total_parts, i, total_slides)
        elif t == "references":
            build_references(slide, s, pal, logo_path, footer_label,
                             current_part, total_parts, i, total_slides)
        elif t == "thanks":
            build_thanks(slide, s, pal, logo_path)
        else:
            raise ValueError(f"slide {i}: unknown type {t!r}")

    return prs


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("spec", help="Path to deck spec JSON")
    ap.add_argument("--out", required=True, help="Output .pptx path")
    ap.add_argument(
        "--allow-unprocessed-assets",
        action="store_true",
        help="Bypass the postprocess sidecar gate for isolated debugging only.",
    )
    args = ap.parse_args()

    spec_path = Path(args.spec).resolve()
    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    prs = build(spec, spec_path.parent, require_processed_assets=not args.allow_unprocessed_assets)
    prs.save(args.out)
    print(f"Wrote {args.out}  ({len(spec.get('slides', []))} slides)")


if __name__ == "__main__":
    main()
