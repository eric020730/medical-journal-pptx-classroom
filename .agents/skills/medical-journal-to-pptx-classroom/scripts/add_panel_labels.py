#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""add_panel_labels.py — stamp A/B/C/D panel labels onto an ALREADY-BUILT .pptx
as native text boxes at a FIXED point size. Because the text is added after the
deck is built (not burned into the image), every label has identical actual size
on every slide, regardless of how each figure was scaled to fit its image box.

Positions come from the geometry JSON produced by recompose_panels_banded.py.
Figure slides are matched to geometry entries via the deck-spec JSON: each
slide of type "figure" whose `image` basename matches a geometry key gets its
labels. The figure picture on the slide is found as the largest picture shape.

Usage
-----
  python add_panel_labels.py IN.pptx OUT.pptx \
      --spec deck_spec.json --geometry panel_geometry.json \
      --label-pt 18 --color 8FA8C8

Requirements: python-pptx.
"""
import argparse, json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE


def biggest_picture(slide):
    best = None
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            area = int(sh.width) * int(sh.height)
            if best is None or area > best[0]:
                best = (area, sh)
    return best[1] if best else None


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("in_pptx")
    ap.add_argument("out_pptx")
    ap.add_argument("--spec", required=True, help="deck spec JSON (slides with type/image)")
    ap.add_argument("--geometry", required=True, help="geometry JSON from recompose_panels_banded.py")
    ap.add_argument("--label-pt", type=float, default=18.0)
    ap.add_argument("--color", default="8FA8C8", help="hex label color, no '#'")
    ap.add_argument("--box-w-in", type=float, default=0.55)
    ap.add_argument("--box-h-in", type=float, default=0.30)
    ap.add_argument("--right-pad-in", type=float, default=0.05)
    ap.add_argument("--bold", action="store_true")
    ap.add_argument("--font", default="Calibri")
    a = ap.parse_args()

    geom = json.load(open(a.geometry))
    spec = json.load(open(a.spec))
    color = RGBColor(int(a.color[0:2], 16), int(a.color[2:4], 16), int(a.color[4:6], 16))
    BOX_W, BOX_H = Inches(a.box_w_in), Inches(a.box_h_in)
    RIGHT_PAD = Inches(a.right_pad_in)

    prs = Presentation(a.in_pptx)
    slides = list(prs.slides)
    sslides = spec.get("slides", spec) if isinstance(spec, dict) else spec
    if len(slides) != len(sslides):
        raise SystemExit(f"slide count mismatch: pptx={len(slides)} spec={len(sslides)}")

    added = 0
    for slide, sp in zip(slides, sslides):
        if sp.get("type") != "figure":
            continue
        img = sp.get("image", "")
        name = img.split("/")[-1].rsplit(".", 1)[0]
        if name not in geom:
            continue
        pic = biggest_picture(slide)
        if pic is None:
            continue
        L, T, W, H = int(pic.left), int(pic.top), int(pic.width), int(pic.height)
        for p in geom[name]:
            if not p.get("label"):
                continue
            rx = L + int(p["fx_right"] * W) - int(RIGHT_PAD)
            cy = T + int(p["fy_center"] * H)
            left = rx - int(BOX_W)
            top = cy - int(BOX_H) // 2
            tb = slide.shapes.add_textbox(left, top, BOX_W, BOX_H)
            tf = tb.text_frame
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.RIGHT
            run = para.add_run(); run.text = p["label"]
            run.font.size = Pt(a.label_pt); run.font.bold = a.bold
            run.font.name = a.font; run.font.color.rgb = color
            added += 1

    prs.save(a.out_pptx)
    print(f"added {added} native panel labels -> {a.out_pptx}")


if __name__ == "__main__":
    main()
