from __future__ import annotations

import importlib.util
import json
import os
import shutil
import subprocess
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest import mock

from PIL import Image, ImageDraw
import pptx
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE

from test_advanced_qa import (
    add_figure,
    create_synthetic_inversion_fixture,
    full_spec,
    patterned_image,
    write_asset,
    write_spec,
)
import package_release


PROJECT_ROOT = Path(__file__).resolve().parents[1]
SKILL_NAME = "medical-journal-to-pptx-integrated"
SKILL_ROOT = PROJECT_ROOT / ".agents" / "skills" / SKILL_NAME
SKILL_RUNNER = SKILL_ROOT / "scripts" / "run.py"
INSTALLER = PROJECT_ROOT / "install-global.py"
EXISTING_SKILLS = (
    "existing-medical-skill-a",
    "existing-medical-skill-b",
    "existing-medical-skill-c",
)


def invoke(*arguments: object, cwd: Path | None = None) -> subprocess.CompletedProcess[str]:
    environment = dict(os.environ)
    environment["PYTHONUTF8"] = "1"
    environment.pop("MEDICAL_JOURNAL_PPTX_PYTHON", None)
    return subprocess.run(
        [sys.executable, *(str(argument) for argument in arguments)],
        cwd=cwd or PROJECT_ROOT,
        env=environment,
        capture_output=True,
        text=True,
        encoding="utf-8",
        check=False,
    )


def load_installer():
    spec = importlib.util.spec_from_file_location("medical_journal_integrated_installer", INSTALLER)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


class IntegratedSkillStructureTests(unittest.TestCase):
    def test_skill_identity_version_and_concise_entrypoint(self) -> None:
        content = (SKILL_ROOT / "SKILL.md").read_text(encoding="utf-8")
        self.assertIn(f"name: {SKILL_NAME}", content)
        self.assertLess(len(content.splitlines()), 110)
        self.assertIn("40–55 slides", content)
        self.assertIn("`full` is the only supported content mode", content)
        self.assertIn("--mode full", content)
        self.assertEqual(
            (SKILL_ROOT / "VERSION").read_text(encoding="utf-8").strip(),
            "v4.0.0",
        )

    def test_skill_preserves_all_image_building_and_quality_helpers(self) -> None:
        expected = {
            "extract_from_pdf.py", "postprocess_assets.py", "crop_vector_figure.py",
            "recompose_panels_aligned.py", "recompose_panels_banded.py", "add_panel_labels.py",
            "measure_label_gaps.py", "build_deck.py", "build_deck_standard.py",
            "build_deck_nice.py", "image_polarity.py", "deck_quality.py", "qa_gate.py",
            "workflow.py", "run.py", "test_panel_layout.py",
        }
        actual = {path.name for path in (SKILL_ROOT / "scripts").glob("*.py")}
        self.assertFalse(expected.difference(actual), msg=str(expected.difference(actual)))

    def test_slide_aware_panel_layout_passes_synthetic_regressions(self) -> None:
        result = invoke(SKILL_ROOT / "scripts" / "test_panel_layout.py")
        self.assertEqual(result.returncode, 0, msg=result.stderr)
        self.assertRegex(result.stderr, r"Ran \d+ tests")

    def test_skill_is_independent_of_classroom_project_configuration(self) -> None:
        for path in (SKILL_ROOT / "scripts").glob("*.py"):
            self.assertNotIn(".classroom-project.json", path.read_text(encoding="utf-8"), msg=str(path))

    def test_full_reference_and_both_visual_references_are_bundled(self) -> None:
        for name in (
            "full_workflow.md", "visual_style.md", "visual_style_nice.md",
            "quality_gates.md", "script_quality_expectations.md",
            "deck_spec_schema.md", "notes_style.md",
        ):
            with self.subTest(reference=name):
                self.assertTrue((SKILL_ROOT / "references" / name).is_file())

    def test_runtime_preserves_current_virtual_environment_executable(self) -> None:
        spec = importlib.util.spec_from_file_location("medical_journal_integrated_runner", SKILL_RUNNER)
        assert spec is not None and spec.loader is not None
        module = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(module)
        with mock.patch.object(module, "runtime_python", return_value=Path("missing-python")):
            self.assertEqual(module.choose_python(), Path(sys.executable).absolute())


class GlobalInstallationTests(unittest.TestCase):
    def test_portable_runner_generates_patient_free_demo_without_repository(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            output = root / "clean workspace" / "synthetic-demo.pdf"
            generated = invoke(SKILL_RUNNER, "demo", "--out", output, "--json", cwd=root)
            self.assertEqual(generated.returncode, 0, msg=generated.stderr)
            payload = json.loads(generated.stdout)
            self.assertEqual(Path(payload["pdf"]).resolve(), output.resolve())
            self.assertTrue(payload["synthetic"])
            self.assertFalse(payload["patient_data"])
            self.assertEqual(output.read_bytes()[:5], b"%PDF-")
            repeated = invoke(SKILL_RUNNER, "demo", "--out", output, cwd=root)
            self.assertNotEqual(repeated.returncode, 0)
            self.assertIn("already exists", repeated.stderr)

    def test_macos_and_windows_runtime_paths_are_user_scoped(self) -> None:
        installer = load_installer()
        with tempfile.TemporaryDirectory() as temporary:
            cache = Path(temporary) / "portable cache"
            unix = installer.runtime_directory(
                operating_system="Darwin", environ={"XDG_CACHE_HOME": str(cache)}
            )
            windows = installer.runtime_directory(
                operating_system="Windows", environ={"LOCALAPPDATA": str(cache)}
            )
            self.assertEqual(unix, cache / SKILL_NAME / "venv")
            self.assertEqual(windows, cache / SKILL_NAME / "venv")

    def test_codex_home_and_default_global_directories_are_supported(self) -> None:
        installer = load_installer()
        with tempfile.TemporaryDirectory() as temporary:
            configured = Path(temporary) / "codex-global"
            self.assertEqual(
                installer.global_skills_directory({"CODEX_HOME": str(configured)}),
                configured / "skills",
            )
        self.assertEqual(installer.global_skills_directory({}), Path.home() / ".agents" / "skills")

    def test_install_upgrade_and_uninstall_preserve_all_existing_skills(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            skills = root / "global skills"
            for name in EXISTING_SKILLS:
                existing = skills / name
                existing.mkdir(parents=True)
                (existing / "preserve.txt").write_text("keep this skill", encoding="utf-8")

            installed = invoke(INSTALLER, "install", "--target", skills, "--skip-deps", "--json", cwd=root)
            self.assertEqual(installed.returncode, 0, msg=installed.stderr)
            payload = json.loads(installed.stdout)
            destination = skills / SKILL_NAME
            self.assertEqual(Path(payload["install_path"]), destination)
            self.assertEqual(payload["other_skills_removed"], [])

            unrelated = root / "unrelated workspace"
            unrelated.mkdir()
            doctor = invoke(destination / "scripts" / "run.py", "doctor", "--json", cwd=unrelated)
            self.assertEqual(doctor.returncode, 0, msg=doctor.stderr)
            self.assertEqual(Path(json.loads(doctor.stdout)["skill_root"]).resolve(), destination.resolve())

            upgraded = invoke(INSTALLER, "upgrade", "--target", skills, "--skip-deps", "--json", cwd=root)
            self.assertEqual(upgraded.returncode, 0, msg=upgraded.stderr)
            self.assertEqual(json.loads(upgraded.stdout)["action"], "upgrade")

            removed = invoke(INSTALLER, "uninstall", "--target", skills, "--json", cwd=root)
            self.assertEqual(removed.returncode, 0, msg=removed.stderr)
            self.assertFalse(destination.exists())
            for name in EXISTING_SKILLS:
                self.assertTrue((skills / name / "preserve.txt").is_file(), msg=name)

    def test_install_refuses_to_overwrite_an_existing_integrated_skill(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            skills = Path(temporary) / "skills"
            first = invoke(INSTALLER, "install", "--target", skills, "--skip-deps")
            self.assertEqual(first.returncode, 0, msg=first.stderr)
            second = invoke(INSTALLER, "install", "--target", skills, "--skip-deps")
            self.assertNotEqual(second.returncode, 0)
            self.assertIn("Use the upgrade command", second.stderr)

    def test_uninstall_refuses_an_unrelated_same_name_directory(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            destination = Path(temporary) / SKILL_NAME
            destination.mkdir()
            (destination / "SKILL.md").write_text("---\nname: unrelated-skill\n---\n", encoding="utf-8")
            removed = invoke(INSTALLER, "uninstall", "--target", Path(temporary))
            self.assertNotEqual(removed.returncode, 0)
            self.assertTrue(destination.is_dir())


class FullDeckVisualStyleIntegrationTests(unittest.TestCase):
    def check_combination(self, mode: str, style: str, expected_slides: int) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            workspace = Path(temporary) / "independent workspace"
            result = invoke(
                SKILL_RUNNER, "smoke-test", "--workspace", workspace,
                "--mode", mode, "--style", style, "--keep", "--json",
                cwd=Path(temporary),
            )
            self.assertEqual(result.returncode, 0, msg=result.stderr)
            payload = json.loads(result.stdout)
            self.assertEqual(payload["slides"], expected_slides)
            self.assertTrue(payload["prebuild_qa"])
            self.assertTrue(payload["postbuild_qa"])
            output = Path(payload["work_dir"]) / f"synthetic_{style}_{mode}.pptx"
            self.assertTrue(output.is_file())
            presentation = Presentation(str(output))
            self.assertEqual(len(presentation.slides), expected_slides)
            self.assertTrue(all(slide.notes_slide.notes_text_frame.text for slide in presentation.slides))

    def test_standard_full_generates_verified_bilingual_deck(self) -> None:
        self.check_combination("full", "standard", 40)

    def test_nice_full_accepts_canonical_outline_references_and_part_slides(self) -> None:
        self.check_combination("full", "nice", 40)

    def test_smoke_test_defaults_to_the_complete_full_deck(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            result = invoke(
                SKILL_RUNNER, "smoke-test", "--workspace", Path(temporary),
                "--style", "standard", "--json",
                cwd=Path(temporary),
            )
            self.assertEqual(result.returncode, 0, msg=result.stderr)
            payload = json.loads(result.stdout)
            self.assertEqual(payload["mode"], "full")
            self.assertEqual(payload["slides"], 40)

    def test_unsupported_content_mode_is_rejected_by_every_public_entrypoint(self) -> None:
        commands = (
            (SKILL_RUNNER, "init-run", "missing.pdf", "--mode", "unsupported"),
            (SKILL_RUNNER, "prepare", "missing.pdf", "--mode", "unsupported"),
            (SKILL_RUNNER, "qa-spec", "missing.json", "--mode", "unsupported"),
            (
                SKILL_RUNNER, "qa", "missing.pptx", "--spec", "missing.json",
                "--mode", "unsupported",
            ),
            (
                SKILL_RUNNER, "build", "missing.json", "--out", "missing.pptx",
                "--mode", "unsupported",
            ),
            (SKILL_RUNNER, "smoke-test", "--mode", "unsupported"),
            (
                SKILL_ROOT / "scripts" / "qa_gate.py", "spec", "missing.json",
                "--mode", "unsupported",
            ),
            (
                SKILL_ROOT / "scripts" / "qa_check.py", "missing.json", "--spec-only",
                "--mode", "unsupported",
            ),
            (
                SKILL_ROOT / "scripts" / "deck_quality.py", "spec", "missing.json",
                "--content-mode", "unsupported",
            ),
        )
        for command in commands:
            with self.subTest(command=[str(value) for value in command]):
                rejected = invoke(*command)
                self.assertNotEqual(rejected.returncode, 0)
                self.assertIn("invalid choice: 'unsupported'", rejected.stderr)
                self.assertIn("{full}", rejected.stderr)

    def check_native_panel_labels(self, style: str) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            _, _, rendered = create_synthetic_inversion_fixture(root)
            first = write_asset(root, "panel-A.png", source=str(rendered))
            second = write_asset(root, "panel-B.png", source=str(rendered))
            assets = root / "final_assets"
            composite = assets / "Figure_01.png"
            geometry = root / "panel_geometry.json"
            width, height = ("12.13", "4.95") if style == "nice" else ("12.10", "4.85")
            recomposed = invoke(
                SKILL_ROOT / "scripts" / "recompose_panels_banded.py", composite,
                "--inputs", first, second, "--cols", "2", "--labels", "A,B",
                "--geometry", geometry, "--slide-box-w-in", width, "--slide-box-h-in", height,
            )
            self.assertEqual(recomposed.returncode, 0, msg=recomposed.stderr)
            provenance = json.loads(composite.with_suffix(".png.postprocess.json").read_text())
            self.assertEqual(provenance["source_inputs"], [str(first), str(second)])

            spec = full_spec()
            add_figure(
                spec, composite, caption="Figure 1. Synthetic multi-panel example.",
                notes="【圖片說明】🖼️【A 圖】與【B 圖】為虛構教學影像。",
            )
            spec_path = write_spec(root, spec)
            built = root / f"{style}-before-labels.pptx"
            build_result = invoke(
                SKILL_RUNNER, "build", spec_path, "--out", built,
                "--mode", "full", "--style", style,
            )
            self.assertEqual(build_result.returncode, 0, msg=build_result.stderr)
            unfinished = invoke(
                SKILL_RUNNER, "qa", built, "--spec", spec_path,
                "--mode", "full", "--style", style, "--json",
            )
            self.assertNotEqual(unfinished.returncode, 0)
            self.assertTrue(any(
                "missing visible native panel" in failure
                for failure in json.loads(unfinished.stdout)["failures"]
            ))
            labeled = root / f"{style}-native-labels.pptx"
            stamped = invoke(
                SKILL_ROOT / "scripts" / "add_panel_labels.py", built, labeled,
                "--spec", spec_path, "--geometry", geometry, "--label-pt", "18",
            )
            self.assertEqual(stamped.returncode, 0, msg=stamped.stderr)
            figure_slide = Presentation(str(labeled)).slides[3]
            labels = {
                shape.text_frame.text: shape.text_frame.paragraphs[0].runs[0].font.size.pt
                for shape in figure_slide.shapes
                if getattr(shape, "has_text_frame", False)
                and shape.text_frame.text in {"A", "B"}
            }
            self.assertEqual(labels, {"A": 18.0, "B": 18.0})
            final = invoke(
                SKILL_RUNNER, "qa", labeled, "--spec", spec_path,
                "--mode", "full", "--style", style, "--json",
            )
            self.assertEqual(final.returncode, 0, msg=final.stderr)

    def test_standard_supports_fixed_size_native_a_b_panel_labels(self) -> None:
        self.check_native_panel_labels("standard")

    def test_nice_supports_fixed_size_native_a_b_panel_labels(self) -> None:
        self.check_native_panel_labels("nice")

    def check_preserved_embedded_panel_labels(self, style: str) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            _, _, rendered = create_synthetic_inversion_fixture(root)
            inputs = []
            for label in ("A", "B"):
                panel = write_asset(
                    root,
                    f"embedded-{label}.png",
                    source=str(rendered),
                    source_label_placement="embedded",
                    embedded_label=label,
                )
                with Image.open(panel) as source:
                    image = source.convert("RGB")
                ImageDraw.Draw(image).text((8, image.height - 20), label, fill="white")
                image.save(panel)
                inputs.append(panel)

            composite = root / "final_assets" / "Figure_01.png"
            geometry = root / "panel_geometry.json"
            recomposed = invoke(
                SKILL_ROOT / "scripts" / "recompose_panels_banded.py",
                composite,
                "--inputs",
                *inputs,
                "--labels",
                "A,B",
                "--geometry",
                geometry,
            )
            self.assertEqual(recomposed.returncode, 0, msg=recomposed.stderr)
            metadata = json.loads(composite.with_suffix(".png.postprocess.json").read_text())
            self.assertEqual(metadata["source_label_policy"], "preserve")
            self.assertFalse(metadata["native_labels"])
            self.assertEqual(metadata["embedded_labels"], ["A", "B"])
            self.assertEqual(json.loads(geometry.read_text())[composite.stem], [])

            spec = full_spec()
            add_figure(
                spec,
                composite,
                caption="Figure 1. Synthetic embedded source labels.",
                notes="【圖片說明】🖼️【A 圖】與【B 圖】保留原始影像內標示。",
            )
            spec_path = write_spec(root, spec)
            output = root / f"{style}-preserved-embedded-labels.pptx"
            built = invoke(
                SKILL_RUNNER,
                "build",
                spec_path,
                "--out",
                output,
                "--mode",
                "full",
                "--style",
                style,
            )
            self.assertEqual(built.returncode, 0, msg=built.stderr)
            final = invoke(
                SKILL_RUNNER,
                "qa",
                output,
                "--spec",
                spec_path,
                "--mode",
                "full",
                "--style",
                style,
                "--json",
            )
            self.assertEqual(final.returncode, 0, msg=final.stderr or final.stdout)

            geometry.write_text(json.dumps({composite.stem: [
                {"label": "A", "fx_right": 0.5, "fy_center": 0.8},
                {"label": "B", "fx_right": 1.0, "fy_center": 0.8},
            ]}))
            stamped = root / f"{style}-preserved-no-duplicate.pptx"
            result = invoke(
                SKILL_ROOT / "scripts" / "add_panel_labels.py",
                output,
                stamped,
                "--spec",
                spec_path,
                "--geometry",
                geometry,
            )
            self.assertEqual(result.returncode, 0, msg=result.stderr)
            self.assertIn("added 0 native panel labels", result.stdout)
            duplicates = [
                shape.text_frame.text
                for shape in Presentation(str(stamped)).slides[3].shapes
                if getattr(shape, "has_text_frame", False)
                and shape.text_frame.text in {"A", "B"}
            ]
            self.assertEqual(duplicates, [])

    def test_standard_preserves_embedded_labels_without_duplicate_native_text(self) -> None:
        self.check_preserved_embedded_panel_labels("standard")

    def test_nice_preserves_embedded_labels_without_duplicate_native_text(self) -> None:
        self.check_preserved_embedded_panel_labels("nice")

    def check_emf_vector_table(self, style: str) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            asset = root / "final_assets" / "Table_1.emf"
            asset.parent.mkdir()
            bundled_emf = Path(pptx.__file__).parent / "templates" / "generic-icon.emf"
            self.assertTrue(bundled_emf.is_file())
            shutil.copyfile(bundled_emf, asset)
            spec = full_spec()
            add_figure(spec, asset, caption="Table 1. Synthetic vector results.")
            spec_path = write_spec(root, spec)
            output = root / f"{style}-vector-table.pptx"
            built = invoke(
                SKILL_RUNNER, "build", spec_path, "--out", output,
                "--mode", "full", "--style", style,
            )
            self.assertEqual(built.returncode, 0, msg=built.stderr)
            slide = Presentation(str(output)).slides[3]
            vector_pictures = [
                shape for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                and shape.image.ext in {"emf", "wmf"}
            ]
            white_cards = [
                shape for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and shape.fill.type == MSO_FILL.SOLID
                and str(shape.fill.fore_color.rgb) == "FFFFFF"
            ]
            self.assertEqual(len(vector_pictures), 1)
            self.assertTrue(white_cards)
            final = invoke(
                SKILL_RUNNER, "qa", output, "--spec", spec_path,
                "--mode", "full", "--style", style, "--json",
            )
            self.assertEqual(final.returncode, 0, msg=final.stderr)

    def test_standard_preserves_emf_vector_table_on_white_card(self) -> None:
        self.check_emf_vector_table("standard")

    def test_nice_preserves_emf_vector_table_on_white_card(self) -> None:
        self.check_emf_vector_table("nice")

    def check_specified_panel_labels(self, style: str) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            _, _, rendered = create_synthetic_inversion_fixture(root)
            figure = write_asset(root, "Figure_01.png", source=str(rendered))
            spec = full_spec()
            add_figure(
                spec, figure, caption="Figure 1. Synthetic A/B panels.",
                panel_labels=["A", "B"], panel_label_x_fracs=[0.45, 0.95],
                notes="【圖片說明】🖼️【A 圖】與【B 圖】為虛構教學影像。",
            )
            spec_path = write_spec(root, spec)
            output = root / f"{style}-spec-labels.pptx"
            built = invoke(
                SKILL_RUNNER, "build", spec_path, "--out", output,
                "--mode", "full", "--style", style,
            )
            self.assertEqual(built.returncode, 0, msg=built.stderr)
            presentation = Presentation(str(output))
            figure_slide = presentation.slides[3]
            visible = {
                shape.text_frame.text
                for shape in figure_slide.shapes
                if getattr(shape, "has_text_frame", False)
            }
            self.assertTrue({"A", "B"}.issubset(visible))
            passed = invoke(
                SKILL_RUNNER, "qa", output, "--spec", spec_path,
                "--mode", "full", "--style", style, "--json",
            )
            self.assertEqual(passed.returncode, 0, msg=passed.stderr)

            for shape in figure_slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text == "B":
                    shape._element.getparent().remove(shape._element)
                    break
            damaged = root / f"{style}-missing-B.pptx"
            presentation.save(str(damaged))
            rejected = invoke(
                SKILL_RUNNER, "qa", damaged, "--spec", spec_path,
                "--mode", "full", "--style", style, "--json",
            )
            self.assertNotEqual(rejected.returncode, 0)
            report = json.loads(rejected.stdout)
            self.assertTrue(any("missing visible native panel" in value for value in report["failures"]))
            self.assertTrue(any(
                check["check"] == "pptx_panel_labels" and check["level"] == "FAIL"
                for check in report["presentation"]["deck_checks"]
            ))

    def test_standard_renders_spec_panel_labels_and_final_gate_detects_removal(self) -> None:
        self.check_specified_panel_labels("standard")

    def test_nice_renders_spec_panel_labels_and_final_gate_detects_removal(self) -> None:
        self.check_specified_panel_labels("nice")


class IntegratedImageRegressionTests(unittest.TestCase):
    def check_panel_cleanup_rejection(self, *, metadata: dict, notes: str | None = None,
                                      slide: dict | None = None) -> dict:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            _, _, rendered = create_synthetic_inversion_fixture(root)
            asset = write_asset(root, "Figure_01.png", source=str(rendered), **metadata)
            spec = full_spec()
            figure = add_figure(
                spec,
                asset,
                caption="Figure 1. Synthetic panel-integrity regression.",
                notes=notes or "【圖片說明】🖼️【A 圖】與【B 圖】為虛構影像。",
            )
            if slide:
                figure.update(slide)
            result = invoke(
                SKILL_RUNNER,
                "qa-spec",
                write_spec(root, spec),
                "--mode",
                "full",
                "--json",
            )
            self.assertNotEqual(result.returncode, 0, msg=result.stderr or result.stdout)
            return json.loads(result.stdout)

    def test_preserved_panel_notes_cannot_reference_nonexistent_embedded_label(self) -> None:
        report = self.check_panel_cleanup_rejection(
            metadata={
                "source_label_policy": "preserve",
                "native_labels": False,
                "embedded_labels": ["A", "B"],
            },
            notes="【圖片說明】🖼️【C 圖】並不存在。",
        )
        self.assertTrue(any("reference missing panel" in failure for failure in report["failures"]))

    def test_preserved_source_labels_cannot_also_request_native_labels(self) -> None:
        report = self.check_panel_cleanup_rejection(
            metadata={
                "source_label_policy": "preserve",
                "native_labels": True,
                "embedded_labels": ["A", "B"],
            }
        )
        self.assertTrue(any("duplicate" in failure for failure in report["failures"]))

    def test_preserved_source_labels_cannot_also_use_slide_spec_labels(self) -> None:
        report = self.check_panel_cleanup_rejection(
            metadata={
                "source_label_policy": "preserve",
                "native_labels": False,
                "embedded_labels": ["A", "B"],
            },
            slide={"panel_labels": ["A", "B"], "panel_label_x_fracs": [0.4, 0.9]},
        )
        self.assertTrue(any("duplicate" in failure for failure in report["failures"]))

    def test_panel_label_cleanup_cannot_overwrite_clinical_pixels(self) -> None:
        report = self.check_panel_cleanup_rejection(
            metadata={
                "source_label_policy": "preserve",
                "native_labels": False,
                "embedded_labels": ["A", "B"],
                "max_edge_px": 4,
                "panel_cleanup": [
                    {"label_overwritten_pixels": 91, "edge_trim_px": {"left": 1}},
                ],
            }
        )
        self.assertTrue(any("overwrites 91" in failure for failure in report["failures"]))

    def test_panel_edge_cleanup_cannot_exceed_its_declared_pixel_limit(self) -> None:
        report = self.check_panel_cleanup_rejection(
            metadata={
                "source_label_policy": "preserve",
                "native_labels": False,
                "embedded_labels": ["A", "B"],
                "max_edge_px": 4,
                "panel_cleanup": [
                    {"label_overwritten_pixels": 0, "edge_trim_px": {"left": 5}},
                ],
            }
        )
        self.assertTrue(any("4px edge-cleanup limit" in failure for failure in report["failures"]))

    def test_panel_boundary_adjustment_cannot_exceed_its_declared_pixel_limit(self) -> None:
        report = self.check_panel_cleanup_rejection(
            metadata={
                "source_label_policy": "preserve",
                "native_labels": False,
                "embedded_labels": ["A", "B"],
                "max_boundary_shift_px": 12,
                "panel_cleanup": [{
                    "label_overwritten_pixels": 0,
                    "edge_trim_px": {"left": 0},
                    "boundary_adjustments": [{
                        "shift_px": 13,
                        "reason": "preserve-complete-embedded-label-frame",
                    }],
                }],
            }
        )
        self.assertTrue(any("12px label-safe boundary" in failure for failure in report["failures"]))

    def test_panel_boundary_adjustment_requires_a_verified_label_frame(self) -> None:
        report = self.check_panel_cleanup_rejection(
            metadata={
                "source_label_policy": "preserve",
                "native_labels": False,
                "embedded_labels": ["A", "B"],
                "max_boundary_shift_px": 12,
                "panel_cleanup": [{
                    "label_overwritten_pixels": 0,
                    "edge_trim_px": {"left": 0},
                    "boundary_adjustments": [{
                        "shift_px": 4,
                        "reason": "arbitrary-clinical-image-crop",
                    }],
                }],
            }
        )
        self.assertTrue(any("label-safe boundary" in failure for failure in report["failures"]))

    def test_pdf_decode_array_is_detected_by_the_bundled_global_skill(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            manifest, raw, rendered = create_synthetic_inversion_fixture(Path(temporary))
            result = invoke(SKILL_RUNNER, "image-qa", manifest, "--json")
            self.assertEqual(result.returncode, 0, msg=result.stderr)
            report = json.loads(result.stdout)
            self.assertEqual(report["unsafe_raw_streams"], 1)
            self.assertEqual(report["figures"][0]["source_path"], str(raw.resolve()))
            self.assertEqual(report["figures"][0]["rendered_path"], str(rendered.resolve()))

    def test_synthetic_multiasset_chain_rejects_inverted_intermediate_sources(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            _, raw, rendered = create_synthetic_inversion_fixture(root)
            intermediate = write_asset(root, "Figure_01_panel_b_intermediate.png", source=str(raw))
            panel_figure = write_asset(
                root, "Figure_01.png", source_inputs=[str(rendered), str(intermediate)], native_labels=True
            )
            single_figure = write_asset(root, "Figure_02.png", source=str(raw))
            spec = full_spec()
            first = add_figure(spec, panel_figure, caption="Figure 1. Synthetic A/B panels.")
            slides = spec["slides"]
            assert isinstance(slides, list)
            slides[4] = {
                **first,
                "image": f"final_assets/{single_figure.name}",
                "caption": "Figure 2. Synthetic grayscale regression.",
            }
            path = write_spec(root, spec)
            result = invoke(SKILL_RUNNER, "qa-spec", path, "--mode", "full", "--style", "nice", "--json")
            self.assertNotEqual(result.returncode, 0)
            report = json.loads(result.stdout)
            self.assertTrue(any("inverted raw PDF image" in value for value in report["failures"]))
            self.assertTrue(any(raw.name in value for value in report["failures"]))

    def test_safe_decoded_panel_sources_pass_the_integrated_gate(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            _, _, rendered = create_synthetic_inversion_fixture(root)
            figure = write_asset(root, "Figure_01.png", source=str(rendered))
            spec = full_spec()
            add_figure(spec, figure, caption="Figure 1. Synthetic corrected grayscale.")
            result = invoke(SKILL_RUNNER, "qa-spec", write_spec(root, spec), "--mode", "full", "--json")
            self.assertEqual(result.returncode, 0, msg=result.stderr)
            self.assertEqual(json.loads(result.stdout)["image_polarity"]["unsafe_raw_streams"], 1)

    def test_synthetic_orphan_panel_without_sidecar_cannot_bypass_polarity(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            _, raw, rendered = create_synthetic_inversion_fixture(root)
            orphan = root / "orphan_Figure_01_panel_b_intermediate.png"
            with Image.open(raw) as inverted:
                inverted.crop((0, 0, 80, inverted.height)).save(orphan)
            figure = write_asset(
                root, "Figure_01.png", source_inputs=[str(rendered), str(orphan)], native_labels=True
            )
            spec = full_spec()
            add_figure(
                spec, figure, caption="Figure 1. Synthetic A/B panels.",
                panel_labels=["A", "B"], panel_label_x_fracs=[0.5, 1.0],
            )
            result = invoke(
                SKILL_RUNNER, "qa-spec", write_spec(root, spec), "--mode", "full", "--json"
            )
            self.assertNotEqual(result.returncode, 0)
            failures = json.loads(result.stdout)["failures"]
            self.assertTrue(any("no provenance sidecar" in value for value in failures))
            self.assertTrue(any(orphan.name in value for value in failures))

    def test_raster_figure_without_extraction_manifest_fails_closed(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            figure = write_asset(root, "Figure_01.png")
            spec = full_spec()
            add_figure(spec, figure, caption="Figure 1. Synthetic unverified panel.")
            result = invoke(
                SKILL_RUNNER, "qa-spec", write_spec(root, spec), "--mode", "full", "--json"
            )
            self.assertNotEqual(result.returncode, 0)
            self.assertTrue(any(
                "require a readable extraction manifest" in value
                for value in json.loads(result.stdout)["failures"]
            ))

    def test_pdf_rendered_vector_flowchart_preserves_audited_pdf_source(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            manifest, _, _ = create_synthetic_inversion_fixture(root)
            source_pdf = json.loads(manifest.read_text(encoding="utf-8"))["pdf"]
            figure = write_asset(
                root, "Figure_01.png", source=source_pdf, asset_type="flowchart"
            )
            spec = full_spec()
            add_figure(spec, figure, caption="Figure 1. PDF-rendered vector flowchart.")
            result = invoke(
                SKILL_RUNNER, "qa-spec", write_spec(root, spec), "--mode", "full", "--json"
            )
            self.assertEqual(result.returncode, 0, msg=result.stderr)

    def test_strict_full_slide_budget_is_enforced_before_build(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            specification = full_spec()
            slides = specification["slides"]
            assert isinstance(slides, list)
            slides.pop(3)
            result = invoke(
                SKILL_RUNNER, "qa-spec", write_spec(Path(temporary), specification),
                "--mode", "full", "--json",
            )
            self.assertNotEqual(result.returncode, 0)
            self.assertTrue(any("40-55" in value for value in json.loads(result.stdout)["failures"]))


class StandaloneReleaseTests(unittest.TestCase):
    def test_public_release_filter_denies_unrecognized_patient_and_credential_files(self) -> None:
        skill_prefix = Path(".agents") / "skills" / SKILL_NAME
        for private in (
            Path("patients.csv"), Path("private-patient-image.png"), Path("credentials.yaml"),
            Path("api-token.txt"), Path("id_rsa"), Path("notes/patient-history.docx"),
            Path("skill-work/synthetic-panel.png"), skill_prefix / "assets" / "private-clinical-image.png",
            skill_prefix / "agents" / "credentials.yaml", skill_prefix / "references" / "patient.csv",
            skill_prefix / "scripts" / "token.txt", Path("docs") / "patient-history.md",
        ):
            with self.subTest(private=private):
                self.assertFalse(package_release.should_package(private))

    def test_release_discovery_never_enumerates_unrecognized_private_files(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            skill = root / ".agents" / "skills" / SKILL_NAME
            paths = {
                Path("README.md"): "public documentation",
                Path("patients.csv"): "fictional private row",
                Path("credentials.yaml"): "fictional-token",
                Path("docs/patient-history.md"): "fictional patient history",
                Path(".agents/skills") / SKILL_NAME / "SKILL.md": "public skill",
                Path(".agents/skills") / SKILL_NAME / "assets/dr_leether_logo.png": "public logo",
                Path(".agents/skills") / SKILL_NAME / "assets/private-patient-image.png": "fictional image",
                Path(".agents/skills") / SKILL_NAME / "agents/credentials.yaml": "fictional secret",
            }
            for relative, contents in paths.items():
                path = root / relative
                path.parent.mkdir(parents=True, exist_ok=True)
                path.write_text(contents, encoding="utf-8")
            with mock.patch.object(package_release, "PROJECT_ROOT", root):
                selected = package_release.release_files()
            self.assertEqual(
                set(selected),
                {
                    Path("README.md"),
                    skill.relative_to(root) / "SKILL.md",
                    skill.relative_to(root) / "assets" / "dr_leether_logo.png",
                },
            )

    def test_global_release_contains_full_skill_and_no_pdf_pptx_or_private_files(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            destination = Path(temporary) / "integrated.zip"
            result = invoke(PROJECT_ROOT / "tools" / "package_release.py", "--kind", "skill", "--out", destination)
            self.assertEqual(result.returncode, 0, msg=result.stderr)
            payload = json.loads(result.stdout)
            self.assertTrue(Path(payload["checksum_file"]).is_file())
            with zipfile.ZipFile(destination) as archive:
                entries = archive.namelist()
                self.assertTrue(any(name.endswith("/skill/SKILL.md") for name in entries))
                self.assertTrue(any(name.endswith("/skill/scripts/build_deck_nice.py") for name in entries))
                self.assertTrue(any(name.endswith("/skill/scripts/qa_gate.py") for name in entries))
                self.assertTrue(any(name.endswith("/install-global.py") for name in entries))
                self.assertFalse(any(name.lower().endswith((".pdf", ".pptx")) for name in entries))
                self.assertFalse(any("/.skill-work/" in name for name in entries))

    def test_global_release_is_reproducible(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            first = invoke(
                PROJECT_ROOT / "tools" / "package_release.py", "--kind", "skill",
                "--out", Path(temporary) / "first.zip",
            )
            second = invoke(
                PROJECT_ROOT / "tools" / "package_release.py", "--kind", "skill",
                "--out", Path(temporary) / "second.zip",
            )
            self.assertEqual(first.returncode, 0, msg=first.stderr)
            self.assertEqual(second.returncode, 0, msg=second.stderr)
            self.assertEqual(json.loads(first.stdout)["sha256"], json.loads(second.stdout)["sha256"])

    def test_release_installs_after_extraction_outside_the_repository(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            archive_path = root / "skill.zip"
            packaged = invoke(
                PROJECT_ROOT / "tools" / "package_release.py", "--kind", "skill", "--out", archive_path
            )
            self.assertEqual(packaged.returncode, 0, msg=packaged.stderr)
            extracted = root / "release"
            with zipfile.ZipFile(archive_path) as archive:
                archive.extractall(extracted)
            release_root = next(path for path in extracted.iterdir() if path.is_dir())
            target = root / "globally-installed"
            installed = invoke(
                release_root / "install-global.py", "install", "--target", target,
                "--skip-deps", "--json", cwd=root,
            )
            self.assertEqual(installed.returncode, 0, msg=installed.stderr)
            self.assertTrue((target / SKILL_NAME / "scripts" / "workflow.py").is_file())


class SyntheticPrepareRegression(unittest.TestCase):
    def test_prepare_detects_and_corrects_a_synthetic_decode_inversion(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            create_synthetic_inversion_fixture(root, source_size=(280, 240), source_format="JPEG")
            paper = root / "synthetic-inverted-image.pdf"
            result = invoke(
                SKILL_RUNNER, "prepare", paper, "--workspace", root,
                "--output-dir", root, "--mode", "full", "--style", "standard", "--json",
                cwd=root,
            )
            self.assertEqual(result.returncode, 0, msg=result.stderr)
            payload = json.loads(result.stdout)
            report = json.loads(Path(payload["image_polarity_audit"]["report"]).read_text(encoding="utf-8"))
            unsafe = [entry for entry in report["figures"] if entry["raw"]["status"] == "inverted"]
            self.assertGreaterEqual(len(unsafe), 1)
            self.assertEqual({entry["page"] for entry in unsafe}, {1})
            self.assertTrue(all(entry["rendered_polarity"]["status"] == "correct" for entry in unsafe))


if __name__ == "__main__":
    unittest.main()
