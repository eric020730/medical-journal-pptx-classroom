#!/usr/bin/env python3
"""Validate classroom deck specifications and finished PowerPoint files."""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import image_polarity


PROJECT_ROOT = Path(__file__).resolve().parents[1]
CONFIG = json.loads((PROJECT_ROOT / ".classroom-project.json").read_text(encoding="utf-8"))
SKILL_ROOT = PROJECT_ROOT / ".agents" / "skills" / CONFIG["skill_name"]
LOGO_PATH = SKILL_ROOT / "assets" / "dr_leether_logo.png"
CHINESE_RE = re.compile(r"[\u3400-\u4dbf\u4e00-\u9fff]")
FIGURE_RE = re.compile(r"\bfig(?:ure)?\.?\s*(\d+)\b", re.IGNORECASE)
PANEL_FRAGMENT_RE = re.compile(r"(?:^|[_-])panel[_-]?[a-z0-9]+$", re.IGNORECASE)
RAW_STREAM_RE = re.compile(r"^image_p\d+_\d+$", re.IGNORECASE)
SPLIT_TABLE_RE = re.compile(r"^(table[_-]?\d+)[_-]?([a-z])$", re.IGNORECASE)
KEYCAP_RE = re.compile(r"[0-9]\ufe0f?\u20e3")
TABLE_MARGIN_MIN = 8
POSTPROCESS_SUFFIX = ".postprocess.json"
EMOJI_RANGES = ((0x1F300, 0x1FAFF), (0x2600, 0x27BF), (0x2B00, 0x2BFF))


def has_emoji(text: str) -> bool:
    return any(
        any(start <= ord(character) <= end for start, end in EMOJI_RANGES)
        for character in text
    )


def visible_texts(slide: Any) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
    return texts


def visible_spec_fields(slide: dict[str, Any]) -> list[tuple[str, str]]:
    fields: list[tuple[str, str]] = []
    for key in ("title", "caption", "authors", "citation", "subtitle", "kicker"):
        value = slide.get(key)
        if isinstance(value, str) and value.strip():
            fields.append((key, value))
    for key in ("bullets", "items", "body"):
        value = slide.get(key)
        if isinstance(value, str) and value.strip():
            fields.append((key, value))
        elif isinstance(value, list):
            fields.extend(
                (f"{key}[{index}]", item)
                for index, item in enumerate(value)
                if isinstance(item, str) and item.strip()
            )
    return fields


def referenced_panel_letters(notes: str) -> set[str]:
    letters = {
        match.group(1).upper()
        for match in re.finditer(r"【\s*([A-Za-z])\s*[圖图]", notes)
    }
    letters.update(
        match.group(1).upper()
        for match in re.finditer(r"\bpanel\s*([A-Za-z])\b", notes, re.IGNORECASE)
    )
    return letters


def resolve_asset(value: str | Path, spec_path: Path) -> Path:
    asset = Path(value).expanduser()
    return asset.resolve() if asset.is_absolute() else (spec_path.parent / asset).resolve()


def read_sidecar(image_path: Path) -> dict[str, Any] | None:
    sidecar = image_path.with_suffix(image_path.suffix + POSTPROCESS_SUFFIX)
    if not sidecar.is_file():
        return None
    try:
        value = json.loads(sidecar.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        return {}
    return value if isinstance(value, dict) else {}


def split_table_group(image_path: Path) -> str | None:
    match = SPLIT_TABLE_RE.match(image_path.stem)
    return match.group(1).replace("-", "_").lower() if match else None


def extraction_manifest(spec_path: Path, spec: dict[str, Any]) -> Path | None:
    meta = spec.get("meta", {})
    if isinstance(meta, dict):
        for key in ("extraction_manifest", "extracted_manifest"):
            value = meta.get(key)
            if isinstance(value, str) and value.strip():
                candidate = resolve_asset(value, spec_path)
                if candidate.is_file():
                    return candidate
    for candidate in (
        spec_path.parent / "extracted" / "manifest.json",
        spec_path.parent.parent / "extracted" / "manifest.json",
    ):
        if candidate.is_file():
            return candidate.resolve()
    return None


def _check_split_tables(
    groups: dict[str, list[tuple[int, dict[str, Any], Path]]],
    failures: list[str],
    warnings: list[str],
) -> None:
    for group, entries in sorted(groups.items()):
        if len(entries) < 2:
            continue
        widths: dict[str, int] = {}
        on_screen: list[float | None] = []
        for _, entry, path in entries:
            try:
                with Image.open(path) as image:
                    widths[path.name] = image.width
            except (OSError, ValueError) as error:
                failures.append(f"Split table asset {path.name} cannot be opened: {error}")
                continue
            value = entry.get("image_width_in")
            try:
                on_screen.append(float(value) if value is not None else None)
            except (TypeError, ValueError):
                failures.append(f"Split table {group} has an invalid image_width_in: {value!r}.")

        if len(set(widths.values())) > 1:
            failures.append(f"Split table {group} has unequal image pixel widths: {widths}.")
        configured = {value for value in on_screen if value is not None}
        if len(configured) > 1:
            failures.append(
                f"Split table {group} has unequal on-screen image_width_in values: "
                f"{sorted(configured)}."
            )
        elif None in on_screen:
            warnings.append(
                f"Split table {group} should set the same image_width_in on every slide."
            )


def _check_content_structure(
    slides: list[dict[str, Any]], failures: list[str], warnings: list[str], mode: str
) -> None:
    content = [slide for slide in slides if slide.get("type") == "content"]
    flat: list[int] = []
    for index, slide in enumerate(slides, start=1):
        if slide.get("type") != "content":
            continue
        items: list[str] = []
        for key in ("bullets", "body", "items"):
            value = slide.get(key)
            if isinstance(value, list):
                items = [str(item) for item in value if str(item).strip()]
                break
            if isinstance(value, str):
                items = [line for line in value.splitlines() if line.strip()]
                break
        if len(items) < 3:
            continue
        joined = "\n".join(items)
        has_heading = any(re.search(r"[A-Za-z][\w /-]{2,}:", item) for item in items)
        has_logic = any(marker in joined for marker in ("→", "✅", "⚠"))
        if not has_heading and not has_logic:
            flat.append(index)

    if flat and mode == "full" and len(flat) > len(content) / 2:
        failures.append(
            f"{len(flat)}/{len(content)} teaching slides are flat bullet lists "
            f"without structured headings or clinical takeaways: {flat[:10]}."
        )
    elif flat:
        warnings.append(
            f"Teaching slide(s) {flat} would benefit from section headings or clinical takeaways."
        )


def validate_specification(
    spec_path: Path, *, mode: str = "full", audit_images: bool = True
) -> dict[str, Any]:
    spec_path = spec_path.expanduser().resolve()
    failures: list[str] = []
    warnings: list[str] = []
    if not spec_path.is_file():
        return {
            "ok": False,
            "spec": str(spec_path),
            "mode": mode,
            "failures": [f"Deck specification not found: {spec_path}"],
            "warnings": [],
        }
    try:
        spec = json.loads(spec_path.read_text(encoding="utf-8"))
        slides = spec.get("slides", [])
        if not isinstance(spec, dict) or not isinstance(slides, list):
            raise ValueError("The specification must contain a slides array.")
        if not all(isinstance(slide, dict) for slide in slides):
            raise ValueError("Each slide must be a JSON object.")
    except (json.JSONDecodeError, OSError, TypeError, ValueError, AttributeError) as error:
        return {
            "ok": False,
            "spec": str(spec_path),
            "mode": mode,
            "failures": [f"Deck specification is invalid: {error}"],
            "warnings": [],
        }

    if mode in CONFIG["modes"]:
        limits = CONFIG["modes"][mode]
        if not limits["minimum_slides"] <= len(slides) <= limits["maximum_slides"]:
            failures.append(
                f"{mode} mode expects {limits['minimum_slides']}-{limits['maximum_slides']} "
                f"slides but the specification contains {len(slides)}."
            )
    elif mode != "smoke":
        failures.append(f"Unknown QA mode: {mode}")

    counts: dict[str, int] = {}
    figures: dict[str, list[int]] = {}
    split_tables: dict[str, list[tuple[int, dict[str, Any], Path]]] = {}
    emoji_notes = 0

    for index, slide in enumerate(slides, start=1):
        slide_type = str(slide.get("type", "missing"))
        counts[slide_type] = counts.get(slide_type, 0) + 1
        for field, value in visible_spec_fields(slide):
            if CHINESE_RE.search(value):
                failures.append(
                    f"Slide {index} ({slide_type}) visible field {field} contains Chinese text."
                )

        notes = str(slide.get("notes", "")).strip()
        if not notes:
            failures.append(f"Slide {index} has no speaker notes in its specification.")
        elif not CHINESE_RE.search(notes):
            failures.append(f"Slide {index} speaker notes contain no Chinese text.")
        if has_emoji(notes):
            emoji_notes += 1

        if slide_type != "figure":
            continue
        image_name = slide.get("image")
        if not isinstance(image_name, str) or not image_name.strip():
            failures.append(f"Figure slide {index} does not reference an image.")
            continue
        image = resolve_asset(image_name, spec_path)
        if not image.is_file():
            failures.append(f"Figure slide {index} image is missing: {image}")
            continue
        if PANEL_FRAGMENT_RE.search(image.stem):
            failures.append(f"Figure slide {index} references a panel fragment: {image.name}")
        if RAW_STREAM_RE.match(image.stem):
            failures.append(
                f"Figure slide {index} references raw PDF image {image.name}; use the "
                "color-decoded extracted/figures image instead."
            )

        sidecar_path = image.with_suffix(image.suffix + POSTPROCESS_SUFFIX)
        sidecar = read_sidecar(image)
        if image.suffix.lower() in {".png", ".jpg", ".jpeg"} and sidecar is None:
            failures.append(f"Figure slide {index} is missing its image sidecar: {sidecar_path}")
        elif sidecar == {}:
            failures.append(f"Figure slide {index} has an invalid image sidecar: {sidecar_path}")

        labels = slide.get("panel_labels") or []
        if not isinstance(labels, list):
            failures.append(f"Figure slide {index} panel_labels must be a list.")
            labels = []
        if len(labels) > 1 and not slide.get("panel_geometry_exception"):
            fractions = slide.get("panel_label_x_fracs") or []
            boxes = slide.get("panel_boxes") or []
            sidecar_boxes = (sidecar or {}).get("panel_boxes") or []
            valid_geometry = (
                isinstance(fractions, list) and len(fractions) >= len(labels)
            ) or (
                isinstance(boxes, list) and len(boxes) >= len(labels)
            ) or (
                isinstance(sidecar_boxes, list) and len(sidecar_boxes) >= len(labels)
            ) or bool((sidecar or {}).get("native_labels"))
            if not valid_geometry:
                failures.append(
                    f"Figure slide {index} has {len(labels)} panel labels but no panel geometry."
                )
        label_set = {str(label).strip().upper() for label in labels}
        unknown = referenced_panel_letters(notes) - label_set if label_set else set()
        if unknown:
            failures.append(
                f"Figure slide {index} notes reference missing panel(s): {sorted(unknown)}."
            )

        caption = f"{slide.get('caption', '')} {slide.get('title', '')}"
        is_table = (
            "table" in image.name.lower()
            or str((sidecar or {}).get("asset_type", "")).lower() == "table"
            or bool(re.search(r"\btable\b", caption, re.IGNORECASE))
        )
        if is_table:
            margin = (sidecar or {}).get("margin")
            if isinstance(margin, int) and margin < TABLE_MARGIN_MIN:
                failures.append(
                    f"Figure slide {index} table {image.name} has an unsafe "
                    f"{margin}px edge margin; at least {TABLE_MARGIN_MIN}px is required."
                )
            group = split_table_group(image)
            if group:
                split_tables.setdefault(group, []).append((index, slide, image))
        else:
            match = FIGURE_RE.search(caption)
            if match:
                figures.setdefault(match.group(1), []).append(index)

    if counts.get("title", 0) != 1:
        failures.append(f"Expected exactly one title slide; found {counts.get('title', 0)}.")
    if counts.get("thanks", 0) != 1:
        failures.append(f"Expected exactly one closing slide; found {counts.get('thanks', 0)}.")
    if counts.get("content", 0) == 0:
        failures.append("Presentation contains no teaching-content slides.")
    if mode == "full":
        for required in ("outline", "references"):
            if counts.get(required, 0) == 0:
                failures.append(f"Full teaching deck is missing its {required} slide.")
        if counts.get("part", 0) == 0:
            warnings.append("Full teaching deck does not contain any numbered section divider.")

    title = next((slide for slide in slides if slide.get("type") == "title"), None)
    if title is not None:
        missing = [key for key in ("authors", "citation") if not str(title.get(key, "")).strip()]
        if missing:
            failures.append(f"Title slide is missing source-paper metadata: {', '.join(missing)}.")
    if slides and emoji_notes == 0:
        failures.append("Speaker notes do not contain any structural emoji markers.")

    meta = spec.get("meta", {})
    if isinstance(meta, dict):
        footer = meta.get("footer_label", "")
        if isinstance(footer, str) and CHINESE_RE.search(footer):
            failures.append("Visible meta.footer_label contains Chinese text.")

    outline = next((slide for slide in slides if slide.get("type") == "outline"), None)
    if outline is not None:
        items = outline.get("items") or []
        if isinstance(items, list) and items:
            if sum(bool(KEYCAP_RE.search(str(item))) for item in items) < len(items):
                warnings.append("Outline items should use numbered keycap markers.")
            if any(not re.search(r"slides?\s*\d+\s*[-–]\s*\d+", str(item), re.IGNORECASE)
                   for item in items):
                warnings.append("Outline items should indicate their corresponding slide ranges.")

    for number, indexes in sorted(figures.items()):
        if len(indexes) > 1:
            failures.append(
                f"Paper Figure {number} appears on multiple slides: "
                + ", ".join(str(index) for index in indexes)
            )
    _check_split_tables(split_tables, failures, warnings)
    _check_content_structure(slides, failures, warnings, mode)

    polarity_summary: dict[str, Any] | None = None
    manifest = extraction_manifest(spec_path, spec)
    if audit_images and manifest is not None:
        polarity = image_polarity.audit_extraction(manifest)
        failures.extend(polarity.get("failures", []))
        warnings.extend(polarity.get("warnings", []))
        final_assets = image_polarity.audit_final_assets(spec_path, polarity)
        failures.extend(final_assets.get("failures", []))
        polarity_summary = {
            "checked_figures": polarity.get("checked_figures", 0),
            "unsafe_raw_streams": polarity.get("unsafe_raw_streams", 0),
            "checked_final_assets": final_assets.get("checked_assets", 0),
        }
    elif audit_images and counts.get("figure", 0):
        warnings.append(
            "No extraction manifest was found; original-PDF image polarity could not be verified."
        )

    return {
        "ok": not failures,
        "stage": "specification",
        "spec": str(spec_path),
        "mode": mode,
        "slides": len(slides),
        "slides_with_emoji_notes": emoji_notes,
        "slide_types": counts,
        "image_polarity": polarity_summary,
        "failures": failures,
        "warnings": warnings,
    }


def logo_hash() -> str | None:
    if not LOGO_PATH.is_file():
        return None
    return hashlib.sha256(LOGO_PATH.read_bytes()).hexdigest()


def slide_has_logo(slide: Any, expected_hash: str) -> bool:
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            if hashlib.sha256(shape.image.blob).hexdigest() == expected_hash:
                return True
        except (AttributeError, ValueError):
            continue
    return False


def validate_presentation(
    pptx_path: Path, *, spec_path: Path | None = None, mode: str = "full"
) -> dict[str, Any]:
    failures: list[str] = []
    warnings: list[str] = []
    if not pptx_path.is_file():
        return {"ok": False, "failures": [f"Presentation not found: {pptx_path}"], "warnings": []}
    try:
        presentation = Presentation(str(pptx_path))
    except Exception as error:
        return {"ok": False, "failures": [f"PowerPoint cannot be opened: {error}"], "warnings": []}

    slides = list(presentation.slides)
    count = len(slides)
    if not slides:
        failures.append("Presentation does not contain any slides.")
    if mode in CONFIG["modes"]:
        limits = CONFIG["modes"][mode]
        if not limits["minimum_slides"] <= count <= limits["maximum_slides"]:
            failures.append(
                f"{mode} mode expects {limits['minimum_slides']}-{limits['maximum_slides']} "
                f"slides but the presentation contains {count}."
            )
    elif mode != "smoke":
        failures.append(f"Unknown QA mode: {mode}")

    width, height = int(presentation.slide_width), int(presentation.slide_height)
    if height <= 0 or abs(width / height - 16 / 9) > 0.03:
        failures.append(f"Presentation is not widescreen 16:9 ({width} × {height}).")

    spec_slides: list[dict[str, Any]] = []
    spec_report: dict[str, Any] | None = None
    if spec_path is not None:
        spec_report = validate_specification(spec_path, mode=mode)
        failures.extend(spec_report["failures"])
        warnings.extend(spec_report["warnings"])
        if spec_path.is_file():
            try:
                spec_slides = list(json.loads(spec_path.read_text(encoding="utf-8")).get("slides", []))
            except (json.JSONDecodeError, OSError, TypeError, AttributeError):
                spec_slides = []
            if spec_slides and len(spec_slides) != count:
                failures.append(
                    f"Deck specification contains {len(spec_slides)} slides, "
                    f"but the PowerPoint contains {count}."
                )

    expected_logo = logo_hash()
    if expected_logo is None:
        failures.append(f"Bundled presentation logo is missing: {LOGO_PATH}")

    note_count = 0
    emoji_note_count = 0
    logo_count = 0
    visible_cjk_slides: list[int] = []
    for index, slide in enumerate(slides, start=1):
        if any(CHINESE_RE.search(text) for text in visible_texts(slide)):
            visible_cjk_slides.append(index)
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except (AttributeError, ValueError):
            notes = ""
        if not notes:
            failures.append(f"Slide {index} has no speaker notes.")
        else:
            note_count += 1
            if not CHINESE_RE.search(notes):
                failures.append(f"Slide {index} speaker notes contain no Chinese text.")
            if "**" in notes:
                failures.append(f"Slide {index} speaker notes contain unrendered ** markup.")
            if has_emoji(notes):
                emoji_note_count += 1

        present_logo = bool(expected_logo and slide_has_logo(slide, expected_logo))
        logo_count += int(present_logo)
        if index <= len(spec_slides):
            slide_type = spec_slides[index - 1].get("type")
            if slide_type not in {"title", "thanks"} and not present_logo:
                failures.append(f"Slide {index} ({slide_type}) does not contain the bundled logo.")

    if visible_cjk_slides:
        failures.append(
            "English-only slide text contains Chinese on slide(s): "
            + ", ".join(str(index) for index in visible_cjk_slides)
        )
    if slides and emoji_note_count == 0:
        failures.append("Speaker notes do not contain any structural emoji markers.")
    if spec_path is None and count > 2 and logo_count == 0:
        failures.append("No embedded presentation logo was found.")

    return {
        "ok": not failures,
        "stage": "presentation",
        "pptx": str(pptx_path),
        "spec": str(spec_path) if spec_path else None,
        "mode": mode,
        "slides": count,
        "slides_with_notes": note_count,
        "slides_with_emoji_notes": emoji_note_count,
        "slides_with_logo": logo_count,
        "image_polarity": spec_report.get("image_polarity") if spec_report else None,
        "failures": list(dict.fromkeys(failures)),
        "warnings": list(dict.fromkeys(warnings)),
    }


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("target", type=Path, help="PowerPoint or, with --spec-only, deck JSON")
    parser.add_argument("--spec", type=Path)
    parser.add_argument("--spec-only", action="store_true", help="Validate a deck before building")
    parser.add_argument("--mode", choices=(*CONFIG["modes"], "smoke"), default="full")
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args(argv)

    target = args.target.expanduser().resolve()
    if args.spec_only:
        report = validate_specification(target, mode=args.mode)
    else:
        report = validate_presentation(
            target,
            spec_path=args.spec.expanduser().resolve() if args.spec else None,
            mode=args.mode,
        )

    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        if "slides" in report:
            if args.spec_only:
                print(f"Specification slides: {report['slides']} | Mode: {report['mode']}")
            else:
                print(
                    f"Slides: {report['slides']} | Notes: {report['slides_with_notes']} | "
                    f"Logo: {report['slides_with_logo']} | Mode: {report['mode']}"
                )
        polarity = report.get("image_polarity")
        if polarity:
            print(
                f"PDF image polarity: {polarity['checked_figures']} checked; "
                f"{polarity['unsafe_raw_streams']} unsafe raw stream(s) isolated."
            )
        for failure in report["failures"]:
            print(f"[FAIL] {failure}", file=sys.stderr)
        for warning in report["warnings"]:
            print(f"[WARN] {warning}", file=sys.stderr)
        stage = "Specification" if args.spec_only else "Presentation"
        print(f"{stage} QA {'passed' if report['ok'] else 'failed'}.")

    return 0 if report["ok"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
